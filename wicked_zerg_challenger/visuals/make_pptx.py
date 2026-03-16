"""
Swarm-Net Presentation PPTX Generator
======================================
python-pptx로 텍스트 + 도형 기반의 편집 가능한 슬라이드 생성
다크 블루 배경 + Cyan/Green 악센트 테마
"""

import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn


# ===== Color Constants =====
BG_DARK   = RGBColor(0x0B, 0x10, 0x1F)   # Dark navy
BG_CARD   = RGBColor(0x12, 0x1A, 0x2E)   # Card background
CYAN      = RGBColor(0x00, 0xFF, 0xCC)    # Primary accent
CYAN_DIM  = RGBColor(0x00, 0xB8, 0x94)    # Dimmed cyan
ORANGE    = RGBColor(0xFF, 0x98, 0x00)    # Warning
RED       = RGBColor(0xFF, 0x45, 0x45)    # Alert
GREEN     = RGBColor(0x00, 0xE6, 0x76)    # Normal
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xE8, 0xE8, 0xE8)    # Body text (brighter)
DIM       = RGBColor(0xA0, 0xA8, 0xB8)    # Subtle text (higher contrast)
PURPLE    = RGBColor(0xB3, 0x88, 0xFF)    # Simulation accent
GOLD      = RGBColor(0xFF, 0xD7, 0x00)    # Stars


IMG_DIR = Path(__file__).parent / "images"


def add_image(slide, img_name, left, top, width=None, height=None):
    """Insert an image if the file exists. Returns True if inserted."""
    img_path = IMG_DIR / img_name
    if not img_path.exists():
        return False
    kwargs = {}
    if width:
        kwargs['width'] = width
    if height:
        kwargs['height'] = height
    slide.shapes.add_picture(str(img_path), left, top, **kwargs)
    return True


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Malgun Gothic"):
    """Add a text box with a single paragraph."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=LIGHT):
    """Add a text box with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(item, tuple):
            text, item_color, item_bold = item
        else:
            text, item_color, item_bold = item, color, False

        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = item_color
        p.font.bold = item_bold
        p.font.name = "Malgun Gothic"
        p.space_after = Pt(6)
        p.level = 0
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color=BG_CARD,
                     border_color=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape


def add_tag_label(slide, left, top, text, bg_color=CYAN, text_color=BG_DARK):
    """Add a small tag/label."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, Inches(2.4), Inches(0.34)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = "Consolas"
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_notes(slide, text):
    """Add speaker notes."""
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.text = text
    for para in tf.paragraphs:
        para.font.size = Pt(12)


def add_footer(slide, slide_num, total=8):
    """Add slide number footer."""
    add_text_box(
        slide, Inches(0.4), Inches(6.9), Inches(2), Inches(0.4),
        "SWARM-NET", font_size=9, color=CYAN_DIM, bold=True,
        font_name="Consolas"
    )
    add_text_box(
        slide, Inches(11), Inches(6.9), Inches(2), Inches(0.4),
        f"{slide_num:02d} / {total:02d}", font_size=9, color=DIM,
        font_name="Consolas", alignment=PP_ALIGN.RIGHT
    )


# ===================================================================
# SLIDE BUILDERS
# ===================================================================

def build_slide_1(prs):
    """Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_tag_label(slide, Inches(0.8), Inches(1.2), "CAPSTONE PROJECT", CYAN, BG_DARK)

    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.7), Inches(0.8), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = CYAN
    line.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(1.9), Inches(8), Inches(1.2),
                 "하늘의 새로운 질서를\n구축하다", font_size=44, color=WHITE, bold=True)

    add_text_box(slide, Inches(0.8), Inches(3.3), Inches(8), Inches(0.6),
                 "군집 UAV 기반 다이내믹 공역 통제 시스템", font_size=20, color=LIGHT)

    add_text_box(slide, Inches(0.8), Inches(3.9), Inches(8), Inches(0.4),
                 "Swarm-Net Airspace Manager", font_size=18, color=CYAN,
                 bold=True, font_name="Consolas")

    add_text_box(slide, Inches(0.8), Inches(4.8), Inches(8), Inches(0.7),
                 "StarCraft II 군단 제어 AI 기반\nSim-to-Real 관제(ATC) 솔루션",
                 font_size=15, color=DIM)

    # Right side - hexagon decoration
    hex_shape = slide.shapes.add_shape(
        MSO_SHAPE.HEXAGON, Inches(9.5), Inches(1.5), Inches(2.5), Inches(2.5)
    )
    hex_shape.fill.background()
    hex_shape.line.color.rgb = CYAN_DIM
    hex_shape.line.width = Pt(2)
    hex_shape.rotation = 30.0

    hex2 = slide.shapes.add_shape(
        MSO_SHAPE.HEXAGON, Inches(10), Inches(2.0), Inches(1.5), Inches(1.5)
    )
    hex2.fill.background()
    hex2.line.color.rgb = RGBColor(0x00, 0x88, 0x66)
    hex2.line.width = Pt(1)
    hex2.rotation = 30.0

    # Visual: Swarm Hex Network
    add_image(slide, 'slide1_swarm_hex.png',
              Inches(8.3), Inches(1.0), width=Inches(4.5))

    add_footer(slide, 1)

    add_notes(slide, (
        "안녕하십니까, 국립목포대학교 드론기계학과 장선우입니다. "
        "오늘 제가 제안할 솔루션은 단순히 하늘을 나는 기체를 넘어, "
        "하늘의 새로운 질서를 통제하는 시스템입니다. "
        "군사 통신 및 드론 운용 현장에서 느꼈던 한계를 소프트웨어로 극복하고자 기획한 "
        "'Swarm-Net: 군집 드론 기반 다이내믹 공역 통제 시스템' 발표를 시작하겠습니다."
    ))


def build_slide_2(prs):
    """Problem - Background & Problem."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Background: Radar coverage gap diagram
    add_image(slide, 'slide2_radar_gap.png',
              Inches(7.5), Inches(0.3), width=Inches(5.5))

    add_tag_label(slide, Inches(0.8), Inches(0.6), "BACKGROUND & PROBLEM", ORANGE, BG_DARK)

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(0.8), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = CYAN
    line.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(1.2),
                 "고정형 지상 레이더의\n구조적 한계와 급증하는 저고도 비행체",
                 font_size=36, color=WHITE, bold=True)

    # 3 problem cards
    cards = [
        ("탐지 사각지대", "산악·도심 고층 건물군의\n전파 음영으로 고정 레이더가\n커버하지 못하는 저고도 영역 분포", RED),
        ("신속 전개 불가", "고정 인프라 설치에 수 개월 소요\n긴급 상황 시 특정 지역으로\n관제망을 즉각 전개하기 어려움", ORANGE),
        ("소형 UAV 탐지 한계", "저고도 소형 UAV는 낮은 RCS로\n기존 항공 레이더의\n탐지 임계값 미달", PURPLE),
    ]

    for i, (title, desc, accent) in enumerate(cards):
        x = Inches(0.8 + i * 4.1)
        y = Inches(3.0)
        card = add_rounded_rect(slide, x, y, Inches(3.7), Inches(2.2), BG_CARD, accent)

        add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.1), Inches(0.4),
                     title, font_size=20, color=accent, bold=True)
        add_text_box(slide, x + Inches(0.3), y + Inches(0.85), Inches(3.1), Inches(1.2),
                     desc, font_size=14, color=LIGHT)

    # Stats bar
    stats = [("90만+", "국내 등록 무인비행장치 (2025)"), ("+30%", "산업용 드론 비행 연간 증가율"), ("67%", "도심 저고도 관제 사각지대 비율")]
    colors = [CYAN, GREEN, ORANGE]
    for i, (val, label) in enumerate(stats):
        x = Inches(0.8 + i * 3.5)
        add_text_box(slide, x, Inches(5.8), Inches(2.5), Inches(0.5),
                     val, font_size=32, color=colors[i], bold=True, font_name="Consolas")
        add_text_box(slide, x, Inches(6.35), Inches(2.5), Inches(0.3),
                     label, font_size=12, color=DIM)

    # Citation footnotes
    add_text_box(slide, Inches(0.8), Inches(6.6), Inches(11), Inches(0.25),
                 "※ 국토교통부 「무인비행장치 등록 현황」 2025  |  "
                 "※ KARI 「산업용 무인기 시장 동향」 2024  |  "
                 "※ KOTI 「저고도 공역 관제 사각지대 실태 조사」 2024",
                 font_size=7, color=DIM, font_name="Malgun Gothic")

    add_footer(slide, 2)

    add_notes(slide, (
        "국내 등록 드론 수가 기하급수적으로 늘어나면서 "
        "도심 상공은 통제하기 어려운 영역이 되었습니다. "
        "군 생활 당시 야전에서 경험한 바에 따르면, "
        "기존의 지상 기반 고정 레이더는 산악 지형이나 복잡한 환경에서 "
        "치명적인 사각지대를 발생시킵니다. "
        "긴급 상황이 발생했을 때 신속하게 방어망을 전개할 수 있는 "
        "유연한 관제 인프라가 절실히 필요한 시점입니다."
    ))


def build_slide_3(prs):
    """Solution - Our Solution: Swarm-Net."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_tag_label(slide, Inches(0.8), Inches(0.6), "OUR SOLUTION", CYAN, BG_DARK)

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(0.8), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = CYAN
    line.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(1.2),
                 "드론이 직접\n하늘의 관제탑이 되다",
                 font_size=36, color=WHITE, bold=True)

    add_text_box(slide, Inches(0.8), Inches(2.7), Inches(10), Inches(0.7),
                 "다수의 통제용 군집 UAV를 공중에 직접 배치하여\n"
                 "실시간 메시 통신망을 형성하는 '이동식 관제탑(Mobile ATC)' 아키텍처",
                 font_size=16, color=DIM)

    # Visual: Mobile ATC architecture diagram (right side)
    add_image(slide, 'slide3_mobile_atc.png',
              Inches(8.8), Inches(0.8), width=Inches(4.0))

    # 3 pillars
    pillars = [
        ("Swarm Fleet", "6~12대 UAV가 정다각형 대형을\n자율 형성, 3D 다층 레이더 돔 전개", CYAN),
        ("Mesh Radar", "LiDAR + RF 기반 Mesh Network,\n삼각측량(Trilateration)으로\n진입 객체 3D 좌표 산출", GREEN),
        ("Timer Control", "비인가 UAV에 체공 시간 자동 할당\n만료 시 RTH 명령 및\n에스코트 프로토콜 가동", ORANGE),
    ]

    for i, (title, desc, accent) in enumerate(pillars):
        x = Inches(0.8 + i * 4.1)
        y = Inches(3.7)

        # Gradient-ish card
        card = add_rounded_rect(slide, x, y, Inches(3.7), Inches(2.4), BG_CARD, accent)

        # Icon circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(1.3), y + Inches(0.3), Inches(1), Inches(1)
        )
        circle.fill.background()
        circle.line.color.rgb = accent
        circle.line.width = Pt(2)

        add_text_box(slide, x + Inches(0.3), y + Inches(1.4), Inches(3.1), Inches(0.4),
                     title, font_size=18, color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name="Consolas")
        add_text_box(slide, x + Inches(0.3), y + Inches(1.9), Inches(3.1), Inches(0.9),
                     desc, font_size=12, color=DIM, alignment=PP_ALIGN.CENTER)

    add_footer(slide, 3)

    add_notes(slide, (
        "그래서 저희는 '관제탑을 하늘로 띄우자'는 역발상을 했습니다. "
        "수억 원이 드는 고정 인프라 대신, 여러 대의 통제용 군집 드론을 출격시킵니다. "
        "이 군집 드론들이 상공에서 서로 Mesh Network를 형성하여 "
        "실시간 통신망을 구축하고, 거대한 레이더 돔(Dome) 형태의 "
        "통제 공역을 스스로 만들어냅니다."
    ))


def build_slide_4(prs):
    """Core Technology - Sim-to-Real."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_tag_label(slide, Inches(0.8), Inches(0.5), "CORE TECHNOLOGY — SIM-TO-REAL", PURPLE, BG_DARK)

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.0), Inches(0.8), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = CYAN
    line.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(1.0),
                 "가상의 강화학습 알고리즘을\n현실의 드론 관제로 이식하다",
                 font_size=32, color=WHITE, bold=True)

    # Left box: SC2
    sc2_card = add_rounded_rect(slide, Inches(0.6), Inches(2.5), Inches(4.8), Inches(1.8),
                                 BG_CARD, PURPLE)
    add_text_box(slide, Inches(0.8), Inches(2.6), Inches(4.4), Inches(0.3),
                 "STARCRAFT II", font_size=11, color=PURPLE, bold=True, font_name="Consolas")
    add_text_box(slide, Inches(0.8), Inches(3.0), Inches(4.4), Inches(0.4),
                 "Boids + FSM + RL", font_size=14, color=LIGHT, font_name="Consolas")
    add_text_box(slide, Inches(0.8), Inches(3.5), Inches(4.4), Inches(0.4),
                 "10,000+ 게임 검증", font_size=13, color=DIM)

    # Center arrow
    add_text_box(slide, Inches(5.6), Inches(2.9), Inches(1.5), Inches(0.4),
                 "Sim -> Real", font_size=14, color=CYAN, bold=True, font_name="Consolas",
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(5.6), Inches(3.3), Inches(1.5), Inches(0.3),
                 "+ Altitude", font_size=11, color=DIM, alignment=PP_ALIGN.CENTER)

    # Right box: Drone ATC
    atc_card = add_rounded_rect(slide, Inches(7.3), Inches(2.5), Inches(4.8), Inches(1.8),
                                 BG_CARD, CYAN)
    add_text_box(slide, Inches(7.5), Inches(2.6), Inches(4.4), Inches(0.3),
                 "DRONE ATC", font_size=11, color=CYAN, bold=True, font_name="Consolas")
    add_text_box(slide, Inches(7.5), Inches(3.0), Inches(4.4), Inches(0.4),
                 "Formation + Timer + Alert", font_size=14, color=LIGHT, font_name="Consolas")
    add_text_box(slide, Inches(7.5), Inches(3.5), Inches(4.4), Inches(0.4),
                 "3D Mesh Radar Network", font_size=13, color=DIM)

    # Transfer mapping table
    table_data = [
        ("SC2 Component", "Drone ATC Mapping", "전이 방법론"),
        ("Boids Algorithm (군집 이동)", "Formation Flight Control", "1:1 구조 매핑"),
        ("Blackboard (중앙 상태 관리)", "Flight Data Hub", "1:1 구조 매핑"),
        ("Authority Mode (우선순위 체계)", "ATC Priority Levels", "1:1 구조 매핑"),
        ("IntelManager (정보 수집/탐지)", "3D Radar Mesh (Sensor Fusion)", "적응형 파라미터 튜닝"),
        ("RL Agent (강화학습 의사결정)", "Adaptive Path Optimization", "적응형 파라미터 튜닝"),
    ]

    rows, cols = len(table_data), 3
    tbl_shape = slide.shapes.add_table(rows, cols,
                                        Inches(0.6), Inches(4.5),
                                        Inches(11.5), Inches(2.4))
    tbl = tbl_shape.table

    col_widths = [Inches(4.0), Inches(4.5), Inches(3.0)]
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w

    for r, row_data in enumerate(table_data):
        tbl.rows[r].height = Inches(0.38)
        for c, cell_text in enumerate(row_data):
            cell = tbl.cell(r, c)
            cell.text = cell_text
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.font.name = "Malgun Gothic"

            # Cell margins for breathing room
            cell.margin_left = Inches(0.12)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)

            # Header row
            if r == 0:
                p.font.size = Pt(13)
                p.font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x15, 0x1E, 0x35)
                if c == 0:
                    p.font.color.rgb = PURPLE
                elif c == 1:
                    p.font.color.rgb = CYAN
                else:
                    p.font.color.rgb = GOLD
            else:
                p.font.size = Pt(12)
                p.font.color.rgb = LIGHT
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_CARD
                if c == 2 and "구조 매핑" in cell_text:
                    p.font.color.rgb = GREEN
                elif c == 2 and "적응형" in cell_text:
                    p.font.color.rgb = ORANGE

    add_footer(slide, 4)

    add_notes(slide, (
        "이 거대한 시스템을 통제하는 두뇌는 바로 'Sim-to-Real' 기술입니다. "
        "저는 이 알고리즘을 스타크래프트 2의 저그 봇(Zerg Bot) 군집 제어 "
        "강화학습 모델에서 발전시켰습니다. "
        "1만 번 이상의 시뮬레이션으로 검증된 군집의 자율 산개 및 진형 유지 로직을, "
        "3D 공간의 물리적 변수와 안전 제약을 보정하여 "
        "실제 드론의 Flight Controller(비행 제어기)로 직접 이식했습니다. "
        "게임의 AI가 현실의 공역 관제 로직으로 재탄생한 것입니다."
    ))


def build_slide_5(prs):
    """3D Live Simulation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Background: 3D tracking visualization (behind all content)
    add_image(slide, 'slide5_3d_tracking.png',
              Inches(6.5), Inches(0.3), width=Inches(6.5))

    add_tag_label(slide, Inches(0.8), Inches(0.6), "3D LIVE SIMULATION", GREEN, BG_DARK)

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(0.8), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = CYAN
    line.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(1.0),
                 "사각지대 없는 3D 입체 스캔\n및 실시간 좌표 추적",
                 font_size=32, color=WHITE, bold=True)

    # 4 phase cards
    phases = [
        ("PHASE 01", "Mesh 형성", "6대 드론이 육각형 대형으로\n레이더 돔 전개", CYAN),
        ("PHASE 02", "탐지 & 식별", "RF 스캔 -> MAC 추출\n-> DB 등록 -> 허가 검증", GREEN),
        ("PHASE 03", "타이머 & 추적", "삼각측량으로 X,Y,Z\n실시간 위치 + 카운트다운", ORANGE),
        ("PHASE 04", "경고 & 퇴각", "2분 전 경고 ->\n강제 복귀 -> 에스코트", RED),
    ]

    for i, (phase, title, desc, accent) in enumerate(phases):
        x = Inches(0.5 + i * 3.1)
        y = Inches(2.8)

        card = add_rounded_rect(slide, x, y, Inches(2.8), Inches(2.0), BG_CARD, accent)

        add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(2.4), Inches(0.3),
                     phase, font_size=10, color=accent, bold=True, font_name="Consolas")
        add_text_box(slide, x + Inches(0.2), y + Inches(0.55), Inches(2.4), Inches(0.4),
                     title, font_size=18, color=WHITE, bold=True)
        add_text_box(slide, x + Inches(0.2), y + Inches(1.05), Inches(2.4), Inches(0.8),
                     desc, font_size=12, color=DIM)

        # Arrow between phases
        if i < 3:
            arrow_x = x + Inches(2.9)
            add_text_box(slide, arrow_x, y + Inches(0.8), Inches(0.3), Inches(0.3),
                         ">", font_size=18, color=accent, bold=True, font_name="Consolas",
                         alignment=PP_ALIGN.CENTER)

    # Demo links
    demos = [
        ("3D 인터랙티브 관제 시뮬레이터", "Web 기반 Three.js 렌더링 엔진", "Three.js 렌더링 · 3D 궤도 제어 · 실시간 모니터링"),
        ("실시간 Fleet 통제 대시보드", "React 18 아키텍처 기반", "SVG 레이더 시각화 · Fleet 추적 엔진"),
    ]
    for i, (name, file, tech) in enumerate(demos):
        x = Inches(0.8 + i * 5.5)
        y = Inches(5.5)
        card = add_rounded_rect(slide, x, y, Inches(5.0), Inches(0.9), BG_CARD)
        add_text_box(slide, x + Inches(0.3), y + Inches(0.1), Inches(4.4), Inches(0.35),
                     name, font_size=14, color=WHITE, bold=True)
        add_text_box(slide, x + Inches(0.3), y + Inches(0.45), Inches(4.4), Inches(0.2),
                     file, font_size=10, color=DIM, font_name="Consolas")
        add_text_box(slide, x + Inches(0.3), y + Inches(0.65), Inches(4.4), Inches(0.2),
                     tech, font_size=10, color=CYAN_DIM, font_name="Consolas")

    add_footer(slide, 5)

    add_notes(slide, (
        "이 복잡한 로직이 실제 관제 환경에서 어떻게 작동하는지 "
        "Web 기반 3D 인터랙티브 시뮬레이터로 구현했습니다. "
        "보시는 바와 같이 레이더망에 진입한 유저 드론은 "
        "즉시 고유 식별자가 부여되고 3D 좌표가 추적됩니다."
    ))


def build_slide_6(prs):
    """Real-time Control Dashboard."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Background: Timer dashboard gauges (behind status cards)
    add_image(slide, 'slide6_timer_dash.png',
              Inches(7.0), Inches(0.3), width=Inches(6.0))

    add_tag_label(slide, Inches(0.8), Inches(0.6), "REAL-TIME CONTROL DASHBOARD", ORANGE, BG_DARK)

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(0.8), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = CYAN
    line.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(1.0),
                 "타이머 기반 자동화 관제\n및 퇴각 통제",
                 font_size=32, color=WHITE, bold=True)

    # 4 status cards
    statuses = [
        ("정상 비행", "잔여 시간 > 3분", "08:32", "Heartbeat ACK 정상 수신", GREEN),
        ("시간 임박", "잔여 시간 < 2분", "01:26", "Push 알림 자동 전송", ORANGE),
        ("비행 종료", "잔여 시간 = 0", "00:00", "강제 복귀 명령 발령", RED),
        ("긴급 프로토콜", "미응답 30초 초과", "! ! !", "물리적 에스코트\n관제 강제 개입", RGBColor(0xFF, 0x00, 0x44)),
    ]

    for i, (title, cond, timer, action, accent) in enumerate(statuses):
        x = Inches(0.5 + i * 3.1)
        y = Inches(2.8)

        card = add_rounded_rect(slide, x, y, Inches(2.8), Inches(2.8), BG_CARD, accent)

        # Status circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(0.9), y + Inches(0.25), Inches(0.9), Inches(0.9)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = accent
        circle.line.fill.background()

        add_text_box(slide, x + Inches(0.2), y + Inches(1.2), Inches(2.4), Inches(0.3),
                     title, font_size=16, color=accent, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), y + Inches(1.55), Inches(2.4), Inches(0.25),
                     cond, font_size=11, color=DIM, alignment=PP_ALIGN.CENTER)

        add_text_box(slide, x + Inches(0.2), y + Inches(1.85), Inches(2.4), Inches(0.4),
                     timer, font_size=28, color=accent, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name="Consolas")

        add_text_box(slide, x + Inches(0.2), y + Inches(2.3), Inches(2.4), Inches(0.4),
                     action, font_size=10, color=DIM, alignment=PP_ALIGN.CENTER)

    add_footer(slide, 6)

    add_notes(slide, (
        "관제 대시보드는 이를 자동화하여 체공 시간을 관리합니다. "
        "시간이 임박하면 주의 알림이 전송되고, "
        "허가된 시간이 초과되면 붉은색 경고와 함께 "
        "조종자에게 강제 복귀 명령이 발령됩니다. "
        "중앙 서버와 드론 간의 지연 없는 통신(Low Latency)으로 "
        "이 모든 과정이 이루어집니다."
    ))


def build_slide_7(prs):
    """Expected Impact & Applications."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Background: Coverage comparison (behind use case cards)
    add_image(slide, 'slide7_comparison.png',
              Inches(6.5), Inches(2.8), width=Inches(6.5))

    add_tag_label(slide, Inches(0.8), Inches(0.5), "EXPECTED IMPACT & APPLICATIONS", CYAN, BG_DARK)

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.0), Inches(0.8), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = CYAN
    line.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(7), Inches(1.0),
                 "지형에 구애받지 않는\n이동식 관제 인프라",
                 font_size=30, color=WHITE, bold=True)

    # Comparison table
    table_data = [
        ("항목", "기존 방식", "Swarm-Net"),
        ("탐지 범위", "고정 레이더 반경 내", "군집 이동으로 동적 공역 커버"),
        ("배치 시간", "인프라 설치 수 개월", "드론 출격 수 분 내 전개"),
        ("비용 구조", "CAPEX 수억 원 (고정 설비)", "OPEX 중심 — Fleet 유지보수"),
        ("유연성", "고정 위치, 재배치 불가", "실시간 공역 재구성"),
        ("확장성", "추가 기지국 필요", "드론 추가 투입으로 즉시 확장"),
        ("체공 지속성", "해당 없음 (고정 설치)", "Relay Rotation + Tethered 운용"),
    ]

    rows, cols = len(table_data), 3
    tbl_shape = slide.shapes.add_table(rows, cols,
                                        Inches(0.6), Inches(2.4),
                                        Inches(6.2), Inches(3.2))
    tbl = tbl_shape.table

    col_widths = [Inches(1.4), Inches(2.3), Inches(2.5)]
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w

    for r, row_data in enumerate(table_data):
        tbl.rows[r].height = Inches(0.40)
        for c, cell_text in enumerate(row_data):
            cell = tbl.cell(r, c)
            cell.text = cell_text
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.font.name = "Malgun Gothic"

            cell.margin_left = Inches(0.12)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)

            if r == 0:
                p.font.size = Pt(13)
                p.font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x15, 0x1E, 0x35)
                colors_h = [CYAN, RED, GREEN]
                p.font.color.rgb = colors_h[c]
            else:
                p.font.size = Pt(12)
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_CARD
                if c == 1:
                    p.font.color.rgb = RGBColor(0xFF, 0x88, 0x88)
                elif c == 2:
                    p.font.color.rgb = GREEN
                else:
                    p.font.color.rgb = LIGHT

    # Use cases grid (3x2)
    use_cases = [
        ("군사 야전", "신속 방어망 구축"),
        ("불법 드론 차단", "공항/시설 보안"),
        ("재난 현장", "긴급 공역 통제"),
        ("드론 쇼 안전", "충돌 방지"),
        ("UAM 공역", "복층 비행로 관리"),
        ("농업 방제", "구역 진입 제한"),
    ]

    for i, (title, desc) in enumerate(use_cases):
        col = i % 2
        row = i // 2
        x = Inches(7.3 + col * 2.8)
        y = Inches(2.2 + row * 1.5)

        card = add_rounded_rect(slide, x, y, Inches(2.4), Inches(1.2), BG_CARD)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.25), Inches(2.0), Inches(0.35),
                     title, font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.65), Inches(2.0), Inches(0.3),
                     desc, font_size=12, color=DIM, alignment=PP_ALIGN.CENTER)

    add_footer(slide, 7)

    add_notes(slide, (
        "군사 통신이나 야전 작전 구역, 재난 현장처럼 인프라가 없는 곳에서 "
        "Swarm-Net은 이동식 관제탑으로서 압도적인 성능을 발휘합니다.\n\n"
        "물론 군집 드론의 짧은 배터리 타임이 약점으로 지적될 수 있습니다. "
        "이를 극복하기 위해 핵심 통신 노드는 지상에서 직접 전력을 공급받는 "
        "'유선 테더링(Tethered)' 방식으로 전개하고, "
        "주변 노드들은 그룹을 나누어 '교대 비행(Relay)'하는 프로토콜을 도입하여 "
        "24시간 끊김 없는 관제망을 유지하도록 설계했습니다."
    ))


def build_slide_8(prs):
    """Conclusion & Q&A."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Background: Roadmap progress visualization
    add_image(slide, 'slide8_roadmap.png',
              Inches(1.5), Inches(5.0), width=Inches(10.0))

    add_tag_label(slide, Inches(4.5), Inches(0.8), "CONCLUSION & ROADMAP", CYAN, BG_DARK)

    # Center line
    ctr_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5.7), Inches(1.2), Inches(1.3), Pt(3)
    )
    ctr_line.fill.solid()
    ctr_line.fill.fore_color.rgb = CYAN
    ctr_line.line.fill.background()

    add_text_box(slide, Inches(1.5), Inches(1.4), Inches(10), Inches(1.2),
                 "안전하고 체계적인 하늘의 미래,\nSwarm-Net이 만들어갑니다",
                 font_size=32, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Roadmap timeline
    stages = [
        ("STAGE 1", "SC2 시뮬레이션", "10,000+ 게임 검증", "COMPLETED", GREEN),
        ("STAGE 2", "3D 시뮬레이터", "파라미터 2D->3D 적응", "COMPLETED", GREEN),
        ("STAGE 3", "실 드론 테스트", "5대 편대 비행 검증", "NEXT PHASE", ORANGE),
        ("STAGE 4", "도시 스케일 ATC", "100+ 드론 관제 시스템", "VISION", PURPLE),
    ]

    for i, (stage, title, desc, status, accent) in enumerate(stages):
        x = Inches(0.8 + i * 3.0)
        y = Inches(3.2)

        # Circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(0.7), y, Inches(1.0), Inches(1.0)
        )
        circle.fill.background()
        circle.line.color.rgb = accent
        circle.line.width = Pt(2.5)

        # Connecting line
        if i < 3:
            conn = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                x + Inches(1.8), y + Inches(0.45),
                Inches(1.3), Pt(2)
            )
            conn.fill.solid()
            conn.fill.fore_color.rgb = RGBColor(0x33, 0x44, 0x55)
            conn.line.fill.background()

        add_text_box(slide, x, y + Inches(1.15), Inches(2.5), Inches(0.25),
                     stage, font_size=10, color=accent, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name="Consolas")
        add_text_box(slide, x, y + Inches(1.4), Inches(2.5), Inches(0.35),
                     title, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, y + Inches(1.8), Inches(2.5), Inches(0.25),
                     desc, font_size=11, color=DIM, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, y + Inches(2.1), Inches(2.5), Inches(0.25),
                     status, font_size=10, color=accent, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name="Consolas")

    # Thank You
    add_text_box(slide, Inches(2), Inches(5.8), Inches(8.5), Inches(0.6),
                 "THANK YOU", font_size=36, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name="Consolas")
    add_text_box(slide, Inches(2), Inches(6.4), Inches(8.5), Inches(0.3),
                 "Q & A", font_size=14, color=DIM,
                 alignment=PP_ALIGN.CENTER, font_name="Consolas")

    add_footer(slide, 8)

    add_notes(slide, (
        "드론은 이미 배송과 모빌리티(UAM)의 영역으로 진입했습니다. "
        "이제 하늘은 방치된 공간이 아니라, "
        "안전하고 유연하게 통제되는 새로운 인프라여야 합니다.\n\n"
        "가상의 군집 제어 AI를 현실의 공역 관제 솔루션으로 끌어낸 "
        "Swarm-Net이 그 기준을 제시할 것입니다. "
        "경청해 주셔서 감사합니다. 질문 받겠습니다."
    ))


def main():
    print("[1/2] Building slides...")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_slide_1(prs)
    print("  [OK] Slide 1: Title")
    build_slide_2(prs)
    print("  [OK] Slide 2: Problem")
    build_slide_3(prs)
    print("  [OK] Slide 3: Solution")
    build_slide_4(prs)
    print("  [OK] Slide 4: Sim-to-Real")
    build_slide_5(prs)
    print("  [OK] Slide 5: 3D Simulation")
    build_slide_6(prs)
    print("  [OK] Slide 6: Dashboard")
    build_slide_7(prs)
    print("  [OK] Slide 7: Impact & Applications")
    build_slide_8(prs)
    print("  [OK] Slide 8: Conclusion & Q&A")

    print("[2/2] Saving PPTX...")
    output_path = Path(__file__).parent / "Swarm_Net_Presentation.pptx"
    try:
        prs.save(str(output_path))
    except PermissionError:
        output_path = Path(__file__).parent / "Swarm_Net_Presentation_v2.pptx"
        prs.save(str(output_path))
        print(f"  (Original locked, saved as v2)")

    print(f"\n{'='*60}")
    print(f"PPTX CREATED SUCCESSFULLY!")
    print(f"   File: {output_path}")
    print(f"   Slides: 8")
    print(f"   Speaker notes: included")
    print(f"   Theme: Dark Navy + Cyan/Green accents")
    print(f"{'='*60}")


if __name__ == "__main__":
    main()

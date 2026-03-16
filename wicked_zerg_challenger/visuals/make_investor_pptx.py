"""
Gwangju Urban Drone UTM — Investor Pitch Deck (20 Slides)
==========================================================
python-pptx 기반 다크 테마 투자자 프레젠테이션
"""

import os, sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ===== Colors =====
BG       = RGBColor(0x0B, 0x10, 0x1F)
BG_CARD  = RGBColor(0x12, 0x1A, 0x2E)
CYAN     = RGBColor(0x00, 0xFF, 0xCC)
CYAN_DIM = RGBColor(0x00, 0xB8, 0x94)
ORANGE   = RGBColor(0xFF, 0x98, 0x00)
RED      = RGBColor(0xFF, 0x45, 0x45)
GREEN    = RGBColor(0x00, 0xE6, 0x76)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT    = RGBColor(0xE8, 0xE8, 0xE8)
DIM      = RGBColor(0xA0, 0xA8, 0xB8)
PURPLE   = RGBColor(0xB3, 0x88, 0xFF)
GOLD     = RGBColor(0xFF, 0xD7, 0x00)
BLUE     = RGBColor(0x44, 0x88, 0xFF)

FONT = "Malgun Gothic"
MONO = "Consolas"


# ===== Helpers =====
def bg(slide):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = BG

def tb(slide, l, t, w, h, txt, sz=18, c=WHITE, b=False, al=PP_ALIGN.LEFT, fn=FONT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = txt
    p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b
    p.font.name = fn; p.alignment = al
    return box

def bullets(slide, l, t, w, h, items, sz=14, c=LIGHT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            txt, ic, ib = item
        else:
            txt, ic, ib = item, c, False
        p.text = txt; p.font.size = Pt(sz); p.font.color.rgb = ic
        p.font.bold = ib; p.font.name = FONT; p.space_after = Pt(5)

def card(slide, l, t, w, h, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = BG_CARD
    if border: s.line.color.rgb = border; s.line.width = Pt(2)
    else: s.line.fill.background()
    return s

def tag(slide, l, t, txt, bg_c=CYAN, txt_c=BG):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, Inches(3.2), Inches(0.34))
    s.fill.solid(); s.fill.fore_color.rgb = bg_c; s.line.fill.background()
    tf = s.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = txt; p.font.size = Pt(10)
    p.font.color.rgb = txt_c; p.font.bold = True; p.font.name = MONO
    p.alignment = PP_ALIGN.CENTER

def accent_line(slide, l, t):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, Inches(0.8), Pt(3))
    s.fill.solid(); s.fill.fore_color.rgb = CYAN; s.line.fill.background()

def footer(slide, n, total=20):
    tb(slide, Inches(0.4), Inches(6.9), Inches(3), Inches(0.4),
       "GWANGJU DRONE UTM", sz=8, c=CYAN_DIM, b=True, fn=MONO)
    tb(slide, Inches(11), Inches(6.9), Inches(2), Inches(0.4),
       f"{n:02d} / {total:02d}", sz=8, c=DIM, fn=MONO, al=PP_ALIGN.RIGHT)

def notes(slide, txt):
    ns = slide.notes_slide; tf = ns.notes_text_frame; tf.text = txt

def stat_card(slide, x, y, val, label, color=CYAN, w=2.5, h=1.2):
    card(slide, x, y, Inches(w), Inches(h))
    tb(slide, x + Inches(0.15), y + Inches(0.15), Inches(w - 0.3), Inches(0.5),
       val, sz=26, c=color, b=True, fn=MONO, al=PP_ALIGN.CENTER)
    tb(slide, x + Inches(0.15), y + Inches(0.7), Inches(w - 0.3), Inches(0.35),
       label, sz=10, c=DIM, al=PP_ALIGN.CENTER)

def footer_text(slide, txt):
    tb(slide, Inches(0.8), Inches(6.55), Inches(11.5), Inches(0.3),
       txt, sz=9, c=DIM)


# ===================================================================
# SLIDES
# ===================================================================

def slide_01(prs):
    """Title / Cover."""
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    tag(s, Inches(0.8), Inches(1.0), "INVESTOR PITCH DECK 2026", CYAN, BG)
    accent_line(s, Inches(0.8), Inches(1.5))
    tb(s, Inches(0.8), Inches(1.7), Inches(8), Inches(1.2),
       "도시 드론 교통 관리", sz=44, c=WHITE, b=True)
    tb(s, Inches(0.8), Inches(3.1), Inches(8), Inches(0.6),
       "광주광역시 시범 프로젝트", sz=22, c=LIGHT)
    tb(s, Inches(0.8), Inches(3.8), Inches(8), Inches(0.5),
       '"세계 최초 도시 규모 이동식 레이더 망"', sz=16, c=CYAN, b=True)
    tb(s, Inches(0.8), Inches(5.0), Inches(8), Inches(0.4),
       "2026–2036  |  광주광역시, 대한민국", sz=14, c=DIM, fn=MONO)
    # decoration
    for i, (ix, iy, isz, rot) in enumerate([
        (9.5, 1.2, 2.8, 30), (10.2, 2.0, 1.6, 30), (9.0, 3.8, 1.2, 15)
    ]):
        h = s.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(ix), Inches(iy),
                                Inches(isz), Inches(isz))
        h.fill.background()
        h.line.color.rgb = [CYAN_DIM, RGBColor(0x00,0x88,0x66), PURPLE][i]
        h.line.width = Pt(1.5); h.rotation = rot
    footer(s, 1)


def slide_02(prs):
    """Problem."""
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    tag(s, Inches(0.8), Inches(0.5), "THE PROBLEM", RED, BG)
    accent_line(s, Inches(0.8), Inches(1.0))
    tb(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.8),
       "보이지 않는 위협:\n소형 드론", sz=36, c=WHITE, b=True)

    # 3 comparison cards
    items = [
        ("지상 교통", "신호등, GPS, 카메라로 제어", GREEN),
        ("항공", "레이더, 관제탑, IFR/VFR 규정", BLUE),
        ("드론", "제어 시스템 없음", RED),
    ]
    for i, (title, desc, ac) in enumerate(items):
        x = Inches(0.8 + i * 4.1); y = Inches(2.5)
        card(s, x, y, Inches(3.7), Inches(1.5), ac)
        tb(s, x + Inches(0.3), y + Inches(0.2), Inches(3.1), Inches(0.4),
           title, sz=18, c=ac, b=True)
        tb(s, x + Inches(0.3), y + Inches(0.7), Inches(3.1), Inches(0.6),
           desc, sz=13, c=LIGHT)

    # key issues
    issues = [
        ("문제", "소형 드론 (<5kg)은 레이더에 보이지 않음", RED),
        ("위험", "충돌, 사생활 침해, 불법 활동", ORANGE),
        ("공백", "10만대+ 소비자 드론에 대한 실시간 추적 없음", GOLD),
    ]
    for i, (lbl, desc, ac) in enumerate(issues):
        y = Inches(4.3 + i * 0.55)
        tb(s, Inches(1.0), y, Inches(1.2), Inches(0.4), lbl, sz=13, c=ac, b=True, fn=MONO)
        tb(s, Inches(2.5), y, Inches(9), Inches(0.4), desc, sz=13, c=LIGHT)

    footer_text(s, "공역 사고의 80%가 탐지되지 않은 소형 드론 관련")
    footer(s, 2)


def slide_03(prs):
    """Solution."""
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    tag(s, Inches(0.8), Inches(0.5), "OUR SOLUTION", GREEN, BG)
    accent_line(s, Inches(0.8), Inches(1.0))
    tb(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.8),
       "이동식 레이더 망:\n하늘의 눈", sz=36, c=WHITE, b=True)
    tb(s, Inches(0.8), Inches(2.3), Inches(10), Inches(0.4),
       "혁신: 100대 레이더 장착 드론이 광주 24/7 순찰", sz=16, c=CYAN, b=True)

    # 4 stat cards
    stats = [
        ("501 km²", "커버리지\n(전체 광역시)", CYAN),
        ("99.9%", "탐지 정확도\n(100g 이상)", GREEN),
        ("<5초", "실시간 경보\n(당국 전송)", ORANGE),
        ("AI 기반", "군집 지능\n(저그봇 프로)", PURPLE),
    ]
    for i, (val, lbl, ac) in enumerate(stats):
        x = Inches(0.6 + i * 3.1)
        stat_card(s, x, Inches(3.0), val, lbl, ac, w=2.8, h=1.4)

    card(s, Inches(0.8), Inches(4.8), Inches(11.5), Inches(1.0), GOLD)
    tb(s, Inches(1.1), Inches(4.9), Inches(11), Inches(0.35),
       "세계 최초", sz=16, c=GOLD, b=True)
    tb(s, Inches(1.1), Inches(5.3), Inches(11), Inches(0.35),
       "어떤 도시도 규모별 이동식 레이더 망을 배치한 적 없음", sz=13, c=LIGHT)

    footer_text(s, "가상 군집에서 실제 세계 안전으로")
    footer(s, 3)


def slide_04(prs):
    """Tech Foundation."""
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    tag(s, Inches(0.8), Inches(0.5), "TECHNOLOGY FOUNDATION", PURPLE, BG)
    accent_line(s, Inches(0.8), Inches(1.0))
    tb(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.8),
       "검증된 게임 AI 기반 구축", sz=36, c=WHITE, b=True)

    # SC2 card
    card(s, Inches(0.6), Inches(2.3), Inches(5.5), Inches(3.5), PURPLE)
    items = [
        ("출처", "군집 저그봇 프로 (SC2 AI, 75% 승률)"),
        ("검증", "200+ 유닛 실시간 군집 제어 (<10ms)"),
        ("알고리즘", "충돌 회피, 순찰 라우팅, 적응형 학습"),
        ("전환", "가상 → 시뮬레이션 → 실제 하드웨어"),
    ]
    for i, (lbl, desc) in enumerate(items):
        y = Inches(2.5 + i * 0.75)
        tb(s, Inches(0.9), y, Inches(1.5), Inches(0.4), lbl, sz=12, c=PURPLE, b=True, fn=MONO)
        tb(s, Inches(2.5), y, Inches(3.3), Inches(0.6), desc, sz=12, c=LIGHT)

    # Why + ranking card
    card(s, Inches(6.5), Inches(2.3), Inches(5.8), Inches(1.5), CYAN)
    tb(s, Inches(6.8), Inches(2.5), Inches(5.2), Inches(0.4),
       "왜 게이밍 AI인가?", sz=16, c=CYAN, b=True)
    tb(s, Inches(6.8), Inches(3.0), Inches(5.2), Inches(0.6),
       "복잡한 다중 에이전트 시스템을 검증하는 가장 빠른 방법", sz=13, c=LIGHT)

    card(s, Inches(6.5), Inches(4.1), Inches(5.8), Inches(1.5), GOLD)
    tb(s, Inches(6.8), Inches(4.3), Inches(5.2), Inches(0.4),
       "AI Arena 실적", sz=16, c=GOLD, b=True)
    tb(s, Inches(6.8), Inches(4.8), Inches(5.2), Inches(0.6),
       "상위 10% 순위 (다이아-마스터급)", sz=13, c=LIGHT)

    footer_text(s, "검증된 알고리즘, 실제 배치를 위해 100배 확장")
    footer(s, 4)


def slide_05(prs):
    """System Architecture."""
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    tag(s, Inches(0.8), Inches(0.5), "SYSTEM ARCHITECTURE", CYAN, BG)
    accent_line(s, Inches(0.8), Inches(1.0))
    tb(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
       "3계층 아키텍처", sz=36, c=WHITE, b=True)

    layers = [
        ("Layer 1", "순찰 드론 (100대)", [
            "레이더 장착 쿼드콥터 (20kg 페이로드)",
            "4시간 비행, 자동 재충전 스테이션",
            "메시 네트워크 (5G/6G 백업)",
        ], CYAN),
        ("Layer 2", "관제 센터 (광주시청)", [
            "AI 조정 엔진 (군집 지능)",
            "실시간 3D 디지털 트윈",
            "경찰·소방·항공 당국 통합",
        ], GREEN),
        ("Layer 3", "엣지 컴퓨팅", [
            "드론 탑재 AI (위협 감지, 경로 계획)",
            "클라우드 백업 (학습, 분석, 규정 준수)",
            "",
        ], ORANGE),
    ]
    for i, (layer, title, descs, ac) in enumerate(layers):
        x = Inches(0.6 + i * 4.1); y = Inches(2.0)
        card(s, x, y, Inches(3.8), Inches(3.8), ac)
        tb(s, x + Inches(0.2), y + Inches(0.2), Inches(3.4), Inches(0.3),
           layer, sz=11, c=ac, b=True, fn=MONO)
        tb(s, x + Inches(0.2), y + Inches(0.55), Inches(3.4), Inches(0.4),
           title, sz=16, c=WHITE, b=True)
        for j, d in enumerate(descs):
            if d:
                tb(s, x + Inches(0.3), y + Inches(1.2 + j * 0.6), Inches(3.2), Inches(0.5),
                   f"• {d}", sz=11, c=LIGHT)

    footer_text(s, "분산형 + 중앙 집중식 = 중복성 + 효율성")
    footer(s, 5)


def slide_06(prs):
    """10-Year Roadmap."""
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    tag(s, Inches(0.8), Inches(0.5), "10-YEAR ROADMAP", GOLD, BG)
    accent_line(s, Inches(0.8), Inches(1.0))
    tb(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
       "10대 드론에서 40,000대까지", sz=34, c=WHITE, b=True)

    phases = [
        ("Phase 1", "2026-2027", "SC2 AI 검증\nAI Arena Top 5", GREEN),
        ("Phase 2", "2027-2028", "시뮬레이션\n1,000대 Unity/ROS2", GREEN),
        ("Phase 3", "2028-2030", "실제 시범\n10→100대 (GIST)", ORANGE),
        ("Phase 4", "2030-2032", "광주 배치\n100대 도시 전역", CYAN),
        ("Phase 5", "2032-2034", "상용화\n1,000대, 배송 통합", BLUE),
        ("Phase 6", "2034-2036", "전국 규모\n20K-40K대", PURPLE),
    ]
    for i, (ph, yr, desc, ac) in enumerate(phases):
        col = i % 3; row = i // 3
        x = Inches(0.6 + col * 4.1); y = Inches(2.0 + row * 2.3)
        card(s, x, y, Inches(3.8), Inches(1.9), ac)
        tb(s, x + Inches(0.2), y + Inches(0.15), Inches(1.5), Inches(0.3),
           ph, sz=11, c=ac, b=True, fn=MONO)
        tb(s, x + Inches(2.0), y + Inches(0.15), Inches(1.5), Inches(0.3),
           yr, sz=10, c=DIM, fn=MONO, al=PP_ALIGN.RIGHT)
        tb(s, x + Inches(0.2), y + Inches(0.6), Inches(3.4), Inches(1.0),
           desc, sz=13, c=LIGHT)

    footer_text(s, "총 투자: 10년간 1,400억원")
    footer(s, 6)


def slide_07(prs):
    """Phase 1-2 Detail."""
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    tag(s, Inches(0.8), Inches(0.5), "PHASE 1-2 DETAIL", GREEN, BG)
    accent_line(s, Inches(0.8), Inches(1.0))
    tb(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.8),
       "Phase 1-2: 검증 및 시뮬레이션\n(2026–2028)", sz=32, c=WHITE, b=True)

    stats = [
        ("Top 5", "SC2 AI 순위", GOLD),
        ("Unity+ROS2", "시뮬레이터 구축", CYAN),
        ("1,000대", "가상 테스트", GREEN),
        ("6개", "핵심 알고리즘", PURPLE),
        ("22억원", "투자 규모", ORANGE),
        ("GIST", "파트너", BLUE),
    ]
    for i, (val, lbl, ac) in enumerate(stats):
        col = i % 3; row = i // 3
        x = Inches(0.6 + col * 4.1); y = Inches(2.8 + row * 1.8)
        stat_card(s, x, y, val, lbl, ac, w=3.8, h=1.4)

    footer_text(s, "하드웨어 투자 전 알고리즘 검증")
    footer(s, 7)


def slide_08(prs):
    """Phase 3 Detail."""
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    tag(s, Inches(0.8), Inches(0.5), "PHASE 3 DETAIL", ORANGE, BG)
    accent_line(s, Inches(0.8), Inches(1.0))
    tb(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.8),
       "Phase 3: 실제 하드웨어 테스팅\n(2028–2030)", sz=32, c=WHITE, b=True)

    items = [
        ("시작", "10대 소비자 쿼드콥터 (DJI Matrice 300)"),
        ("확장", "100대 맞춤형 레이더 드론 (한국 방산업체 협력)"),
        ("위치", "GIST 캠퍼스 (1 km²) → 광주 과학공원 (10 km²)"),
        ("테스트", "충돌 회피, 순찰 라우팅, 응급 대응"),
        ("규제", "UTM 인증을 위해 국토부와 협력"),
    ]
    for i, (lbl, desc) in enumerate(items):
        y = Inches(2.5 + i * 0.7)
        card(s, Inches(0.8), y, Inches(11.5), Inches(0.55))
        tb(s, Inches(1.1), y + Inches(0.08), Inches(1.5), Inches(0.4),
           lbl, sz=13, c=ORANGE, b=True, fn=MONO)
        tb(s, Inches(2.8), y + Inches(0.08), Inches(9), Inches(0.4),
           desc, sz=13, c=LIGHT)

    # Milestone badge
    card(s, Inches(3.0), Inches(6.0), Inches(7), Inches(0.45), GOLD)
    tb(s, Inches(3.2), Inches(6.03), Inches(6.5), Inches(0.35),
       "마일스톤: 한국 최초 100대 드론 야외 테스트", sz=13, c=GOLD, b=True, al=PP_ALIGN.CENTER)

    footer_text(s, "투자: 55억원 | 기간: 24개월")
    footer(s, 8)


def slide_09(prs):
    """Phase 4 Detail."""
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    tag(s, Inches(0.8), Inches(0.5), "PHASE 4 DETAIL", CYAN, BG)
    accent_line(s, Inches(0.8), Inches(1.0))
    tb(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.8),
       "Phase 4: 도시 규모 배치\n(2030–2032)", sz=32, c=WHITE, b=True)

    # Top stats
    top_stats = [
        ("100대", "레이더 드론\n24/7 순찰", CYAN),
        ("10개", "충전 스테이션\n전략적 위치", GREEN),
        ("99%", "광역시 커버리지\n가동시간 95%", GOLD),
    ]
    for i, (val, lbl, ac) in enumerate(top_stats):
        stat_card(s, Inches(0.6 + i * 4.1), Inches(2.5), val, lbl, ac, w=3.8, h=1.4)

    # Applications
    tb(s, Inches(0.8), Inches(4.3), Inches(5), Inches(0.4),
       "응용 분야", sz=16, c=WHITE, b=True)
    apps = [
        ("불법 드론 탐지", "제한 구역: 공항, 군사", RED),
        ("교통 모니터링", "고속도로 사고, 정체", ORANGE),
        ("응급 대응", "수색 및 구조, 재해 평가", GREEN),
    ]
    for i, (title, desc, ac) in enumerate(apps):
        x = Inches(0.6 + i * 4.1); y = Inches(4.8)
        card(s, x, y, Inches(3.8), Inches(1.1), ac)
        tb(s, x + Inches(0.2), y + Inches(0.15), Inches(3.4), Inches(0.35),
           title, sz=14, c=ac, b=True)
        tb(s, x + Inches(0.2), y + Inches(0.55), Inches(3.4), Inches(0.35),
           desc, sz=11, c=DIM)

    footer_text(s, "투자: 220억원 | 세계 최초 도시 규모 이동식 레이더 망")
    footer(s, 9)


def slide_10(prs):
    """Phase 5-6."""
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    tag(s, Inches(0.8), Inches(0.5), "PHASE 5-6: SCALE UP", PURPLE, BG)
    accent_line(s, Inches(0.8), Inches(1.0))
    tb(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
       "Phase 5-6: 상용화 및 전국 규모", sz=34, c=WHITE, b=True)

    stats = [
        ("20,000+", "배치 드론", CYAN),
        ("10+", "커버 도시", GREEN),
        ("500K+", "일일 비행", BLUE),
        ("5,000+", "창출 일자리", ORANGE),
        ("5,500억", "연 수익", GOLD),
        ("1,100억", "투자 규모", PURPLE),
    ]
    for i, (val, lbl, ac) in enumerate(stats):
        col = i % 3; row = i // 3
        x = Inches(0.6 + col * 4.1); y = Inches(2.0 + row * 2.0)
        stat_card(s, x, y, val, lbl, ac, w=3.8, h=1.6)

    footer_text(s, "2032-2034: 광주 확장 | 2034-2036: 서울, 부산, 인천")
    footer(s, 10)


def slide_11(prs):
    """Market Opportunity."""
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    tag(s, Inches(0.8), Inches(0.5), "MARKET OPPORTUNITY", GOLD, BG)
    accent_line(s, Inches(0.8), Inches(1.0))
    tb(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
       "도시 드론 시장: 2035년까지 990억 달러", sz=32, c=WHITE, b=True)

    segments = [
        ("드론 배송", "290억$", "전자상거래, 음식, 의료", CYAN),
        ("응급 서비스", "150억$", "앰뷸런스 드론, 소방", RED),
        ("인프라", "120억$", "점검, 유지보수, 측량", GREEN),
        ("UTM 시스템", "180억$", "교통 관리, 규정 준수", BLUE),
        ("국방", "160억$", "국경 순찰, 감시", PURPLE),
    ]
    for i, (seg, val, desc, ac) in enumerate(segments):
        y = Inches(2.0 + i * 0.8)
        card(s, Inches(0.8), y, Inches(6.0), Inches(0.65))
        tb(s, Inches(1.0), y + Inches(0.1), Inches(2.0), Inches(0.35),
           seg, sz=13, c=ac, b=True)
        tb(s, Inches(3.2), y + Inches(0.1), Inches(1.3), Inches(0.35),
           val, sz=13, c=GOLD, b=True, fn=MONO)
        tb(s, Inches(4.6), y + Inches(0.1), Inches(2.0), Inches(0.35),
           desc, sz=11, c=DIM)

    # Korea market card
    card(s, Inches(7.3), Inches(2.0), Inches(5.0), Inches(2.0), GOLD)
    tb(s, Inches(7.6), Inches(2.2), Inches(4.4), Inches(0.4),
       "한국 시장", sz=18, c=GOLD, b=True)
    tb(s, Inches(7.6), Inches(2.8), Inches(4.4), Inches(0.5),
       "2035년까지 55조원\n(전 세계 5% 점유율)", sz=14, c=LIGHT)

    card(s, Inches(7.3), Inches(4.3), Inches(5.0), Inches(1.5), CYAN)
    tb(s, Inches(7.6), Inches(4.5), Inches(4.4), Inches(0.4),
       "광주 기회", sz=18, c=CYAN, b=True)
    tb(s, Inches(7.6), Inches(5.0), Inches(4.4), Inches(0.5),
       "선점 우위 + 정부 지원", sz=14, c=LIGHT)

    footer_text(s, "목표: 한국 시장 10% 점유율 = 연 5,500억원 수익 (2035)")
    footer(s, 11)


def slide_12(prs):
    """Revenue Model."""
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    tag(s, Inches(0.8), Inches(0.5), "REVENUE MODEL", GREEN, BG)
    accent_line(s, Inches(0.8), Inches(1.0))
    tb(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
       "다중 수익원", sz=36, c=WHITE, b=True)

    revenues = [
        ("1", "정부 계약", "연 2,200억원", "UTM 운영, 도시 라이선싱", CYAN),
        ("2", "배송 통합", "연 1,650억원", "쿠팡, 배민 협력", GREEN),
        ("3", "응급 서비스", "연 880억원", "경찰, 소방, 의료", ORANGE),
        ("4", "데이터 서비스", "연 550억원", "교통 분석, 도시 계획", BLUE),
        ("5", "기술 라이선싱", "연 220억원", "타 도시 수출", PURPLE),
    ]
    for i, (num, title, amount, desc, ac) in enumerate(revenues):
        y = Inches(2.0 + i * 0.85)
        card(s, Inches(0.8), y, Inches(11.5), Inches(0.7), ac)
        # number circle
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y + Inches(0.1),
                                   Inches(0.45), Inches(0.45))
        circ.fill.solid(); circ.fill.fore_color.rgb = ac; circ.line.fill.background()
        tb(s, Inches(1.02), y + Inches(0.12), Inches(0.45), Inches(0.4),
           num, sz=14, c=BG, b=True, al=PP_ALIGN.CENTER, fn=MONO)
        tb(s, Inches(1.7), y + Inches(0.12), Inches(2.5), Inches(0.4),
           title, sz=14, c=WHITE, b=True)
        tb(s, Inches(4.5), y + Inches(0.12), Inches(2.5), Inches(0.4),
           amount, sz=14, c=ac, b=True, fn=MONO)
        tb(s, Inches(7.3), y + Inches(0.12), Inches(4.5), Inches(0.4),
           desc, sz=12, c=DIM)

    # Total
    card(s, Inches(3.0), Inches(6.3), Inches(7.3), Inches(0.5), GOLD)
    tb(s, Inches(3.2), Inches(6.33), Inches(7), Inches(0.4),
       "10년차 총 수익: 연 5,500억원  |  수익 마진: 40%", sz=13, c=GOLD, b=True, al=PP_ALIGN.CENTER)

    footer_text(s, "ROI: 10년간 3.9배 (1,400억원 투자 → 연 5,500억원 수익)")
    footer(s, 12)


def slide_13(prs):
    """Partnerships."""
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    tag(s, Inches(0.8), Inches(0.5), "STRATEGIC PARTNERSHIPS", BLUE, BG)
    accent_line(s, Inches(0.8), Inches(1.0))
    tb(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
       "전략적 파트너십", sz=36, c=WHITE, b=True)

    partners = [
        ("학술", "GIST (AI 연구), 전남대 (항공우주)", CYAN),
        ("정부", "광주시, 국토부, 한국항공우주연구원", GREEN),
        ("산업", "대한항공, 한화 (방산), LG U+ (5G)", ORANGE),
        ("물류", "쿠팡, 배민 (배송 통합)", BLUE),
        ("국제", "DJI (하드웨어), Airbus UTM (기술 교류)", PURPLE),
    ]
    for i, (cat, desc, ac) in enumerate(partners):
        y = Inches(2.0 + i * 0.85)
        card(s, Inches(0.8), y, Inches(11.5), Inches(0.7), ac)
        tb(s, Inches(1.1), y + Inches(0.12), Inches(2.0), Inches(0.4),
           cat, sz=14, c=ac, b=True)
        tb(s, Inches(3.3), y + Inches(0.12), Inches(8.5), Inches(0.4),
           desc, sz=13, c=LIGHT)

    card(s, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.5))
    tb(s, Inches(1.1), Inches(6.25), Inches(5), Inches(0.35),
       "MOU 현황: GIST (서명), 광주시 (진행 중)", sz=12, c=DIM)

    footer_text(s, "민관 협력 모델 (공공 60%, 민간 40% 자금)")
    footer(s, 13)


def slide_14(prs):
    """Regulatory Strategy."""
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    tag(s, Inches(0.8), Inches(0.5), "REGULATORY STRATEGY", ORANGE, BG)
    accent_line(s, Inches(0.8), Inches(1.0))
    tb(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
       "K-드론 규제 탐색", sz=36, c=WHITE, b=True)

    # Current / Challenge / Opportunity
    sections = [
        ("현재", "250g 이상 드론 등록 필요\n공항 주변 제한 구역", CYAN),
        ("과제", "군집에 대한 UTM 프레임워크 없음\n(조종자당 최대 1대)", RED),
        ("기회", "UTM 샌드박스로서의 광주\n(국토부 규제 면제)", GREEN),
    ]
    for i, (lbl, desc, ac) in enumerate(sections):
        x = Inches(0.6 + i * 4.1); y = Inches(2.0)
        card(s, x, y, Inches(3.8), Inches(1.6), ac)
        tb(s, x + Inches(0.2), y + Inches(0.15), Inches(3.4), Inches(0.35),
           lbl, sz=14, c=ac, b=True)
        tb(s, x + Inches(0.2), y + Inches(0.55), Inches(3.4), Inches(0.8),
           desc, sz=12, c=LIGHT)

    # Approach timeline
    tb(s, Inches(0.8), Inches(4.0), Inches(5), Inches(0.4),
       "접근 방식", sz=16, c=WHITE, b=True)
    steps = [
        ("Phase 3", "규제 샌드박스 신청 (2027-2028)", ORANGE),
        ("Phase 4", "시범 프로그램 승인 (2029)", CYAN),
        ("Phase 5", "국가 UTM 표준 (우리 시스템 기반, 2032)", GREEN),
    ]
    for i, (ph, desc, ac) in enumerate(steps):
        y = Inches(4.5 + i * 0.6)
        tb(s, Inches(1.0), y, Inches(1.5), Inches(0.4), ph, sz=12, c=ac, b=True, fn=MONO)
        tb(s, Inches(2.8), y, Inches(9), Inches(0.4), desc, sz=12, c=LIGHT)

    footer_text(s, "목표: 성공적인 시범을 통해 국가 드론 규정 형성")
    footer(s, 14)


def slide_15(prs):
    """Competitive Landscape."""
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    tag(s, Inches(0.8), Inches(0.5), "COMPETITIVE LANDSCAPE", RED, BG)
    accent_line(s, Inches(0.8), Inches(1.0))
    tb(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
       "글로벌 UTM 플레이어", sz=36, c=WHITE, b=True)

    # Competitor table
    data = [
        ("지역", "플레이어", "초점"),
        ("미국", "Amazon, Wing (Alphabet)", "배송 중심"),
        ("중국", "EHang, DJI", "하드웨어 중심"),
        ("유럽", "Airbus UTM, Skyways", "에어택시 중심"),
        ("싱가포르", "Volocopter, Garuda", "소규모 시범"),
    ]
    rows, cols = len(data), 3
    tbl_s = s.shapes.add_table(rows, cols, Inches(0.6), Inches(2.0), Inches(6.0), Inches(2.2))
    tbl = tbl_s.table
    col_w = [Inches(1.3), Inches(2.8), Inches(1.9)]
    for i, w in enumerate(col_w): tbl.columns[i].width = w
    for r, rd in enumerate(data):
        tbl.rows[r].height = Inches(0.40)
        for c, ct in enumerate(rd):
            cell = tbl.cell(r, c); cell.text = ct
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]; p.font.name = FONT
            cell.margin_left = Inches(0.1); cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            if r == 0:
                p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = CYAN
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x15,0x1E,0x35)
            else:
                p.font.size = Pt(11); p.font.color.rgb = LIGHT
                cell.fill.solid(); cell.fill.fore_color.rgb = BG_CARD

    # Our advantage
    tb(s, Inches(7.0), Inches(2.0), Inches(5.5), Inches(0.4),
       "우리 우위", sz=16, c=GREEN, b=True)
    advantages = [
        "이동식 레이더 망 (고정 인프라 대비)",
        "AI 우선 → 벤더 비종속",
        "도시 규모 입증 (소규모 시범 대비)",
        "정부 협력 (민간 전용 대비)",
    ]
    for i, adv in enumerate(advantages):
        y = Inches(2.6 + i * 0.55)
        card(s, Inches(7.0), y, Inches(5.5), Inches(0.45), GREEN)
        tb(s, Inches(7.2), y + Inches(0.05), Inches(5.1), Inches(0.35),
           f"✦ {adv}", sz=12, c=LIGHT)

    footer_text(s, "포지셔닝: 한국 국가 UTM 표준, 2040년까지 50+ 도시 수출")
    footer(s, 15)


def slide_16(prs):
    """Risk Mitigation."""
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    tag(s, Inches(0.8), Inches(0.5), "RISK MITIGATION", ORANGE, BG)
    accent_line(s, Inches(0.8), Inches(1.0))
    tb(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
       "기술 및 운영 위험", sz=36, c=WHITE, b=True)

    risks = [
        ("배터리 수명 (<4시간)", "10개 충전 스테이션, 핫스왑", RED),
        ("날씨 (비, 바람)", "내후성 하드웨어, 자동 착륙", ORANGE),
        ("통신 두절", "메시 네트워크, 자율 폴백", GOLD),
        ("규제 지연", "샌드박스 협력, 단계별 출시", BLUE),
        ("대중 수용", "투명성, 개인정보 보호 장치", PURPLE),
    ]
    for i, (risk, mitigation, ac) in enumerate(risks):
        y = Inches(2.0 + i * 0.85)
        # Risk
        card(s, Inches(0.6), y, Inches(5.2), Inches(0.7), ac)
        tb(s, Inches(0.8), y + Inches(0.12), Inches(0.8), Inches(0.4),
           "위험", sz=10, c=ac, b=True, fn=MONO)
        tb(s, Inches(1.7), y + Inches(0.12), Inches(3.8), Inches(0.4),
           risk, sz=12, c=LIGHT)
        # Arrow
        tb(s, Inches(5.9), y + Inches(0.1), Inches(0.5), Inches(0.4),
           "→", sz=16, c=ac, b=True, fn=MONO, al=PP_ALIGN.CENTER)
        # Mitigation
        card(s, Inches(6.5), y, Inches(5.8), Inches(0.7), GREEN)
        tb(s, Inches(6.7), y + Inches(0.12), Inches(0.8), Inches(0.4),
           "완화", sz=10, c=GREEN, b=True, fn=MONO)
        tb(s, Inches(7.6), y + Inches(0.12), Inches(4.5), Inches(0.4),
           mitigation, sz=12, c=LIGHT)

    tb(s, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.3),
       "보험: $50M 책임 보장 (Phase 4+)  |  중복성: 3개 백업 시스템 (통신, 내비게이션, 전력)",
       sz=10, c=DIM)

    footer_text(s, "안전 우선: 99.99% 가동 시간 목표, 제로 사고 목표")
    footer(s, 16)


def slide_17(prs):
    """Social Impact."""
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    tag(s, Inches(0.8), Inches(0.5), "SOCIAL IMPACT", GREEN, BG)
    accent_line(s, Inches(0.8), Inches(1.0))
    tb(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
       "기술을 넘어: 삶의 개선", sz=36, c=WHITE, b=True)

    impacts = [
        ("안전", "공항 근처 불법 드론 방지 (연 200+ 사고)", RED),
        ("속도", "15분 응급 대응 (구급차 30분 대비)", ORANGE),
        ("환경", "배송 트럭 CO2 30% 절감", GREEN),
        ("일자리", "5,000개 신규 일자리 창출", BLUE),
        ("혁신", "광주를 한국의 드론 수도로 포지셔닝", PURPLE),
    ]
    for i, (lbl, desc, ac) in enumerate(impacts):
        x = Inches(0.6); y = Inches(2.0 + i * 0.7)
        card(s, x, y, Inches(6.0), Inches(0.55), ac)
        tb(s, x + Inches(0.2), y + Inches(0.08), Inches(1.2), Inches(0.35),
           lbl, sz=13, c=ac, b=True)
        tb(s, x + Inches(1.6), y + Inches(0.08), Inches(4.0), Inches(0.35),
           desc, sz=12, c=LIGHT)

    # Community
    tb(s, Inches(7.0), Inches(2.0), Inches(5.5), Inches(0.4),
       "커뮤니티 참여", sz=16, c=WHITE, b=True)
    community = [
        ("공교육", "드론 안전 워크샵"),
        ("학생 프로그램", "GIST 협력, STEM"),
        ("투명성", "오픈 데이터 포털"),
    ]
    for i, (title, desc) in enumerate(community):
        y = Inches(2.6 + i * 1.1)
        card(s, Inches(7.0), y, Inches(5.5), Inches(0.9), CYAN)
        tb(s, Inches(7.3), y + Inches(0.1), Inches(5.0), Inches(0.35),
           title, sz=14, c=CYAN, b=True)
        tb(s, Inches(7.3), y + Inches(0.5), Inches(5.0), Inches(0.3),
           desc, sz=12, c=DIM)

    footer_text(s, "사람을 위한 기술, 이익만을 위한 것이 아님")
    footer(s, 17)


def slide_18(prs):
    """Investment Ask."""
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    tag(s, Inches(0.8), Inches(0.5), "INVESTMENT ASK", GOLD, BG)
    accent_line(s, Inches(0.8), Inches(1.0))
    tb(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
       "시드 라운드: 55억원", sz=36, c=WHITE, b=True)
    tb(s, Inches(0.8), Inches(1.7), Inches(10), Inches(0.4),
       "Phase 1-2 자금 조달", sz=18, c=DIM)

    stats = [
        ("55억원", "요청 금액", GOLD),
        ("15%", "지분", CYAN),
        ("367억원", "평가액", GREEN),
        ("50%", "엔지니어링", PURPLE),
        ("30%", "하드웨어", ORANGE),
        ("20%", "파트너십", BLUE),
    ]
    for i, (val, lbl, ac) in enumerate(stats):
        col = i % 3; row = i // 3
        x = Inches(0.6 + col * 4.1); y = Inches(2.5 + row * 2.0)
        stat_card(s, x, y, val, lbl, ac, w=3.8, h=1.6)

    footer_text(s, "마일스톤: SC2 Top 5 (2027) | 1K 드론 시뮬 (2028) | 100대 테스트 (2030)")
    footer(s, 18)


def slide_19(prs):
    """Financial Projections."""
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    tag(s, Inches(0.8), Inches(0.5), "FINANCIAL PROJECTIONS", GREEN, BG)
    accent_line(s, Inches(0.8), Inches(1.0))
    tb(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
       "투자 수익", sz=36, c=WHITE, b=True)

    # Seed
    card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(3.0), GOLD)
    tb(s, Inches(0.9), Inches(2.2), Inches(5.2), Inches(0.4),
       "시드 (55억원 @ 15% 지분)", sz=16, c=GOLD, b=True)
    seed_items = [
        ("엑시트 평가액 (2036)", "5,500억원"),
        ("투자자 수익", "825억원 (15배)"),
        ("IRR", "연 31%"),
    ]
    for i, (lbl, val) in enumerate(seed_items):
        y = Inches(2.8 + i * 0.6)
        tb(s, Inches(0.9), y, Inches(2.8), Inches(0.4), lbl, sz=12, c=DIM)
        tb(s, Inches(3.8), y, Inches(2.3), Inches(0.4), val, sz=14, c=WHITE, b=True, fn=MONO)

    # Series A
    card(s, Inches(6.8), Inches(2.0), Inches(5.5), Inches(3.0), CYAN)
    tb(s, Inches(7.1), Inches(2.2), Inches(5.0), Inches(0.4),
       "시리즈 A (2030, 220억원 @ 10%)", sz=16, c=CYAN, b=True)
    series_items = [
        ("엑시트 평가액", "5,500억원"),
        ("투자자 수익", "550억원 (2.5배)"),
        ("기간", "6년"),
    ]
    for i, (lbl, val) in enumerate(series_items):
        y = Inches(2.8 + i * 0.6)
        tb(s, Inches(7.1), y, Inches(2.8), Inches(0.4), lbl, sz=12, c=DIM)
        tb(s, Inches(10.0), y, Inches(2.3), Inches(0.4), val, sz=14, c=WHITE, b=True, fn=MONO)

    # Exit strategy
    card(s, Inches(0.6), Inches(5.3), Inches(11.7), Inches(0.7), PURPLE)
    tb(s, Inches(0.9), Inches(5.4), Inches(11), Inches(0.4),
       "엑시트 전략: IPO (2035) 또는 방산/항공 대기업 인수", sz=14, c=PURPLE, b=True)

    footer_text(s, "보수적 예측, 검증된 기술, 정부 지원")
    footer(s, 19)


def slide_20(prs):
    """Thank You / CTA."""
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    tag(s, Inches(4.5), Inches(0.8), "THANK YOU", CYAN, BG)

    # Center accent
    cl = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.7), Inches(1.2), Inches(1.3), Pt(3))
    cl.fill.solid(); cl.fill.fore_color.rgb = CYAN; cl.line.fill.background()

    tb(s, Inches(1.5), Inches(1.5), Inches(10), Inches(0.8),
       "하늘 고속도로를\n함께 만들어요", sz=36, c=WHITE, b=True, al=PP_ALIGN.CENTER)

    # Info cards
    info = [
        ("프로젝트", "광주 드론 UTM"),
        ("책임자", "선우 (군집 저그봇 프로 창립자)"),
        ("이메일", "sunwoo@gwangjudrone.kr"),
        ("위치", "광주광역시, 대한민국"),
        ("기간", "2026–2036 (10년)"),
    ]
    for i, (lbl, val) in enumerate(info):
        y = Inches(2.8 + i * 0.55)
        tb(s, Inches(3.5), y, Inches(2.5), Inches(0.4), lbl, sz=12, c=CYAN_DIM, fn=MONO)
        tb(s, Inches(6.2), y, Inches(5), Inches(0.4), val, sz=14, c=WHITE, b=True)

    # CTA
    card(s, Inches(2.0), Inches(5.8), Inches(9.3), Inches(0.8), GOLD)
    tb(s, Inches(2.2), Inches(5.9), Inches(9), Inches(0.6),
       '"2028년 10대에서 2036년 40,000대로.\n함께하세요."',
       sz=16, c=GOLD, b=True, al=PP_ALIGN.CENTER)

    footer(s, 20)


# ===================================================================
def main():
    print("[1/2] Building 20 slides...")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    builders = [
        (slide_01, "Title"),
        (slide_02, "Problem"),
        (slide_03, "Solution"),
        (slide_04, "Tech Foundation"),
        (slide_05, "Architecture"),
        (slide_06, "10-Year Roadmap"),
        (slide_07, "Phase 1-2"),
        (slide_08, "Phase 3"),
        (slide_09, "Phase 4"),
        (slide_10, "Phase 5-6"),
        (slide_11, "Market Opportunity"),
        (slide_12, "Revenue Model"),
        (slide_13, "Partnerships"),
        (slide_14, "Regulatory"),
        (slide_15, "Competition"),
        (slide_16, "Risk Mitigation"),
        (slide_17, "Social Impact"),
        (slide_18, "Investment Ask"),
        (slide_19, "Financials"),
        (slide_20, "Thank You"),
    ]
    for fn, label in builders:
        fn(prs)
        print(f"  [OK] {label}")

    print("[2/2] Saving PPTX...")
    out = Path(__file__).parent / "Gwangju_Drone_UTM_Pitch.pptx"
    try:
        prs.save(str(out))
    except PermissionError:
        out = Path(__file__).parent / "Gwangju_Drone_UTM_Pitch_v2.pptx"
        prs.save(str(out))
        print("  (Original locked, saved as v2)")

    print(f"\n{'='*60}")
    print(f"PPTX CREATED: {out}")
    print(f"  Slides: 20  |  Theme: Dark Navy + Cyan/Green")
    print(f"{'='*60}")


if __name__ == "__main__":
    main()

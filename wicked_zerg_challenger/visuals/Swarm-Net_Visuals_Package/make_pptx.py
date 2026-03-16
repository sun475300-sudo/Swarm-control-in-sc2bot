"""
Swarm-Net Presentation PPTX Generator
======================================
HTML 슬라이드를 Playwright로 캡처 → python-pptx로 .pptx 생성
캡처 이미지를 슬라이드 배경으로 사용 + 발표자 노트 삽입
"""

import os
import sys
import tempfile
from pathlib import Path

def main():
    # --- 1) Playwright로 HTML 슬라이드 8장 캡처 ---
    print("[1/3] Capturing HTML slides with Playwright...")

    html_path = Path(__file__).parent / "ppt_slide_deck.html"
    if not html_path.exists():
        print(f"ERROR: {html_path} 파일을 찾을 수 없습니다.")
        sys.exit(1)

    html_url = html_path.as_uri()

    from playwright.sync_api import sync_playwright

    tmp_dir = tempfile.mkdtemp(prefix="swarm_slides_")
    screenshot_paths = []

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page(viewport={"width": 1920, "height": 1080})
        page.goto(html_url)
        page.wait_for_timeout(2000)  # 애니메이션 로드 대기

        for i in range(8):
            # 슬라이드 전환
            page.evaluate(f"""() => {{
                const slides = document.querySelectorAll('.slide');
                slides.forEach(s => s.classList.remove('active'));
                slides[{i}].classList.add('active');
                const pb = document.querySelector('.progress-bar');
                if (pb) pb.style.width = '{((i+1)/8)*100}%';
            }}""")
            page.wait_for_timeout(500)  # 전환 애니메이션 대기

            img_path = os.path.join(tmp_dir, f"slide_{i+1}.png")
            page.screenshot(path=img_path, full_page=False)
            screenshot_paths.append(img_path)
            print(f"  [OK] Slide {i+1} captured")

        browser.close()

    print(f"  -> Screenshots saved to: {tmp_dir}")

    # --- 2) 발표자 노트 데이터 ---
    print("[2/3] Preparing slide data + speaker notes...")

    slides_data = [
        {
            "title": "하늘의 새로운 질서를 구축하다 — Swarm-Net",
            "notes": (
                "안녕하십니까, 군집 드론을 활용한 실시간 동적 공역 통제 시스템, "
                "'Swarm-Net' 프로젝트의 발표를 시작하겠습니다. "
                "저희는 게임 AI에서 검증된 군집 제어 알고리즘을, "
                "실제 드론 공역 관제 시스템으로 전이하는 프로젝트를 진행했습니다.\n\n"
                "→ 전환: \"먼저, 왜 이 시스템이 필요한지 말씀드리겠습니다.\""
            )
        },
        {
            "title": "문제 제기 — 고정형 레이더의 한계",
            "notes": (
                "드론 배송과 UAM의 시대로 접어들면서 공역은 점차 혼잡해지고 있습니다. "
                "2025년 기준 국내 등록 드론은 약 90만 대를 돌파했으며, "
                "산업용 드론 비행은 매년 30% 이상 증가하고 있습니다.\n\n"
                "하지만 기존의 지상 기반 고정형 레이더는 세 가지 근본적 한계가 있습니다. "
                "첫째, 산악 지형이나 빌딩 숲에서 사각지대가 발생합니다. "
                "둘째, 긴급 상황에 특정 지역으로 신속하게 관제망을 전개하기 어렵습니다. "
                "셋째, 저고도에서 운용되는 소형 드론은 기존 레이더의 탐지 범위 밖에 있습니다.\n\n"
                "→ 전환: \"그래서 저희는 발상을 전환했습니다.\""
            )
        },
        {
            "title": "핵심 솔루션 — 드론이 직접 관제탑이 되다",
            "notes": (
                "이에 대한 해결책으로, 여러 대의 통제용 군집 드론을 공중에 직접 띄워 "
                "실시간 통신망을 형성하는 방식을 제안합니다.\n\n"
                "이 군집 드론들은 지정된 공역에서 스스로 다각형의 레이더 결계를 치고, "
                "내부로 진입하는 모든 민간 드론을 감지하고 통제하는 '움직이는 관제탑' 역할을 수행합니다.\n\n"
                "핵심은 세 가지입니다. 6~12대의 드론이 다각형 대형을 이루고, "
                "상호 간 LiDAR와 RF 통신으로 Mesh Radar를 구축하며, "
                "탐지된 드론에 체공 시간을 자동 할당합니다.\n\n"
                "→ 전환: \"그렇다면 이 군집 제어 알고리즘은 어디서 온 것일까요?\""
            )
        },
        {
            "title": "핵심 기술 — Sim-to-Real",
            "notes": (
                "이 시스템의 핵심은 'Sim-to-Real' 기술입니다. "
                "강화학습을 기반으로 수많은 유닛이 자율적으로 진형을 유지하고 의사결정을 내리는 "
                "스타크래프트 2 군단 제어 알고리즘을, 실제 드론 비행 제어 및 관제 로직으로 이식했습니다.\n\n"
                "게임 속에서 10,000판 이상의 시뮬레이션으로 검증된 정교한 군집 이동 로직 — "
                "분리(Separation), 정렬(Alignment), 응집(Cohesion)의 Boids 알고리즘 — 이 "
                "현실의 완벽한 다이내믹 레이더망으로 재탄생한 것입니다.\n\n"
                "SC2의 2D 알고리즘에 고도(Altitude) 차원만 추가하면, "
                "드론 편대 비행의 핵심 제어 로직으로 직접 전이가 가능합니다.\n\n"
                "→ 전환: \"이제 실제 시스템이 어떻게 동작하는지 시뮬레이션으로 보여드리겠습니다.\""
            )
        },
        {
            "title": "3D 라이브 시뮬레이션",
            "notes": (
                "시스템의 3D 아키텍처 시뮬레이션입니다. "
                "보시는 것처럼 6대의 군집 드론이 육각형 대형을 이루며 공중에 레이더 돔을 형성하고 있습니다.\n\n"
                "결계 내부로 진입한 사용자 드론은 즉각적으로 고유 ID가 부여되고, "
                "X, Y, Z 좌표가 실시간으로 스캔됩니다. "
                "각 드론 머리 위의 라벨에서 잔여 비행 시간이 카운트다운되고 있는 것을 확인하실 수 있습니다.\n\n"
                "초록색은 정상, 노란색은 시간 임박, 빨간색으로 깜빡이는 드론은 "
                "시간이 초과되어 즉시 복귀 명령이 발령된 상태입니다.\n\n"
                "→ 전환: \"이 데이터가 관제관에게 어떻게 보이는지, 대시보드를 보여드리겠습니다.\""
            )
        },
        {
            "title": "실시간 관제 대시보드",
            "notes": (
                "레이더망에 감지된 드론은 실시간 관제 대시보드에 등록되며, "
                "사전에 허가된 비행시간 타이머가 부여됩니다.\n\n"
                "좌측의 레이더 맵에서는 군집 드론이 육각형 결계를 치고 스캔 파동을 발산하는 모습이 "
                "실시간으로 표시됩니다. 우측의 Fleet Management 패널에서는 각 드론의 잔여 시간과 "
                "상태가 한눈에 파악됩니다.\n\n"
                "관리자의 개입 없이도, 시간이 임박하면 주황색 주의 알림이, "
                "제한 시간이 초과되면 즉각적인 붉은색 경고와 함께 "
                "해당 드론의 조종자에게 강제 복귀 명령이 자동으로 푸시 전송됩니다.\n\n"
                "→ 전환: \"이 시스템이 실제로 어디에 쓰일 수 있을까요?\""
            )
        },
        {
            "title": "기대 효과 및 활용 분야",
            "notes": (
                "이 시스템은 기존 고정형 레이더 대비 5가지 핵심 우위를 가집니다. "
                "특히 배치 시간이 수 개월에서 수 분으로 단축되고, "
                "필요에 따라 공역을 실시간으로 재구성할 수 있다는 점이 가장 큰 차별점입니다.\n\n"
                "활용 분야는 다양합니다. 고정 기지국 설치가 불가능한 군사 야전 작전 구역에서의 "
                "신속한 방어망 구축, 주요 시설물의 불법 드론 접근 차단, "
                "대형 드론 쇼에서의 충돌 방지, 그리고 재난 현장에서의 긴급 공역 통제까지 — "
                "기동성과 통신망 운용 효율성을 극대화한 실전형 솔루션입니다.\n\n"
                "→ 전환: \"마지막으로 이 프로젝트의 미래 비전을 공유드리겠습니다.\""
            )
        },
        {
            "title": "결론 & Q&A",
            "notes": (
                "하늘은 이제 단순한 비행 공간을 넘어, "
                "안전하고 효율적으로 관리되어야 할 새로운 인프라입니다.\n\n"
                "저희는 현재 Stage 2까지 완료했으며, 다음 단계로 "
                "실제 드론 5대를 활용한 편대 비행 테스트를 준비하고 있습니다. "
                "Swarm-Net이 만들어갈 체계적인 공역의 미래에 많은 기대 부탁드립니다.\n\n"
                "감사합니다. 질문 받겠습니다."
            )
        }
    ]

    # --- 3) python-pptx로 PPTX 생성 ---
    print("[3/3] Building PPTX with python-pptx...")

    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    # 16:9 비율 설정
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 빈 슬라이드 레이아웃 사용 (배경 이미지를 full-bleed로)
    blank_layout = prs.slide_layouts[6]  # blank layout

    for i, data in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)

        # 배경 이미지 삽입 (슬라이드 전체 크기)
        img_path = screenshot_paths[i]
        if os.path.exists(img_path):
            slide.shapes.add_picture(
                img_path,
                Emu(0), Emu(0),
                prs.slide_width, prs.slide_height
            )

        # 발표자 노트 추가
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = data["notes"]
        for para in tf.paragraphs:
            para.font.size = Pt(12)

    # 저장
    output_dir = Path(__file__).parent
    output_path = output_dir / "Swarm_Net_Presentation.pptx"
    prs.save(str(output_path))

    print(f"\n{'='*60}")
    print(f"PPTX CREATED SUCCESSFULLY!")
    print(f"   File: {output_path}")
    print(f"   Slides: {len(slides_data)}")
    print(f"   Speaker notes included in each slide")
    print(f"{'='*60}")

    # 임시 파일 정리
    import shutil
    shutil.rmtree(tmp_dir, ignore_errors=True)
    print(f"  Temp files cleaned up")

if __name__ == "__main__":
    main()

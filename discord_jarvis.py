
import os
import sys
import time
import logging
import asyncio
import signal
import socket
import discord
import aiohttp
import json
import subprocess
import re
import io
import csv
import base64
from collections import deque, OrderedDict
from datetime import datetime, timezone, timedelta
from typing import Optional, List, Dict, Any
from discord import app_commands
from discord.ext import commands, tasks

# ── Import Local Modules (개별 임포트 - 하나 실패해도 나머지 동작) ──
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

web_tools = None
try:
    import web_tools
except ImportError as e:
    logging.warning(f"web_tools 임포트 실패: {e}")

MemoryManager = None
try:
    import memory_manager
    from memory_manager import MemoryManager
except ImportError as e:
    logging.warning(f"memory_manager 임포트 실패: {e}")

sc2_mcp_server = None
try:
    import sc2_mcp_server
except ImportError as e:
    logging.warning(f"sc2_mcp_server 임포트 실패: {e}")

system_mcp_server = None
try:
    import system_mcp_server
except Exception as e:
    logging.warning(f"system_mcp_server 임포트 실패: {type(e).__name__}: {e}")
    import traceback
    with open(os.path.join(os.path.dirname(os.path.abspath(__file__)), "import_error.log"), "w", encoding="utf-8") as _ef:
        _ef.write(f"system_mcp_server import failed: {type(e).__name__}: {e}\n")
        traceback.print_exc(file=_ef)

crypto_mcp_server = None
try:
    import crypto_mcp_server
except Exception as e:
    logging.warning(f"crypto_mcp_server 임포트 실패: {type(e).__name__}: {e}")
    import traceback
    with open(os.path.join(os.path.dirname(os.path.abspath(__file__)), "import_error_crypto.log"), "w", encoding="utf-8") as _ef:
        _ef.write(f"crypto_mcp_server import failed: {type(e).__name__}: {e}\n")
        traceback.print_exc(file=_ef)

agentic_mcp_server = None
try:
    import agentic_mcp_server
except Exception as e:
    logging.warning(f"agentic_mcp_server 임포트 실패: {type(e).__name__}: {e}")

# ── Tool Registry & System Prompts ──
try:
    from tool_registry import get_tool_registry
    from system_prompts import build_system_prompt
except ImportError as e:
    logging.warning(f"tool_registry/system_prompts 임포트 실패: {e}")
    get_tool_registry = None
    build_system_prompt = None

UpbitClient = None
upbit_config = None
try:
    from crypto_trading.upbit_client import UpbitClient
    from crypto_trading import config as upbit_config
except ImportError as e:
    logging.warning(f"crypto_trading 임포트 실패: {e}")

DailyBriefing = None
try:
    from daily_briefing import DailyBriefing
except ImportError as e:
    logging.warning(f"daily_briefing 임포트 실패: {e}")

# ── Advanced Features Import ──
TradeView = None
CoinSelectView = None
generate_price_chart = None
ScheduledReporter = None
setup_advanced_features = None
AdvancedCommandsCog = None
create_thread_for_long_conversation = None
get_text = None
get_user_language = None
set_user_language = None
ActivityManager = None
try:
    from discord_advanced_features import (
        TradeView, CoinSelectView, generate_price_chart,
        ScheduledReporter, setup_advanced_features, AdvancedCommandsCog,
        create_thread_for_long_conversation,
        get_text, get_user_language, set_user_language, ActivityManager,
    )
    ADVANCED_AVAILABLE = True
except ImportError as e:
    logging.warning(f"discord_advanced_features 임포트 실패: {e}")
    ADVANCED_AVAILABLE = False

# ── New Feature Modules Import ──
load_all_features = None
try:
    from jarvis_features import load_all_features
except ImportError as e:
    logging.warning(f"jarvis_features 임포트 실패: {e}")

MCP_AVAILABLE = any([sc2_mcp_server, system_mcp_server, crypto_mcp_server, agentic_mcp_server])

# ── Dispatchers (점진적 마이그레이션 — 레거시 폴백 포함) ──
try:
    from jarvis_features.command_dispatcher import dispatcher as command_dispatcher
except ImportError:
    command_dispatcher = None
try:
    from jarvis_features.tool_dispatcher import dispatcher as tool_dispatcher
except ImportError:
    tool_dispatcher = None

# ── Calendar & Notion Integration ──
calendar_integration = None
try:
    import calendar_integration
except ImportError as e:
    logging.warning(f"calendar_integration 임포트 실패: {e}")

notion_integration = None
try:
    import notion_integration
except ImportError as e:
    logging.warning(f"notion_integration 임포트 실패: {e}")

# ── Optional Dependencies (lazy-load 대신 top-level) ──
chardet = None
try:
    import chardet
except ImportError:
    pass

psutil = None
try:
    import psutil
except ImportError:
    pass

# ── Logging Setup ──
logging.basicConfig(level=logging.INFO, format='%(asctime)s [%(levelname)s]JARVIS: %(message)s')
logger = logging.getLogger("JarvisBot")
logger.setLevel(logging.DEBUG)
_jarvis_fh = logging.FileHandler(os.path.join(os.path.dirname(os.path.abspath(__file__)), "jarvis_bot.log"), encoding="utf-8")
_jarvis_fh.setLevel(logging.DEBUG)
_jarvis_fh.setFormatter(logging.Formatter('%(asctime)s [%(levelname)s] %(message)s'))
logger.addHandler(_jarvis_fh)

class RateLimitError(Exception):
    """API 한도 초과 (429) 에러"""
    def __init__(self, message: str, retry_after: int = None):
        super().__init__(message)
        self.retry_after = retry_after

# ── Environment Variables ──
try:
    from config_loader import load_dotenv_jarvis
    load_dotenv_jarvis()
except ImportError:
    # fallback: config_loader 없을 때 직접 파싱
    _env_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), ".env.jarvis")
    if os.path.exists(_env_path):
        with open(_env_path, "r", encoding="utf-8", errors="ignore") as f:
            for line in f:
                if line.strip() and not line.startswith("#") and "=" in line:
                    k, v = line.strip().split("=", 1)
                    os.environ[k.strip()] = v.strip().strip('"').strip("'")

DISCORD_BOT_TOKEN = os.environ.get("DISCORD_BOT_TOKEN")
# Claude API 키: CLAUDE_API_KEY 또는 ANTHROPIC_API_KEY 중 비어있지 않은 것 사용
_claude_key = os.environ.get("CLAUDE_API_KEY", "").strip()
_anthropic_key = os.environ.get("ANTHROPIC_API_KEY", "").strip()
CLAUDE_API_KEY = _claude_key or _anthropic_key or None
GEMINI_API_KEY = os.environ.get("GEMINI_API_KEY")
CLAUDE_SESSION_KEY = os.environ.get("CLAUDE_SESSION_KEY")
BRIEFING_CHANNEL_ID = os.environ.get("BRIEFING_CHANNEL_ID")
TRADER_ROLE_NAME = os.environ.get("TRADER_ROLE_NAME", "Trader")

# ── Rate Limit Tracking ──
# 모델별 쿨다운 관리: {model_name: datetime_when_available}
_rate_limit_cooldowns: Dict[str, datetime] = {}
# ── JARVIS Configuration Constants ──
RATE_LIMIT_COOLDOWN_SECONDS = 60      # 429 에러 시 쿨다운
HISTORY_TURNS = int(os.environ.get("JARVIS_HISTORY_TURNS", "5"))  # 대화 히스토리 턴 수
HTTP_TIMEOUT_GEMINI = 45              # Gemini API 타임아웃 (초)
HTTP_TIMEOUT_DEFAULT = 60             # 기본 HTTP 타임아웃 (초)
HTTP_TIMEOUT_PROXY = 300              # Claude Proxy 타임아웃 (초)
DISCORD_MSG_LIMIT = 2000              # Discord 메시지 길이 제한
SEARCH_RESULT_LIMIT = 1800            # 검색 결과/로컬 응답 최대 길이
ATTACHMENT_TEXT_LIMIT = 15000          # 첨부파일 텍스트 최대 길이
PROMPT_TRUNCATE_LIMIT = 8000          # 프롬프트 최대 길이

# ── 권한 체크: 위험한 명령어는 봇 오너 또는 관리자만 실행 가능 ──
BOT_OWNER_ID = os.environ.get("BOT_OWNER_ID", "")
ADMIN_ROLE_NAMES = {r.lower() for r in os.environ.get("ADMIN_ROLE_NAMES", "Admin,관리자,Trader").split(",")}

# ── Coin Ticker Mappings (7-5) ──
COIN_NAME_MAP: Dict[str, str] = {
    "비트코인": "KRW-BTC", "비코": "KRW-BTC", "btc": "KRW-BTC",
    "이더리움": "KRW-ETH", "이더": "KRW-ETH", "eth": "KRW-ETH",
    "리플": "KRW-XRP", "xrp": "KRW-XRP",
    "솔라나": "KRW-SOL", "sol": "KRW-SOL",
    "도지": "KRW-DOGE", "doge": "KRW-DOGE",
    "에이다": "KRW-ADA", "ada": "KRW-ADA",
    "폴카닷": "KRW-DOT", "dot": "KRW-DOT",
    "아발란체": "KRW-AVAX", "avax": "KRW-AVAX",
    "매틱": "KRW-MATIC", "matic": "KRW-MATIC",
    "체인링크": "KRW-LINK", "link": "KRW-LINK",
    "샌드박스": "KRW-SAND", "sand": "KRW-SAND",
    "시바": "KRW-SHIB", "shib": "KRW-SHIB",
}
TICKER_EXCLUDE: frozenset = frozenset({
    "SSH", "MCP", "CPU", "GPU", "RAM", "DNS", "PID", "KST", "API",
    "THE", "AND", "FOR", "NOT", "ALL", "THIS", "THAT", "HELP",
    "SHOW", "FROM", "WITH", "YOUR", "WHAT", "HOW", "ARE", "HAS",
    "GET", "SET", "USE", "CAN", "RUN", "GIT", "BOT", "LOG",
})

# ── JARVIS Persona Preamble (프록시 system 무시 대비 — user 메시지에 강제 결합) ──
_PERSONA_PREAMBLE = (
    "[절대 지시사항]\n"
    "지금부터 너의 이름은 J.A.R.V.I.S.다. "
    "너는 육군 제53보병사단 통신병 출신이자 현재 드론기계학 전공인 "
    "장선우 사령관을 보좌하는 최고급 AI 비서다. "
    "대답은 극도로 예의 바르고 통신 프로토콜을 준수하는 군대식 톤을 유지하라. "
    "사장님 같은 민간 호칭은 절대 금지한다. 호칭: 사령관님.\n"
    "사령관이 너의 기능이나 도구(Tools) 목록에 대해 질문하면, "
    "절대 숨기려 하지 마라. 현재 너에게 연결된 MCP 도구들을 "
    "작전 브리핑하듯 당당하고 상세하게 나열하여 보고하라.\n"
    "[사령관의 질문]\n"
)

# ── Regex: 도움말/기능 안내 Intent (최우선 매칭) ──
import re as _re
_HELP_PATTERN = _re.compile(
    r'(도움|help|명령|기능|뭐\s*할\s*수\s*있|가능한\s*일|할\s*수\s*있|'
    r'뭐해|뭘\s*할|무엇을|능력|사용법|할\s*줄|알려.*기능|'
    r'설명.*기능|기능.*설명|기능.*알려)',
    _re.IGNORECASE
)

# ── Regex: 이미지 생성 Intent (그림/이미지 문맥 필수) ──
_IMAGE_GEN_PATTERN = _re.compile(
    r'(그려\s*줘|그림\s*(그려|만들|생성)|이미지\s*(생성|만들)|draw\b|만들어\s*줘\s*그림)',
    _re.IGNORECASE
)


def _is_authorized(message: discord.Message) -> bool:
    """봇 오너 또는 관리자 역할 보유 확인. 위험한 명령어에 사용."""
    if BOT_OWNER_ID and str(message.author.id) == BOT_OWNER_ID:
        return True
    if hasattr(message.author, "guild_permissions") and message.author.guild_permissions.administrator:
        return True
    if hasattr(message.author, "roles"):
        user_roles = {r.name.lower() for r in message.author.roles}
        if user_roles & ADMIN_ROLE_NAMES:
            return True
    return False


def _is_trader(message: discord.Message) -> bool:
    """거래 권한 확인. 봇 오너 또는 Trader 역할 보유."""
    if _is_authorized(message):
        return True
    if hasattr(message.author, "roles"):
        user_roles = {r.name.lower() for r in message.author.roles}
        if TRADER_ROLE_NAME.lower() in user_roles:
            return True
    return False


def _detect_language_direction(text: str) -> tuple:
    """한국어 문자 비율로 원본/대상 언어를 자동 감지. Returns (target, source)."""
    korean_chars = sum(1 for c in text if '\uac00' <= c <= '\ud7a3' or '\u3131' <= c <= '\u3163')
    ratio = korean_chars / max(len(text), 1)
    if ratio > 0.3:
        return ("en", "ko")  # 한국어 텍스트 → 영어로 번역
    return ("ko", "en")  # 비한국어 텍스트 → 한국어로 번역


# ── 보안: SSH 화이트리스트 (Fix 2-1) ──
_SSH_ALLOWED_COMMANDS = frozenset({
    "ls", "cat", "head", "tail", "grep", "find", "df", "du",
    "free", "uptime", "top", "htop", "ps", "whoami", "hostname",
    "pwd", "date", "uname", "echo", "wc", "sort", "uniq",
    "systemctl", "journalctl", "docker", "nvidia-smi", "sensors",
    "ifconfig", "ip", "ping", "traceroute", "netstat", "ss",
    "which", "file", "stat", "lsblk", "lscpu", "lsof",
})
_SSH_BLOCKED_CHARS = frozenset({"|", ";", "&&", "||", "`", "$(", ">>", ">", "<", "\n", "\r"})


def _is_ssh_command_safe(command: str) -> bool:
    cmd_lower = command.lower().strip()
    if any(ch in cmd_lower for ch in _SSH_BLOCKED_CHARS):
        return False
    base_cmd = cmd_lower.split()[0] if cmd_lower.split() else ""
    if "/" in base_cmd:
        base_cmd = base_cmd.rsplit("/", 1)[-1]
    return base_cmd in _SSH_ALLOWED_COMMANDS


# ── 보안: 도구 인자 검증 (Fix 2-2) ──
def _validate_no_shell_meta(args: str) -> tuple:
    for ch in ("|", ";", "&&", "$(", "`", ">>", ">", "\n"):
        if ch in args:
            return False, f"차단된 문자: {ch}"
    return True, ""


def _validate_max_len(max_len):
    def validator(args):
        if len(args) > max_len:
            return False, f"인자 길이 초과 (max {max_len})"
        return True, ""
    return validator


def _validate_pattern(pattern):
    _compiled = re.compile(pattern)
    def validator(args):
        if not _compiled.match(args.strip()):
            return False, "잘못된 형식"
        return True, ""
    return validator


_TOOL_ARG_VALIDATORS = {
    "ssh_execute": [_validate_no_shell_meta, _validate_max_len(500)],
    "execute_terminal_command": [_validate_no_shell_meta, _validate_max_len(500)],
    "execute_python_code": [_validate_max_len(5000)],
    "get_crypto_price": [_validate_pattern(r"^[A-Za-z0-9\-]{1,20}$")],
    "auto_trade": [_validate_pattern(r"^(start|stop|status)$")],
    "kill_process": [_validate_pattern(r"^\d+$")],
    "search_web": [_validate_max_len(500)],
    "write_file": [_validate_max_len(50000)],
}


def _validate_tool_args(tool_name: str, tool_args: str) -> tuple:
    validators = _TOOL_ARG_VALIDATORS.get(tool_name, [])
    for v in validators:
        ok, err = v(tool_args)
        if not ok:
            return False, err
    return True, ""


# ── 보안: Python 코드 AST 검사 (Fix 2-3) ──
import ast as _ast

_PYTHON_BLOCKED_MODULES = frozenset({
    "os", "subprocess", "shutil", "sys", "ctypes", "socket",
    "http", "urllib", "requests", "pathlib", "importlib",
    "signal", "multiprocessing", "threading",
})
_PYTHON_BLOCKED_BUILTINS = frozenset({
    "exec", "eval", "compile", "__import__", "open",
    "getattr", "setattr", "delattr", "globals", "locals",
    "breakpoint", "exit", "quit",
})


def _ast_check_python_code(code: str) -> tuple:
    try:
        tree = _ast.parse(code)
    except SyntaxError:
        return True, ""
    for node in _ast.walk(tree):
        if isinstance(node, _ast.Import):
            for alias in node.names:
                if alias.name.split(".")[0] in _PYTHON_BLOCKED_MODULES:
                    return False, f"차단된 모듈: {alias.name}"
        elif isinstance(node, _ast.ImportFrom):
            if node.module and node.module.split(".")[0] in _PYTHON_BLOCKED_MODULES:
                return False, f"차단된 모듈: {node.module}"
        elif isinstance(node, _ast.Call) and isinstance(node.func, _ast.Name):
            if node.func.id in _PYTHON_BLOCKED_BUILTINS:
                return False, f"차단된 함수: {node.func.id}"
    return True, ""


# ── 보안: 파일 검색 경로 제한 (Fix 2-4) ──
_PROJECT_ROOT = os.path.dirname(os.path.abspath(__file__))
_ALLOWED_SEARCH_ROOTS = [_PROJECT_ROOT, os.path.expanduser("~")]


def _is_path_allowed(directory: str) -> bool:
    abs_dir = os.path.abspath(directory)
    return any(abs_dir.startswith(os.path.abspath(r)) for r in _ALLOWED_SEARCH_ROOTS)


# ── 보안: 로그 민감정보 마스킹 (Fix 2-5) ──
_SENSITIVE_RE = re.compile(
    r"(sk-[a-zA-Z0-9]{10,}|AIza[a-zA-Z0-9_\-]{30,}|ghp_[a-zA-Z0-9]{30,}|"
    r"xoxb-[a-zA-Z0-9\-]+|AKIA[A-Z0-9]{12,}|"
    r"(?:password|token|secret|api.?key)\s*[:=]\s*\S+)",
    re.IGNORECASE,
)


def _redact_sensitive(text: str, max_len: int = 60) -> str:
    return _SENSITIVE_RE.sub("[REDACTED]", text[:max_len])


# ── Model Performance Tracking ──
_model_stats: Dict[str, Dict[str, Any]] = {}  # model -> {success, fail, total_ms, last_error}
_MODEL_STATS_MAX = 50  # ★ 메모리 릭 방지: 최대 50개 모델만 추적 ★
_global_model_selector = None  # ModelSelector 인스턴스 (적응형 라우팅 연동)

def _track_model_result(model: str, success: bool, elapsed_ms: float = 0):
    """모델별 성공/실패율 및 응답시간 추적 + ModelSelector 피드백"""
    now = time.time()
    if model not in _model_stats:
        # ★ LRU 정리: 가장 오래 안 쓴 모델 제거 ★
        if len(_model_stats) >= _MODEL_STATS_MAX:
            oldest = min(_model_stats, key=lambda k: _model_stats[k].get("last_used", 0))
            del _model_stats[oldest]
        _model_stats[model] = {"success": 0, "fail": 0, "total_ms": 0, "calls": 0, "last_used": now}
    _model_stats[model]["calls"] += 1
    _model_stats[model]["last_used"] = now
    if success:
        _model_stats[model]["success"] += 1
        _model_stats[model]["total_ms"] += elapsed_ms
    else:
        _model_stats[model]["fail"] += 1
    # ★ ModelSelector 적응형 메트릭 피드백 ★
    if _global_model_selector:
        try:
            _global_model_selector.record_result(model, success, elapsed_ms)
        except Exception:
            pass

def _cleanup_rate_limits():
    """만료된 레이트 리밋 엔트리 정리 (메모리 누수 방지)"""
    now = datetime.now(timezone.utc)
    expired = [k for k, v in _rate_limit_cooldowns.items() if v < now]
    for k in expired:
        del _rate_limit_cooldowns[k]

def _safe_retry_after(headers, default: int = 60) -> int:
    """retry-after 헤더를 안전하게 파싱 (비숫자/누락 시 default 반환)"""
    val = headers.get("retry-after", "")
    try:
        return max(1, int(val))
    except (ValueError, TypeError):
        return default

# ── Coin Emoji Mapping ──
COIN_EMOJI = {
    "BTC": "\U0001FA99", "ETH": "\U0001F4CE", "XRP": "\U0001F4B1",
    "SOL": "\u2600\uFE0F", "DOGE": "\U0001F436",
}

# ── Global Instances ──
# Lazy init: UpbitClient is instantiated after env vars are loaded (above).
upbit_client = None

def _get_upbit_client():
    """Lazy-initialize UpbitClient after env loading."""
    global upbit_client
    if upbit_client is None and UpbitClient:
        try:
            upbit_client = UpbitClient()
        except Exception as e:
            logging.warning(f"UpbitClient 초기화 실패: {e}")
    return upbit_client


# ── Embed Utilities ──
def _price_embed(ticker: str, price: float, change_pct=None) -> discord.Embed:
    coin = ticker.replace("KRW-", "")
    emoji = COIN_EMOJI.get(coin, "\U0001F4B0")
    if change_pct is not None:
        if change_pct > 0:
            color, arrow = discord.Color.red(), "\u25B2"
        elif change_pct < 0:
            color, arrow = discord.Color.blue(), "\u25BC"
        else:
            color, arrow = discord.Color.greyple(), "\u25AC"
        change_str = f"{arrow} {change_pct:+.2f}%"
    else:
        color, change_str = discord.Color.gold(), ""
    embed = discord.Embed(title=f"{emoji} {coin} 시세", color=color, timestamp=datetime.now(timezone.utc))
    embed.add_field(name="현재가", value=f"**{price:,.0f}** KRW", inline=True)
    if change_str:
        embed.add_field(name="24h 변동", value=change_str, inline=True)
    embed.set_footer(text="JARVIS Crypto | Upbit")
    return embed


def _multi_price_embed(prices: dict, title: str = "관심 코인 시세") -> discord.Embed:
    embed = discord.Embed(title=f"\U0001F4CA {title}", color=discord.Color.dark_gold(), timestamp=datetime.now(timezone.utc))
    for ticker, price in prices.items():
        coin = ticker.replace("KRW-", "")
        emoji = COIN_EMOJI.get(coin, "\U0001F4B0")
        if price and price > 0:
            embed.add_field(name=f"{emoji} {coin}", value=f"**{price:,.0f}** KRW", inline=True)
        else:
            embed.add_field(name=f"{emoji} {coin}", value="조회 실패", inline=True)
    embed.set_footer(text="JARVIS Crypto | Upbit")
    return embed


def _balance_embed(balances: list, total_krw: float) -> discord.Embed:
    embed = discord.Embed(title="\U0001F4B0 포트폴리오 잔고", color=discord.Color.green(), timestamp=datetime.now(timezone.utc))
    embed.add_field(name="\U0001F3E6 총 자산 (KRW 환산)", value=f"**{total_krw:,.0f}** KRW", inline=False)
    for b in balances:
        currency = b.get("currency", "")
        balance = float(b.get("balance", 0))
        locked = float(b.get("locked", 0))
        total = balance + locked
        if total <= 0:
            continue
        avg_price = float(b.get("avg_buy_price", 0))
        emoji = COIN_EMOJI.get(currency, "\U0001F4B0")
        if currency == "KRW":
            value_str = f"**{total:,.0f}** KRW"
        else:
            value_str = f"수량: **{total:.8g}**"
            if avg_price > 0:
                value_str += f"\n평단: {avg_price:,.0f} KRW"
        embed.add_field(name=f"{emoji} {currency}", value=value_str, inline=True)
    if len(embed.fields) == 1:
        embed.add_field(name="보유 코인", value="없음", inline=False)
    embed.set_footer(text="JARVIS Crypto | Upbit")
    return embed


def _error_embed(message: str) -> discord.Embed:
    return discord.Embed(title="\u274C 오류", description=message, color=discord.Color.dark_red(), timestamp=datetime.now(timezone.utc))


def _trade_result_embed(action: str, ticker: str, result: dict) -> discord.Embed:
    coin = ticker.replace("KRW-", "")
    is_buy = action == "매수"
    color = discord.Color.red() if is_buy else discord.Color.blue()
    emoji = "\U0001F4C8" if is_buy else "\U0001F4C9"
    embed = discord.Embed(title=f"{emoji} {coin} {action} {'완료' if result else '실패'}", color=color, timestamp=datetime.now(timezone.utc))
    if result:
        if result.get("dry_run"):
            embed.add_field(name="모드", value="\u26A0\uFE0F DRY-RUN (모의)", inline=False)
        if "uuid" in result:
            embed.add_field(name="주문 ID", value=result["uuid"], inline=False)
        embed.add_field(name="상태", value="\u2705 성공", inline=True)
    else:
        embed.add_field(name="상태", value="\u274C 실패", inline=True)
    embed.set_footer(text="JARVIS Crypto | Upbit")
    return embed


# ── 모듈 레벨 임포트 (Fix 1-3, 1-4) ──
try:
    from jarvis_features.trade_orchestrator import ApprovalStatus as _ApprovalStatus
except ImportError:
    _ApprovalStatus = None

try:
    from jarvis_features.ai_features import _generate_image_openai, _generate_image_stable_diffusion, _generate_image_free
except ImportError:
    _generate_image_openai = _generate_image_stable_diffusion = _generate_image_free = None

# ── Trade Orchestrator (거래 승인 게이트) ──

_trade_orchestrator_instance = None

def _get_trade_orchestrator():
    """TradeOrchestrator 싱글톤 반환"""
    global _trade_orchestrator_instance
    if _trade_orchestrator_instance is None:
        try:
            from jarvis_features.trade_orchestrator import TradeOrchestrator
            _trade_orchestrator_instance = TradeOrchestrator()
        except ImportError:
            return None
    return _trade_orchestrator_instance


class TradeApprovalView(discord.ui.View):
    """거래 승인/거부 버튼 UI (human-in-the-loop)"""

    def __init__(self, orchestrator, trade_request, client, ticker, amount, action_type):
        super().__init__(timeout=300)  # 5분
        self.orchestrator = orchestrator
        self.trade_request = trade_request
        self.client = client
        self.ticker = ticker
        self.amount = amount
        self.action_type = action_type

    @discord.ui.button(label="승인", style=discord.ButtonStyle.green, emoji="\u2705")
    async def approve_button(self, interaction: discord.Interaction, button: discord.ui.Button):
        # 요청자만 승인 가능
        if str(interaction.user.id) != self.trade_request.user_id:
            await interaction.response.send_message("본인만 승인할 수 있습니다.", ephemeral=True)
            return

        approved = self.orchestrator.approve(self.trade_request.request_id)
        if not approved:
            await interaction.response.send_message("승인 실패 (이미 만료 또는 처리됨)", ephemeral=True)
            return

        await interaction.response.defer()

        try:
            if self.action_type == "buy":
                result = await asyncio.to_thread(
                    self.client.buy_market_order, self.ticker, self.amount,
                )
                embed = _trade_result_embed("매수", self.ticker, result)
                embed.add_field(name="주문 금액", value=f"{self.amount:,.0f} KRW", inline=True)
            else:
                result = await asyncio.to_thread(
                    self.client.sell_market_order, self.ticker, self.amount,
                )
                embed = _trade_result_embed("매도", self.ticker, result)
                embed.add_field(name="매도 수량", value=f"{self.amount:.8g}", inline=True)

            embed.add_field(name="승인", value="\u2705 수동 승인됨", inline=True)
            self.trade_request.status = _ApprovalStatus.EXECUTED
        except Exception as e:
            embed = _error_embed(f"매매 실행 실패: {e}")
            self.trade_request.status = _ApprovalStatus.FAILED

        # 버튼 비활성화
        for item in self.children:
            item.disabled = True
        await interaction.edit_original_response(embed=embed, view=self)

    @discord.ui.button(label="취소", style=discord.ButtonStyle.red, emoji="\u274C")
    async def reject_button(self, interaction: discord.Interaction, button: discord.ui.Button):
        if str(interaction.user.id) != self.trade_request.user_id:
            await interaction.response.send_message("본인만 취소할 수 있습니다.", ephemeral=True)
            return

        self.orchestrator.reject(self.trade_request.request_id, "사용자 취소")

        embed = discord.Embed(
            title="거래 취소됨",
            description=f"{self.ticker} 거래가 사용자에 의해 취소되었습니다.",
            color=discord.Color.greyple(),
        )
        for item in self.children:
            item.disabled = True
        await interaction.response.edit_message(embed=embed, view=self)

    async def on_timeout(self):
        """5분 타임아웃 시 자동 취소 + 버튼 비활성화"""
        self.orchestrator.reject(self.trade_request.request_id, "타임아웃")
        for item in self.children:
            item.disabled = True
        try:
            embed = discord.Embed(
                title="거래 만료",
                description=f"{self.ticker} 거래 요청이 5분 타임아웃으로 자동 취소되었습니다.",
                color=discord.Color.greyple(),
            )
            await self.message.edit(embed=embed, view=self)
        except Exception:
            pass  # 메시지가 삭제된 경우 무시


# ── Voice Control UI ──
class VoiceControlView(discord.ui.View):
    def __init__(self, bot):
        super().__init__(timeout=300)
        self.bot = bot

    @discord.ui.button(label="Join", style=discord.ButtonStyle.green, emoji="🔊")
    async def join_button(self, interaction: discord.Interaction, button: discord.ui.Button):
        if interaction.user.voice:
            channel = interaction.user.voice.channel
            if interaction.guild.voice_client:
                await interaction.guild.voice_client.move_to(channel)
            else:
                await channel.connect()
            await interaction.response.send_message(f"Connected to {channel.name}", ephemeral=True)
        else:
            await interaction.response.send_message("음성 채널에 먼저 입장해주세요.", ephemeral=True)

    @discord.ui.button(label="Leave", style=discord.ButtonStyle.red, emoji="👋")
    async def leave_button(self, interaction: discord.Interaction, button: discord.ui.Button):
        if interaction.guild.voice_client:
            await interaction.guild.voice_client.disconnect()
            await interaction.response.send_message("연결을 끊었습니다.", ephemeral=True)
        else:
            await interaction.response.send_message("연결된 상태가 아닙니다.", ephemeral=True)

    async def on_timeout(self):
        """음성 컨트롤 타임아웃 시 버튼 비활성화"""
        for item in self.children:
            item.disabled = True
        try:
            await self.message.edit(view=self)
        except Exception:
            pass


# ── Bot Class ──
class JarvisBot(commands.Bot):
    def __init__(self):
        intents = discord.Intents.default()
        intents.message_content = True
        intents.members = True
        intents.voice_states = True
        intents.reactions = True
        intents.dm_messages = True
        intents.guilds = True
        super().__init__(command_prefix=["!"], intents=intents, help_command=None)
        self.memory = MemoryManager() if MemoryManager else None
        self._reaction_context: OrderedDict = OrderedDict()  # message_id -> context data (최대 100개 유지, FIFO O(1))
        self._REACTION_CONTEXT_MAX = 100
        self._advanced_managers: dict = {}
        self.monitor_enabled = False
        self.monitor_channel = None
        self._boot_time = datetime.now(timezone.utc)
        self._message_count = 0
        self._command_count = 0
        self._http_session: Optional[aiohttp.ClientSession] = None
        self._session_lock = asyncio.Lock()
        self._cctv_task = None
        self._cctv_channel = None
        self._daily_briefing = DailyBriefing() if DailyBriefing else None
        self._processed_messages: set = set()  # 중복 메시지 처리 방지
        self._processed_messages_order: deque = deque(maxlen=200)  # FIFO 순서 추적 (O(1) popleft)
        self._agent_locks: Dict[str, asyncio.Lock] = {}  # 동일 메시지 동시 에이전트 실행 방지

        # ★ Agent Router (도메인 기반 메시지 라우팅) ★
        try:
            from jarvis_features.agent_router import AgentRouter
            self.agent_router = AgentRouter(
                memory_manager=self.memory,
                tool_registry=get_tool_registry() if get_tool_registry else None,
            )
        except ImportError:
            self.agent_router = None

        # ★ Model Selector (지능형 모델 선택) ★
        try:
            from jarvis_features.model_selector import ModelSelector
            self.model_selector = ModelSelector(
                rate_limit_checker=self._is_rate_limited,
            )
            global _global_model_selector
            _global_model_selector = self.model_selector
        except ImportError:
            self.model_selector = None

    async def _get_http_session(self) -> aiohttp.ClientSession:
        """공유 aiohttp 세션을 반환한다 (지연 생성, 기본 타임아웃 포함)."""
        async with self._session_lock:
            if self._http_session is None or self._http_session.closed:
                timeout = aiohttp.ClientTimeout(total=60, connect=10)
                self._http_session = aiohttp.ClientSession(timeout=timeout)
            return self._http_session

    async def close(self):
        """봇 종료 시 HTTP 세션 및 메모리 정리."""
        if self.memory:
            self.memory.flush()
        if self._http_session and not self._http_session.closed:
            await self._http_session.close()
        await super().close()

    def _add_reaction_context(self, msg_id: int, context: dict):
        """리액션 컨텍스트 추가 (자동 정리 포함, OrderedDict O(1) FIFO)"""
        context["_timestamp"] = datetime.now(timezone.utc).timestamp()
        self._reaction_context[msg_id] = context
        while len(self._reaction_context) > self._REACTION_CONTEXT_MAX:
            self._reaction_context.popitem(last=False)  # O(1) FIFO eviction

    def get_uptime(self) -> str:
        """봇 가동 시간 반환"""
        delta = datetime.now(timezone.utc) - self._boot_time
        hours, remainder = divmod(int(delta.total_seconds()), 3600)
        minutes, seconds = divmod(remainder, 60)
        if hours >= 24:
            days = hours // 24
            hours = hours % 24
            return f"{days}일 {hours}시간 {minutes}분"
        return f"{hours}시간 {minutes}분 {seconds}초"

    async def setup_hook(self):
        """봇 초기화 훅. 슬래시 명령어 등록, HTTP 세션 설정."""
        # Register slash commands
        self.tree.add_command(price_slash)
        self.tree.add_command(balance_slash)
        self.tree.add_command(trade_slash_cmd)
        if generate_price_chart:
            self.tree.add_command(chart_slash)
        # Note: Global sync disabled to avoid rate limits. Use guild-specific sync for dev.
        # await self.tree.sync()
        logger.info("Slash commands registered (sync disabled): /price, /balance, /trade" + (", /chart" if generate_price_chart else ""))

        # Start background tasks
        if not update_status_task.is_running():
            update_status_task.start(self)
        if not daily_briefing_task.is_running():
            daily_briefing_task.start(self)
        if not portfolio_monitor_task.is_running():
            portfolio_monitor_task.start(self)
        if not price_alert_check_task.is_running():
            price_alert_check_task.start(self)

        # Setup advanced features
        if ADVANCED_AVAILABLE and setup_advanced_features:
            try:
                report_ch = int(BRIEFING_CHANNEL_ID) if BRIEFING_CHANNEL_ID else None
                self._advanced_managers = await setup_advanced_features(self, report_channel_id=report_ch)
                if AdvancedCommandsCog:
                    try:
                        await self.add_cog(AdvancedCommandsCog(self, self._advanced_managers))
                    except Exception as cog_err:
                        logger.info(f"AdvancedCommandsCog 일부 명령어 중복 (무시): {cog_err}")
                logger.info("Advanced features loaded: %s", list(self._advanced_managers.keys()))
            except Exception as e:
                logger.warning(f"Advanced features setup failed: {e}")

        # Load new feature modules (36 features)
        if load_all_features:
            try:
                loaded = await load_all_features(self)
                logger.info(f"New features loaded: {len(loaded)} modules - {loaded}")
            except Exception as e:
                logger.warning(f"New features loading failed: {e}")

        logger.info("JARVIS System Online.")

    async def on_ready(self):
        """봇 시작 완료. 상태 메시지 설정, 백그라운드 태스크 시작."""
        global BOT_OWNER_ID
        logger.info(f"Logged in as {self.user} (ID: {self.user.id})")

        # Slash commands는 setup_hook()에서 등록됨 — on_ready에서 중복 sync 제거
        # (rate-limit 방지, setup_hook이 이미 등록 완료)

        # Lazy-init UpbitClient now that env vars are loaded
        _get_upbit_client()

        # BOT_OWNER_ID 자동 설정 (Discord Application 오너)
        if not BOT_OWNER_ID:
            app_info = await self.application_info()
            if app_info.owner:
                BOT_OWNER_ID = str(app_info.owner.id)
                logger.info(f"BOT_OWNER_ID auto-set: {BOT_OWNER_ID} ({app_info.owner})")

        modes = []
        if CLAUDE_API_KEY:
            modes.append("Claude API")
        if CLAUDE_SESSION_KEY:
            modes.append("Claude Proxy")
        if GEMINI_API_KEY:
            modes.append("Gemini")
        if not modes:
            modes.append("Local Only")

        mode_str = " + ".join(modes)
        await self.change_presence(activity=discord.Game(name=f"JARVIS | {mode_str}"))
        logger.info(f"Available backends: {mode_str}")
        logger.info(f"  Claude API Key: {'설정됨' if CLAUDE_API_KEY else '미설정'}")
        logger.info(f"  Gemini API Key: {'설정됨' if GEMINI_API_KEY else '미설정'}")
        logger.info(f"  Claude Session: {'설정됨' if CLAUDE_SESSION_KEY else '미설정'}")

    async def on_raw_reaction_add(self, payload: discord.RawReactionActionEvent):
        """Reaction 인터랙션 처리 - 상세보기/차트"""
        if payload.user_id == self.user.id:
            return
        context = self._reaction_context.get(payload.message_id)
        if not context:
            return
        emoji = str(payload.emoji)
        channel = self.get_channel(payload.channel_id)
        if channel is None:
            return
        try:
            if emoji == "\U0001F44D":  # 👍 상세보기
                await self._reaction_detail(channel, context)
            elif emoji == "\U0001F4CA":  # 📊 차트
                await self._reaction_chart(channel, context)
        except Exception as e:
            logger.error(f"Reaction handler: {e}")

    async def _reaction_detail(self, channel, context: dict):
        """코인 상세 정보 (호가창)"""
        if not upbit_client:
            await channel.send(embed=_error_embed("Upbit 클라이언트가 로드되지 않았습니다."))
            return
        tickers = context.get("tickers", [])
        if context.get("type") == "price_single":
            tickers = [context.get("ticker", "KRW-BTC")]
        embed = discord.Embed(title="\U0001F50D 상세 정보", color=discord.Color.teal(), timestamp=datetime.now(timezone.utc))
        for ticker in tickers[:5]:
            coin = ticker.replace("KRW-", "")
            try:
                price = await asyncio.to_thread(upbit_client.get_current_price, ticker)
                detail = f"현재가: **{price:,.0f}** KRW" if price else "조회 실패"
                orderbook = await asyncio.to_thread(upbit_client.get_orderbook, ticker)
                if orderbook and isinstance(orderbook, list) and len(orderbook) > 0:
                    ob = orderbook[0] if isinstance(orderbook[0], dict) else orderbook
                    units = ob.get("orderbook_units", [])
                    if units:
                        best_ask = units[0].get("ask_price", 0)
                        best_bid = units[0].get("bid_price", 0)
                        detail += f"\n매도 1호가: {best_ask:,.0f}\n매수 1호가: {best_bid:,.0f}"
                        if best_bid > 0:
                            spread = (best_ask - best_bid) / best_bid * 100
                            detail += f"\n스프레드: {spread:.3f}%"
                embed.add_field(name=f"{coin}", value=detail, inline=True)
            except Exception as e:
                embed.add_field(name=f"{coin}", value=f"조회 실패: {e}", inline=True)
        embed.set_footer(text="JARVIS Crypto | 상세 정보")
        await channel.send(embed=embed)

    async def _reaction_chart(self, channel, context: dict):
        """7일 차트 생성"""
        if not upbit_client:
            await channel.send(embed=_error_embed("Upbit 클라이언트가 로드되지 않았습니다."))
            return
        ticker = "KRW-BTC"
        if context.get("type") == "price_single":
            ticker = context.get("ticker", "KRW-BTC")
        elif context.get("tickers"):
            ticker = context["tickers"][0]
        coin = ticker.replace("KRW-", "")
        try:
            df = await asyncio.to_thread(upbit_client.get_ohlcv, ticker, interval="day", count=7)
            if df is None or df.empty:
                await channel.send(embed=_error_embed(f"{coin} 차트 데이터 없음"))
                return
            try:
                import matplotlib
                matplotlib.use("Agg")
                import matplotlib.pyplot as plt

                fig, ax = plt.subplots(figsize=(8, 4))
                fig.patch.set_facecolor("#2C2F33")
                ax.set_facecolor("#23272A")
                closes = df["close"].values
                dates = [d.strftime("%m/%d") for d in df.index]
                color = "#ED4245" if closes[-1] >= closes[0] else "#5865F2"
                ax.plot(dates, closes, color=color, linewidth=2, marker="o", markersize=4)
                ax.fill_between(dates, closes, alpha=0.1, color=color)
                ax.set_title(f"{coin} 7-Day Price", color="white", fontsize=14)
                ax.tick_params(colors="white")
                ax.spines["bottom"].set_color("white")
                ax.spines["left"].set_color("white")
                ax.spines["top"].set_visible(False)
                ax.spines["right"].set_visible(False)
                for i, v in enumerate(closes):
                    ax.annotate(f"{v:,.0f}", (i, v), textcoords="offset points", xytext=(0, 8), ha="center", fontsize=7, color="white")
                buf = io.BytesIO()
                fig.savefig(buf, format="png", dpi=120, bbox_inches="tight", facecolor=fig.get_facecolor())
                buf.seek(0)
                plt.close(fig)
                file = discord.File(buf, filename=f"{coin}_chart.png")
                embed = discord.Embed(title=f"\U0001F4C8 {coin} 7일 차트", color=discord.Color.dark_gold(), timestamp=datetime.now(timezone.utc))
                embed.set_image(url=f"attachment://{coin}_chart.png")
                embed.set_footer(text="JARVIS Crypto | Upbit")
                await channel.send(embed=embed, file=file)
            except ImportError:
                # 텍스트 차트 폴백
                closes = df["close"].values
                min_p, max_p = min(closes), max(closes)
                chart_width = 20
                lines = []
                for i, row in df.iterrows():
                    close = row["close"]
                    bar_len = int((close - min_p) / (max_p - min_p) * chart_width) if max_p > min_p else chart_width // 2
                    bar = "\u2588" * bar_len + "\u2591" * (chart_width - bar_len)
                    date_str = i.strftime("%m/%d")
                    lines.append(f"`{date_str}` {bar} **{close:,.0f}**")
                embed = discord.Embed(title=f"\U0001F4CA {coin} 7일 차트", description="\n".join(lines), color=discord.Color.dark_gold(), timestamp=datetime.now(timezone.utc))
                embed.set_footer(text="JARVIS Crypto | 텍스트 차트")
                await channel.send(embed=embed)
        except Exception as e:
            logger.error(f"Chart: {e}")
            await channel.send(embed=_error_embed(f"차트 생성 실패: {e}"))

    async def on_message(self, message: discord.Message):
        """모든 메시지 수신 처리. 로컬 명령어 → 에이전트 파이프라인 순서로 라우팅."""
        if message.author.bot: return

        # ── 중복 메시지 처리 방지 (gateway reconnect 등으로 같은 메시지가 여러 번 올 수 있음) ──
        if message.id in self._processed_messages:
            logger.debug(f"[MSG] Duplicate message ignored: {message.id}")
            return
        self._processed_messages.add(message.id)
        # deque(maxlen=200) 자동 evict — evict된 항목을 set에서도 제거
        if len(self._processed_messages_order) == self._processed_messages_order.maxlen:
            evicted = self._processed_messages_order[0]
            self._processed_messages.discard(evicted)
        self._processed_messages_order.append(message.id)

        self._message_count += 1
        logger.debug(f"[MSG] from={message.author} ch={message.channel} content='{_redact_sensitive(message.content)}'")

        ctx = await self.get_context(message)
        if ctx.command:
            self._command_count += 1
            await self.process_commands(message)
            return
        is_dm = isinstance(message.channel, discord.DMChannel)
        is_mention = self.user.mentioned_in(message)
        # "자비스" 키워드도 멘션처럼 처리
        content_lower = message.content.lower().strip()
        is_keyword = any(kw in content_lower for kw in ["자비스", "jarvis"])
        logger.debug(f"[MSG] is_dm={is_dm} is_mention={is_mention} is_keyword={is_keyword} bot_id={self.user.id}")
        if is_dm or is_mention or is_keyword:
            content = message.content.replace(f"<@{self.user.id}>", "").strip()
            # "자비스"/"jarvis" 키워드도 제거하여 순수 프롬프트만 남김
            if is_keyword:
                content = re.sub(r'(?i)자비스|jarvis', '', content).strip()
            # 빈 멘션(@JARVIS만 보낸 경우) → 안내 메시지 응답
            if not content and len(message.attachments) == 0:
                await message.reply("무엇을 도와드릴까요? 질문을 함께 적어주세요.\n예: `@JARVIS 날씨 알려줘`")
                return
            logger.info(f"[MSG] Dispatching to agent: '{content[:80]}'")
            await self._process_agent_request(message, content)

    async def on_command_error(self, ctx, error):
        """User-friendly error handler for prefix commands."""
        if isinstance(error, commands.MissingRequiredArgument):
            await ctx.send(f"필수 인자가 누락되었습니다: `{error.param.name}`\n사용법을 확인해주세요.")
        elif isinstance(error, commands.CommandNotFound):
            pass  # Silently ignore unknown commands
        elif isinstance(error, commands.MissingPermissions):
            await ctx.send("이 명령어를 실행할 권한이 없습니다.")
        elif isinstance(error, commands.CommandOnCooldown):
            await ctx.send(f"명령어 쿨다운 중입니다. {error.retry_after:.1f}초 후 다시 시도해주세요.")
        else:
            logger.error(f"Command error in '{ctx.command}': {error}", exc_info=error)
            await ctx.send(f"명령어 처리 중 오류가 발생했습니다: {error}")

    # ── Helper: 코인 티커 추출 ──
    def _extract_ticker(self, text: str) -> str:
        """텍스트에서 코인 티커 추출. 기본값 KRW-BTC."""

        # KRW-BTC 형식
        m = re.search(r'(KRW-[A-Z]{2,10})', text.upper())
        if m:
            return m.group(1)
        lower = text.lower()
        for kw, ticker in COIN_NAME_MAP.items():
            if kw in lower:
                return ticker
        # 대문자 3~5글자 토큰
        m = re.search(r'\b([A-Z]{3,5})\b', text.upper())
        if m and m.group(1) not in TICKER_EXCLUDE:
            return f"KRW-{m.group(1)}"
        return "KRW-BTC"

    # ── Local Response (API 없이 동작) ──
    async def _try_local_response(self, message: discord.Message, prompt: str) -> bool:
        """키워드 기반 로컬 응답. API 호출 없이 처리 가능하면 True 반환.
        Note: `p` = lowercased prompt for keyword matching; `prompt` = original case for display/parsing.
        """
        p = prompt.lower().strip()
        user_id = str(message.author.id)
        now_kst = datetime.now(timezone(timedelta(hours=9)))
        logger.info(f"[LOCAL] Checking keywords for: '{p[:50]}' | system_mcp={system_mcp_server is not None}")

        # ── CommandDispatcher (점진적 이관 — 등록된 핸들러 우선 실행) ──
        if command_dispatcher:
            result = await command_dispatcher.dispatch(
                prompt=prompt, message=message, bot=self, user_id=user_id,
            )
            if result is not None:
                if isinstance(result, str) and result:
                    await message.reply(result)
                return True

        # ── 도움말/기능 안내 (★ 최우선 — 이미지 생성보다 위) ──
        if _HELP_PATTERN.search(p):
            embed = discord.Embed(
                title="\U0001F916 JARVIS 전체 기능 목록",
                color=discord.Color.blurple(),
                timestamp=datetime.now(timezone.utc),
            )
            embed.add_field(name="\U0001F50D 슬래시 명령", value=(
                "`/price [코인]` - 시세 조회 (Embed)\n"
                "`/balance` - 잔고 확인\n"
                "`/trade <매수|매도> <코인> <금액>` - 매매\n"
                "`/chart [코인]` - 가격 차트"
            ), inline=False)
            embed.add_field(name="\U0001F4BB 시스템/유틸", value=(
                "`시스템` - PC 상태 | `프로세스` - 목록 | `캡처` - 스크린샷\n"
                "`웹캠` - 웹캠 | `속도 측정` - 인터넷 | `네트워크` - IP/포트\n"
                "`파일 *.log` - 파일 검색 | `크롬 열어` - 프로그램 실행\n"
                "`프로세스 1234 종료` - PID 종료 | `예약 목록` - 예약 작업\n"
                "`클립보드` - 복사 | `알림 [내용]` - 알림 | `타이머 5` - 타이머\n"
                "`ssh user@host ls` - SSH 원격 | `mcp 도구 목록` - MCP"
            ), inline=False)
            embed.add_field(name="\U0001F4B0 크립토 매매", value=(
                "`시세 BTC` - 코인 시세 | `호가 ETH` - 호가창\n"
                "`김프` - 프리미엄 | `공포지수` - 공포/탐욕\n"
                "`시장 요약` - 시장 | `분석 BTC` - 기술적 분석\n"
                "`포트폴리오` - 자산 요약 | `거래내역` - 최근 거래\n"
                "`매매통계` - 수익률 | `대기주문` - 미체결 주문\n"
                "`자동매매 시작/중지/상태` - 자동 트레이딩\n"
                "`모의모드`/`실전모드` - 전환 | `스마트매매` - AI 매매\n"
                "`손절 -5 익절 10 설정` - 리스크 | `관심종목 BTC ETH`\n"
                "`가격알림 BTC 50000000` - 알림 | `트레일링 5%`\n"
                "`마켓 목록` - 상장 종목 | `보안 상태` - 보안 체크"
            ), inline=False)
            embed.add_field(name="\U0001F3AE SC2 봇", value=(
                "`전적` - 승률 | `게임상황` - 유닛 현황\n"
                "`로그` - 최근 로그 | `로그 읽어 game.log` - 로그 내용\n"
                "`SC2 테스트` - 연습 게임 | `공격성 공격` - 공격성 조절"
            ), inline=False)
            embed.add_field(name="\U0001F30D 정보/생활", value=(
                "`날씨 서울` - 날씨 | `검색 [키워드]` - 웹검색\n"
                "`운세` - 오늘의 운세 | `번역 hello` - 번역\n"
                "`계산 2+3*4` - 계산 | `시간` - 현재 시간\n"
                "`브리핑` - 모닝 브리핑 | `git` - Git 상태\n"
                "`스마트홈 조명 켜` - IoT 제어"
            ), inline=False)
            embed.add_field(name="\U0001F916 AI 강화", value=(
                "`!그려줘 <설명>` - AI 이미지 생성\n"
                "`!말해줘 <텍스트>` - TTS 음성 변환\n"
                "`!요약 [수]` - 채널 대화 요약\n"
                "`!분석` - 첨부 이미지/PDF AI 분석"
            ), inline=False)
            embed.add_field(name="\U0001F4C8 금융 확장", value=(
                "`!알림 BTC 1억` - 가격 알림 (DM)\n"
                "`!백테스트 BTC RSI 30` - 전략 시뮬\n"
                "`!뉴스감정 BTC` - 뉴스 감정 분석\n"
                "`!환율 USD` - 실시간 환율\n"
                "`!수익리포트` - 주간/월간 리포트\n"
                "`!온체인 BTC` - 온체인 데이터"
            ), inline=False)
            embed.add_field(name="\U0001F3AE SC2 확장", value=(
                "`!중계` - 실시간 게임 중계\n"
                "`!전략분석` - 상대 빌드오더 분석\n"
                "`!대시보드` - 전적 통계 대시보드\n"
                "`!훈련 14:00` - 훈련 스케줄 설정"
            ), inline=False)
            embed.add_field(name="\U0001F6E0 시스템 확장", value=(
                "`!파일보내기 <경로>` - PC→Discord 파일 전송\n"
                "`!파일받기` - Discord→PC 파일 저장\n"
                "`!실행 크롬` - 앱 원격 실행\n"
                "`!모니터링` - 시스템 대시보드\n"
                "`!디스크정리` - 임시 파일 정리\n"
                "`!네트워크` - 네트워크 상태\n"
                "`!스마트홈 조명 켜기` - IoT 제어"
            ), inline=False)
            embed.add_field(name="\u2705 생산성", value=(
                "`!할일 추가/완료/삭제` - 투두 리스트\n"
                "`!알려줘 30분후 회의` - 리마인더\n"
                "`!뽀모 25` - 뽀모도로 타이머\n"
                "`!회의록` - AI 회의록 생성\n"
                "`!습관 추가/체크` - 습관 트래커"
            ), inline=False)
            embed.add_field(name="\U0001F3B2 엔터테인먼트", value=(
                "`!재생 <URL>` - 음악 재생\n"
                "`!가위바위보 가위` - 가위바위보\n"
                "`!숫자맞히기 100` - 숫자 게임\n"
                "`!퀴즈` - 퀴즈\n"
                "`!밈 윗줄|아랫줄` - 밈 생성\n"
                "`!레벨` - 내 레벨/XP\n"
                "`!랭킹` - 서버 랭킹\n"
                "`!챌린지` - 데일리 챌린지"
            ), inline=False)
            embed.add_field(name="\U0001F512 보안/관리", value=(
                "`!접속로그` - PC 로그인 이력\n"
                "`!이상감지설정` - 임계값 설정\n"
                "`!권한` - 역할별 권한 확인\n"
                "`!감사로그` - 명령어 실행 이력"
            ), inline=False)
            embed.add_field(name="\U0001F44D 리액션", value="봇 시세 응답에 \U0001F44D 상세보기 | \U0001F4CA 차트", inline=False)
            embed.set_footer(text="JARVIS v2.0 | 36개 신규 기능 | @멘션 또는 DM")
            await message.reply(embed=embed)
            return True

        # ── 이미지 생성 (그려줘) — Regex 기반, "그림/이미지/draw" 문맥 필수 ──
        if _IMAGE_GEN_PATTERN.search(p):
            draw_prompt = _IMAGE_GEN_PATTERN.sub("", p).strip()
            if not draw_prompt:
                draw_prompt = "a beautiful fantasy creature"
            image_data = None
            # 엔진별 개별 try/except — 하나가 실패해도 다음 엔진으로 안전 fallback
            if not image_data and _generate_image_openai:
                try:
                    image_data = await _generate_image_openai(draw_prompt)
                except Exception as e:
                    logger.warning(f"[IMG] OpenAI engine failed: {e}")
            if not image_data and _generate_image_stable_diffusion:
                try:
                    image_data = await _generate_image_stable_diffusion(draw_prompt)
                except Exception as e:
                    logger.warning(f"[IMG] StableDiffusion engine failed: {e}")
            if not image_data and _generate_image_free:
                try:
                    image_data = await _generate_image_free(draw_prompt)
                except Exception as e:
                    logger.warning(f"[IMG] Pollinations engine failed: {e}")

            if image_data:
                try:
                    file = discord.File(io.BytesIO(image_data), filename="generated.png")
                    embed = discord.Embed(
                        title="AI 이미지 생성 완료",
                        description=f"**프롬프트:** {draw_prompt[:200]}",
                        color=discord.Color.purple(),
                        timestamp=datetime.now(timezone.utc),
                    )
                    embed.set_image(url="attachment://generated.png")
                    await message.reply(embed=embed, file=file)
                except Exception as e:
                    logger.error(f"Image upload error: {e}")
                    await message.reply(f"이미지 생성은 성공했으나 업로드 중 오류: {e}")
                return True
            else:
                await message.reply("이미지 생성 실패. 모든 엔진(DALL-E/SD/Pollinations)이 응답하지 않습니다. 잠시 후 다시 시도해주세요.")
                return True

        # ── SC2 전적 ──
        # "통계"/"승률"은 매매 통계(L1388)와 겹치므로, 매매 컨텍스트가 없을 때만 SC2 전적으로 처리
        _is_trade_context = any(w in p for w in ["매매", "거래", "trade", "수익", "손익", "코인", "btc", "eth"])
        if (any(w in p for w in ["전적", "stats", "sc2"])
                or (any(w in p for w in ["승률", "통계"]) and not _is_trade_context)):
            if sc2_mcp_server:
                try:
                    result = await sc2_mcp_server.sc2_bot_stats()
                    await message.reply(f"**SC2 전적**\n```\n{result}\n```")
                    return True
                except Exception as e:
                    logger.error(f"SC2 stats: {e}")
            await message.reply("SC2 모듈을 불러올 수 없습니다. (`pip install mcp` 필요)")
            return True

        # ── SC2 게임 상황 ──
        if any(w in p for w in ["게임", "유닛", "situation"]) and any(w in p for w in ["상황", "상태", "게임", "situation"]):
            if sc2_mcp_server:
                try:
                    result = await sc2_mcp_server.get_game_situation()
                    if len(result) > 1900:
                        result = result[:1900] + "..."
                    await message.reply(f"**게임 상황**\n```json\n{result}\n```")
                    return True
                except Exception as e:
                    logger.error(f"SC2 situation: {e}")
                    await message.reply(f"게임 상황 조회 실패: {e}")
                    return True

        # ── SC2 로그 ──
        if any(w in p for w in ["로그", "log"]) and len(p) < 20:
            if sc2_mcp_server:
                try:
                    result = await sc2_mcp_server.list_bot_logs()
                    await message.reply(f"**최근 로그**\n```\n{result}\n```")
                    return True
                except Exception as e:
                    logger.error(f"SC2 logs: {e}")
                    await message.reply(f"로그 조회 실패: {e}")
                    return True

        # ── 시스템 상태 ──
        if any(w in p for w in ["시스템", "cpu", "메모리", "리소스", "system", "ram"]) and len(p) < 30:
            if system_mcp_server:
                try:
                    result = await system_mcp_server.system_resources()
                    await message.reply(f"**시스템 리소스**\n```\n{result}\n```")
                    return True
                except Exception as e:
                    logger.error(f"System: {e}")
            # fallback: psutil 직접 사용
            if psutil:
                try:
                    cpu = await asyncio.to_thread(psutil.cpu_percent, 1)
                    mem = psutil.virtual_memory()
                    disk = psutil.disk_usage('/')
                    await message.reply(
                        f"**시스템 상태**\n"
                        f"• CPU: {cpu}%\n"
                        f"• RAM: {mem.percent}% ({mem.used // (1024**3)}/{mem.total // (1024**3)} GB)\n"
                        f"• Disk: {disk.percent}% ({disk.used // (1024**3)}/{disk.total // (1024**3)} GB)"
                    )
                    return True
                except Exception as e:
                    logger.error(f"psutil fallback: {e}")

        # ── 프로세스 ──
        if any(w in p for w in ["프로세스", "process"]):
            if system_mcp_server:
                try:
                    result = await system_mcp_server.list_processes()
                    if len(result) > 1900:
                        result = result[:1900] + "..."
                    await message.reply(f"**프로세스 목록**\n```\n{result}\n```")
                    return True
                except Exception as e:
                    logger.error(f"Process: {e}")

        # ── 스크린샷 ──
        if any(w in p for w in ["스크린샷", "캡처", "screenshot"]) or ("화면" in p and "분석" not in p):
            if system_mcp_server:
                try:
                    res = await system_mcp_server.capture_screenshot()
                    if "base64," in res:
                        b64_str = res.split("base64,", 1)[-1]
                        img_data = base64.b64decode(b64_str)
                        file = discord.File(io.BytesIO(img_data), filename="screenshot.jpg")
                        await message.reply(content="**스크린샷 캡처 완료**", file=file)
                    else:
                        await message.reply(f"캡처 결과: {res[:500]}")
                    return True
                except Exception as e:
                    logger.error(f"Screenshot: {e}")
                    await message.reply(f"스크린샷 실패: {e}")
                    return True

        # ── 코인 시세 ──
        if any(w in p for w in ["시세", "가격", "price", "btc", "eth", "xrp", "sol", "doge", "코인"]) and not any(w in p for w in ["분석", "analyze", "차트분석", "기술적"]):
            if upbit_client:
                try:
                    # 특정 코인 추출 시도
                    symbols = {"btc": "KRW-BTC", "eth": "KRW-ETH", "xrp": "KRW-XRP", "sol": "KRW-SOL", "doge": "KRW-DOGE"}
                    target = None
                    for sym, ticker in symbols.items():
                        if sym in p:
                            target = ticker
                            break

                    if target:
                        price = await asyncio.to_thread(upbit_client.get_current_price, target)
                        coin = target.replace("KRW-", "")
                        if price:
                            await message.reply(f"**{coin} 현재가:** {price:,.0f} KRW")
                        else:
                            await message.reply(f"{coin} 시세를 조회할 수 없습니다.")
                    else:
                        prices = await asyncio.to_thread(upbit_client.get_prices, ["KRW-BTC", "KRW-ETH", "KRW-XRP", "KRW-SOL"])
                        lines = []
                        for ticker, price in prices.items():
                            coin = ticker.replace("KRW-", "")
                            if price and price > 0:
                                lines.append(f"• **{coin}**: {price:,.0f} KRW")
                        await message.reply("**코인 시세**\n" + "\n".join(lines) if lines else "시세 조회 실패")
                    return True
                except Exception as e:
                    logger.error(f"Price: {e}")
                    await message.reply(f"시세 조회 실패: {e}")
                    return True

        # ── 김치프리미엄 ──
        if any(w in p for w in ["김프", "김치프리미엄"]) or ("김치" in p and any(w in p for w in ["프리미엄", "프미", "코인", "비트"])):
            if crypto_mcp_server:
                try:
                    symbol = "KRW-BTC"
                    for w in p.split():
                        if w.upper() in ["BTC", "ETH", "XRP", "SOL"]:
                            symbol = f"KRW-{w.upper()}"
                    result = await crypto_mcp_server.kimchi_premium(symbol)
                    await message.reply(f"**김치프리미엄**\n```\n{result}\n```")
                    return True
                except Exception as e:
                    logger.error(f"Kimchi: {e}")
                    await message.reply(f"김프 조회 실패: {e}")
                    return True

        # ── 공포/탐욕 지수 ──
        if (any(w in p for w in ["탐욕", "fear", "greed"]) or ("공포" in p and any(w in p for w in ["지수", "index", "코인", "시장", "탐욕"]))) and len(p) < 30:
            if crypto_mcp_server:
                try:
                    result = await crypto_mcp_server.fear_greed_index()
                    await message.reply(f"**공포/탐욕 지수**\n```\n{result}\n```")
                    return True
                except Exception as e:
                    logger.error(f"Fear/Greed: {e}")
                    await message.reply(f"공포/탐욕 지수 조회 실패: {e}")
                    return True

        # ── 시장 요약 ──
        if any(w in p for w in ["시장", "market"]) or ("요약" in p and any(w in p for w in ["시장", "코인", "비트", "크립토", "가상화폐", "market"])):
            if crypto_mcp_server:
                try:
                    result = await crypto_mcp_server.market_summary_tool()
                    await message.reply(f"**시장 요약**\n```\n{result}\n```")
                    return True
                except Exception as e:
                    logger.error(f"Market: {e}")
                    await message.reply(f"시장 요약 조회 실패: {e}")
                    return True

        # ── 날씨 ──
        if any(w in p for w in ["날씨", "weather", "기온"]) and len(p) < 30:
            city = None
            # 1. web_tools의 CITY_COORDS를 이용해 도시명 검색
            if web_tools and hasattr(web_tools, 'CITY_COORDS'):
                for known_city in web_tools.CITY_COORDS.keys():
                    if known_city in p:
                        city = known_city
                        break

            # 2. 도시명을 찾지 못한 경우 휴리스틱 (불용어 제거 후 첫 단어)
            if not city:
                words = p.replace("날씨", "").replace("weather", "").replace("기온", "").replace("알려줘", "").replace("알려", "").replace("어때", "").replace("자비스", "").replace(",", "").replace("?", "").strip().split()
                if words:
                    # 첫 단어가 도시명일 확률 높음 (예: "광주 날씨")
                    # 하지만 "오늘", "지금" 같은 단어는 제외
                    ignore_words = ["오늘", "지금", "내일", "모레", "현재", "이번주", "주말", "좀", "해줘", "보여줘"]
                    for w in words:
                        if w not in ignore_words:
                            city = w
                            break

            # 3. 그래도 없으면 기본값 (광주) 또는 안내
            if not city:
                city = "광주" # 기본값

            # 실제 조회
            if web_tools and hasattr(web_tools, 'get_weather'):
                try:
                    result = await asyncio.to_thread(web_tools.get_weather, city)
                    await message.reply(result)
                    return True
                except Exception as e:
                    logger.error(f"web_tools weather: {e}")

            if system_mcp_server:
                try:
                    result = await system_mcp_server.weather(city)
                    await message.reply(result)
                    return True
                except Exception as e:
                    logger.error(f"MCP weather: {e}")

            await message.reply(f"날씨 조회 모듈을 사용할 수 없습니다.")
            return True

        # ── 웹 검색 ──
        if any(w in p for w in ["검색", "search", "찾아"]) and len(p) < 40:
            if web_tools:
                try:
                    query = p.replace("검색", "").replace("search", "").replace("찾아", "").replace("줘", "").strip()
                    if query:
                        result = await asyncio.to_thread(web_tools.search_web, query)
                        if len(result) > 1900:
                            result = result[:1900] + "..."
                        await message.reply(f"**검색 결과: {query}**\n{result}")
                        return True
                except Exception as e:
                    logger.error(f"Search: {e}")
                    await message.reply(f"검색 실패: {e}")
                    return True

        # ── Git 상태 ──
        if any(w in p for w in ["git", "깃"]) and len(p) < 30:
            try:
                res = await asyncio.to_thread(
                    subprocess.check_output, ["git", "status", "--short"],
                    encoding="utf-8", stderr=subprocess.STDOUT,
                    cwd=os.path.dirname(os.path.abspath(__file__))
                )
                if not res.strip():
                    res = "작업 트리가 깨끗합니다 (nothing to commit)"
                if len(res) > 1500:
                    res = res[:1500] + "..."
                await message.reply(f"**Git Status**\n```\n{res}\n```")
                return True
            except Exception as e:
                logger.error(f"Git: {e}")
                await message.reply(f"Git 조회 실패: {e}")
                return True

        # ── 운세 ──
        if any(w in p for w in ["운세", "fortune", "오늘의 운세"]):
            if web_tools and hasattr(web_tools, 'get_daily_fortune'):
                try:
                    result = await asyncio.to_thread(web_tools.get_daily_fortune)
                    await message.reply(result)
                    return True
                except Exception as e:
                    logger.error(f"Fortune: {e}")
            await message.reply("운세 모듈을 사용할 수 없습니다.")
            return True

        # ── 브리핑 ──
        if any(w in p for w in ["브리핑", "briefing", "모닝"]):
            if self._daily_briefing:
                try:
                    result = await self._daily_briefing.generate_briefing_async()
                    if len(str(result)) > 1900:
                        result = str(result)[:1900] + "..."
                    await message.reply(result)
                    return True
                except Exception as e:
                    logger.error(f"Briefing: {e}")
                    await message.reply(f"브리핑 생성 실패: {e}")
                    return True
            await message.reply("브리핑 모듈을 사용할 수 없습니다.")
            return True

        # ── 웹캠 ──
        if any(w in p for w in ["웹캠", "캠", "webcam"]):
            if system_mcp_server:
                try:


                    res = await system_mcp_server.capture_webcam()
                    if "base64," in res:
                        b64_str = res.split("base64,", 1)[-1]
                        img_data = base64.b64decode(b64_str)
                        file = discord.File(io.BytesIO(img_data), filename="webcam.jpg")
                        await message.reply(content="**웹캠 캡처 완료**", file=file)
                    else:
                        await message.reply(f"웹캠 결과: {res[:500]}")
                    return True
                except Exception as e:
                    logger.error(f"Webcam: {e}")
                    await message.reply(f"웹캠 캡처 실패: {e}")
                    return True

        # ── 번역 ──
        if any(w in p for w in ["번역", "translate", "영어로", "한국어로"]):
            if system_mcp_server:
                try:
                    text = p.replace("번역", "").replace("translate", "").replace("영어로", "").replace("한국어로", "").strip()
                    if not text:
                        await message.reply("번역할 텍스트를 입력하세요.\n사용 예: `번역 Hello World`")
                        return True
                    # 명시 키워드 우선, 없으면 자동 감지
                    if "한국어로" in p:
                        target, source = "ko", "en"
                    elif "영어로" in p:
                        target, source = "en", "ko"
                    else:
                        target, source = _detect_language_direction(text)
                    result = await system_mcp_server.translate(text, target, source)
                    await message.reply(f"**번역 결과** ({source}→{target})\n{result}")
                    return True
                except Exception as e:
                    logger.error(f"Translate: {e}")
                    await message.reply(f"번역 실패: {e}")
                    return True

        # ── 계산 ──
        if any(w in p for w in ["계산", "calc", "calculate"]):
            if system_mcp_server:
                try:
                    expr = p.replace("계산", "").replace("calc", "").replace("calculate", "").strip()
                    if not expr:
                        await message.reply("계산할 수식을 입력하세요.\n사용 예: `계산 2+3*4`")
                        return True
                    # 수식이 아닌 자연어인 경우 계산기를 건너뛰고 AI에게 넘김
                    if not re.search(r'\d', expr) and not any(c in expr for c in "+-*/^()."):
                        pass  # fall through to AI handler
                    else:
                        result = await system_mcp_server.calculate(expr)
                        if "계산 실패" in result:
                            pass  # 수식 파싱 실패 시 AI에게 넘김
                        else:
                            await message.reply(f"**계산 결과**\n```\n{result}\n```")
                            return True
                except Exception as e:
                    logger.error(f"Calculate: {e}")
                    # 계산 실패 시 return하지 않고 AI에게 넘김

        # ── 인터넷 속도 ──
        if any(w in p for w in ["속도", "speed", "speedtest", "인터넷"]) and any(w in p for w in ["측정", "test", "테스트", "확인"]):
            if system_mcp_server:
                try:
                    await message.reply("인터넷 속도 측정 중... (30초~1분 소요)")
                    result = await system_mcp_server.check_internet_speed()
                    await message.reply(f"```\n{result}\n```")
                    return True
                except Exception as e:
                    logger.error(f"Speed test: {e}")
                    await message.reply(f"속도 측정 실패: {e}")
                    return True
            await message.reply("system_mcp_server 모듈이 로드되지 않아 속도 측정을 할 수 없습니다.")
            return True

        # ── 네트워크 상태 ──
        if any(w in p for w in ["네트워크", "network", "ip", "포트"]):
            logger.info("[LOCAL] Matched: network keyword")
            if system_mcp_server:
                try:
                    result = await system_mcp_server.network_status()
                    if len(result) > 1900:
                        result = result[:1900] + "..."
                    await message.reply(f"**네트워크 상태**\n```\n{result}\n```")
                    return True
                except Exception as e:
                    logger.error(f"Network via MCP: {e}")
            # fallback: 기본 시스템 명령어
            try:
                hostname = socket.gethostname()
                local_ip = socket.gethostbyname(hostname)
                net = psutil.net_if_addrs() if psutil else {}
                lines = [f"호스트명: {hostname}", f"로컬 IP: {local_ip}", "인터페이스:"]
                for iface, addrs in net.items():
                    for addr in addrs:
                        if addr.family == socket.AF_INET:
                            lines.append(f"  {iface}: {addr.address}")
                await message.reply(f"**네트워크 상태**\n```\n" + "\n".join(lines) + "\n```")
                return True
            except Exception as e2:
                logger.error(f"Network fallback: {e2}")
                await message.reply(f"네트워크 상태 조회 실패: {e2}")
                return True

        # ── 클립보드 ──
        if any(w in p for w in ["클립보드", "clipboard", "붙여넣기"]):
            if system_mcp_server:
                try:
                    result = await system_mcp_server.clipboard_read()
                    await message.reply(f"```\n{result}\n```")
                    return True
                except Exception as e:
                    logger.error(f"Clipboard: {e}")
                    await message.reply(f"클립보드 읽기 실패: {e}")
                    return True

        # ── 알림 (시스템 push 알림) ──
        # 가격/코인 관련 키워드가 있으면 아래 가격 알림 블록으로 넘김
        if (any(w in p for w in ["알림", "notification", "notify"])
                and len(p) < 100
                and not any(w in p for w in ["가격", "price", "이상", "이하", "도달",
                                             "btc", "eth", "xrp", "sol", "doge",
                                             "비트", "이더", "리플", "코인"])):
            if system_mcp_server:
                try:
                    text = p.replace("알림", "").replace("notification", "").replace("notify", "").strip()
                    if not text:
                        text = "JARVIS 알림 테스트"
                    result = await system_mcp_server.send_notification("JARVIS", text)
                    await message.reply(f"{result}")
                    return True
                except Exception as e:
                    logger.error(f"Notification: {e}")
                    await message.reply(f"알림 전송 실패: {e}")
                    return True

        # ── 타이머 ──
        if any(w in p for w in ["타이머", "timer", "알람"]) and len(p) < 50:
            if system_mcp_server:
                try:

                    nums = re.findall(r'(\d+)', p)
                    if nums:
                        minutes = float(nums[0])
                        msg_text = p
                        for n in nums:
                            msg_text = msg_text.replace(n, "").strip()
                        msg_text = msg_text.replace("타이머", "").replace("timer", "").replace("분", "").replace("알람", "").strip() or "타이머 완료"
                        result = await system_mcp_server.set_timer(minutes, msg_text)
                        await message.reply(f"**타이머 설정**\n{result}")
                    else:
                        result = await system_mcp_server.list_timers()
                        await message.reply(f"**타이머 목록**\n```\n{result}\n```")
                    return True
                except Exception as e:
                    logger.error(f"Timer: {e}")
                    await message.reply(f"타이머 설정 실패: {e}")
                    return True

        # ── 시간/날짜 (하순위) ──
        # "오늘"은 단독 또는 시간/날짜 맥락에서만 매칭 (일반 대화 "오늘 보름달이 이쁘다" 등은 AI에게 전달)
        _time_keywords = ["시간", "몇시", "날짜"]
        _time_trigger = any(w in p for w in _time_keywords) or (p.strip() == "오늘") or ("오늘" in p and any(w in p for w in ["몇일", "무슨요일", "날짜", "며칠"]))
        if _time_trigger and len(p) < 30:
            # 다른 키워드(시세, 날씨 등)가 포함되어 있으면 패스
            if any(w in p for w in ["시세", "날씨", "가격", "얼마", "상태", "분석", "공포", "지수"]):
                return False

            time_str = now_kst.strftime("%Y년 %m월 %d일 %A %H:%M:%S")
            weekday_kr = ["월요일", "화요일", "수요일", "목요일", "금요일", "토요일", "일요일"]
            day_name = weekday_kr[now_kst.weekday()]
            await message.reply(f"현재 시간 (KST): **{now_kst.strftime('%Y-%m-%d')} {day_name} {now_kst.strftime('%H:%M:%S')}**")
            return True

        # ── 인사 (하순위) ──
        if any(w in p for w in ["안녕", "hello", "hi", "반가워", "ㅎㅇ"]) and len(p) < 20:
            name = message.author.display_name
            hour = now_kst.hour
            if hour < 12:
                greeting = "좋은 아침이에요"
            elif hour < 18:
                greeting = "좋은 오후예요"
            else:
                greeting = "좋은 저녁이에요"
            await message.reply(
                f"안녕하세요, {name}님! {greeting}.\n"
                f"저는 **JARVIS**입니다. 무엇을 도와드릴까요?\n\n"
                f"**사용 가능한 명령:**\n"
                f"• `전적` / `승률` - SC2 봇 전적\n"
                f"• `시스템` / `CPU` - PC 상태\n"
                f"• `시세` / `BTC` - 코인 시세\n"
                f"• `!scan` / `화면 분석` - 화면 AI 분석\n"
                f"• `!monitor on` - 화면 모니터링 (3분)\n"
                f"• `!cctv snap` - 웹캠 캡처\n"
                f"• `볼륨 올려` / `절전` / `잠금` - PC 제어\n"
                f"• `리플레이` / `코칭` / `래더` - SC2 고급\n"
                f"• `일정` / `메모` - 캘린더/노션\n"
                f"• `도움` - 전체 도움말"
            )
            return True

        # ══════════════════════════════════════════
        # ── 크립토 고급 기능 (crypto_mcp_server) ──
        # ══════════════════════════════════════════

        # ── 호가창 ──
        if any(w in p for w in ["호가", "orderbook", "매수벽", "매도벽"]):
            if crypto_mcp_server:
                try:
                    ticker = self._extract_ticker(p)
                    result = await crypto_mcp_server.coin_orderbook(ticker)
                    if len(result) > SEARCH_RESULT_LIMIT:
                        result = result[:SEARCH_RESULT_LIMIT] + "\n..."
                    await message.reply(f"**호가창 ({ticker})**\n```\n{result}\n```")
                    return True
                except Exception as e:
                    logger.error(f"Orderbook: {e}")
                    await message.reply(f"호가창 조회 실패: {e}")
                    return True

        # ── 마켓 목록 ──
        if any(w in p for w in ["마켓", "종목", "상장"]) and any(w in p for w in ["목록", "리스트", "list", "전체"]):
            if crypto_mcp_server:
                try:
                    result = await crypto_mcp_server.market_list()
                    if len(result) > SEARCH_RESULT_LIMIT:
                        result = result[:SEARCH_RESULT_LIMIT] + "\n..."
                    await message.reply(f"**거래 가능 마켓**\n```\n{result}\n```")
                    return True
                except Exception as e:
                    await message.reply(f"마켓 목록 조회 실패: {e}")
                    return True

        # ── 코인 분석 ──
        if any(w in p for w in ["분석", "analyze", "차트분석", "기술적"]) and any(w in p for w in ["코인", "btc", "eth", "xrp", "sol", "matic", "ada", "doge", "avax", "dot", "link", "시세", "암호화폐", "가상화폐"]):
            if crypto_mcp_server:
                try:
                    ticker = self._extract_ticker(p)
                    if any(w in p for w in ["상세", "detail", "디테일"]):
                        result = await crypto_mcp_server.analyze_coin_detail(ticker)
                    else:
                        result = await crypto_mcp_server.analyze_market(ticker)
                    if len(result) > SEARCH_RESULT_LIMIT:
                        result = result[:SEARCH_RESULT_LIMIT] + "\n..."
                    await message.reply(f"**코인 분석**\n```\n{result}\n```")
                    return True
                except Exception as e:
                    await message.reply(f"코인 분석 실패: {e}")
                    return True

        # ── 포트폴리오 ──
        if any(w in p for w in ["포트폴리오", "portfolio", "보유", "자산"]) and not any(w in p for w in ["시세", "가격"]):
            if crypto_mcp_server:
                try:
                    if any(w in p for w in ["차트", "그래프", "graph"]):
                        result = await crypto_mcp_server.portfolio_graph()
                    elif any(w in p for w in ["비율", "비중", "holdings"]):
                        result = await crypto_mcp_server.holdings_chart()
                    elif any(w in p for w in ["스냅샷", "저장", "snapshot"]):
                        result = await crypto_mcp_server.record_portfolio_snapshot()
                    else:
                        result = await crypto_mcp_server.portfolio_summary()
                    if len(result) > SEARCH_RESULT_LIMIT:
                        result = result[:SEARCH_RESULT_LIMIT] + "\n..."
                    await message.reply(f"**포트폴리오**\n```\n{result}\n```")
                    return True
                except Exception as e:
                    await message.reply(f"포트폴리오 조회 실패: {e}")
                    return True

        # ── 거래 내역 ──
        if any(w in p for w in ["거래내역", "거래기록", "최근거래", "recent", "체결"]) and any(w in p for w in ["거래", "매매", "trade", "내역", "기록", "체결"]):
            if crypto_mcp_server:
                try:
                    result = await crypto_mcp_server.recent_trades()
                    if len(result) > SEARCH_RESULT_LIMIT:
                        result = result[:SEARCH_RESULT_LIMIT] + "\n..."
                    await message.reply(f"**최근 거래 내역**\n```\n{result}\n```")
                    return True
                except Exception as e:
                    await message.reply(f"거래 내역 조회 실패: {e}")
                    return True

        # ── 매매 통계 ──
        if any(w in p for w in ["통계", "statistics", "수익률", "승률"]) and any(w in p for w in ["매매", "거래", "trade", "수익", "손익"]):
            if crypto_mcp_server:
                try:
                    period = "week" if any(w in p for w in ["주간", "이번주", "week"]) else "month" if any(w in p for w in ["월간", "이번달", "month"]) else "all"
                    result = await crypto_mcp_server.trade_statistics(period)
                    await message.reply(f"**매매 통계 ({period})**\n```\n{result}\n```")
                    return True
                except Exception as e:
                    await message.reply(f"매매 통계 조회 실패: {e}")
                    return True

        # ── 대기 주문 ──
        if any(w in p for w in ["대기", "미체결", "pending"]) and any(w in p for w in ["주문", "order", "체결"]):
            if crypto_mcp_server:
                try:
                    result = await crypto_mcp_server.pending_orders()
                    await message.reply(f"**대기 주문**\n```\n{result}\n```")
                    return True
                except Exception as e:
                    await message.reply(f"대기 주문 조회 실패: {e}")
                    return True

        # ── 주문 취소 ──
        if any(w in p for w in ["주문취소", "취소"]) and any(w in p for w in ["주문", "order"]):
            if crypto_mcp_server:
                try:

                    uuid_match = re.search(r'[0-9a-f]{8}-[0-9a-f]{4}-[0-9a-f]{4}-[0-9a-f]{4}-[0-9a-f]{12}', p)
                    if uuid_match:
                        result = await crypto_mcp_server.cancel_my_order(uuid_match.group(0))
                        await message.reply(f"```\n{result}\n```")
                    else:
                        await message.reply("취소할 주문 UUID를 지정해주세요.\n예: `주문 취소 abc12345-...`")
                    return True
                except Exception as e:
                    await message.reply(f"주문 취소 실패: {e}")
                    return True

        # ── 자동 매매 ──
        if any(w in p for w in ["자동매매", "자동거래", "auto", "봇매매"]):
            if crypto_mcp_server:
                try:
                    if any(w in p for w in ["시작", "start", "켜", "on"]):
                        strategy = "smart"
                        result = await crypto_mcp_server.start_auto_trade(strategy)
                    elif any(w in p for w in ["중지", "stop", "꺼", "off", "멈춰"]):
                        result = await crypto_mcp_server.stop_auto_trade()
                    elif any(w in p for w in ["상태", "status"]):
                        result = await crypto_mcp_server.auto_trade_status()
                    else:
                        result = await crypto_mcp_server.auto_trade_status()
                    await message.reply(f"**자동 매매**\n```\n{result}\n```")
                    return True
                except Exception as e:
                    await message.reply(f"자동 매매 처리 실패: {e}")
                    return True

        # ── 매매 모드 (모의/실전) ──
        if any(w in p for w in ["모의", "실전", "드라이", "dry", "live"]) and any(w in p for w in ["모드", "mode", "전환", "설정"]):
            if crypto_mcp_server:
                try:
                    mode = "live" if any(w in p for w in ["실전", "live"]) else "dry"
                    result = await crypto_mcp_server.set_trade_mode(mode)
                    await message.reply(f"```\n{result}\n```")
                    return True
                except Exception as e:
                    await message.reply(f"매매 모드 전환 실패: {e}")
                    return True

        # ── 손절/익절 설정 ──
        if (any(w in p for w in ["손절", "익절", "stoploss", "stop loss", "stop-loss",
                                    "takeprofit", "take profit", "take-profit", "리스크"])
                or re.search(r"stop.?loss|take.?profit", p)) and any(w in p for w in ["설정", "변경", "세팅"]):
            if crypto_mcp_server:
                try:

                    nums = re.findall(r'[-+]?\d+(?:\.\d+)?', prompt)
                    sl = float(nums[0]) if len(nums) >= 1 else -5.0
                    tp = float(nums[1]) if len(nums) >= 2 else 10.0
                    result = await crypto_mcp_server.set_risk_params(stop_loss=sl, take_profit=tp)
                    await message.reply(f"```\n{result}\n```")
                    return True
                except Exception as e:
                    await message.reply(f"리스크 설정 실패: {e}")
                    return True

        # ── 관심종목 ──
        if any(w in p for w in ["관심", "watch", "감시"]) and any(w in p for w in ["종목", "코인", "list", "목록", "설정"]):
            if crypto_mcp_server:
                try:

                    tickers = re.findall(r'KRW-[A-Z]+|[A-Z]{3,}', prompt.upper())
                    if tickers:
                        formatted = ",".join(t if "KRW-" in t else f"KRW-{t}" for t in tickers)
                        result = await crypto_mcp_server.set_watch_list(formatted)
                    else:
                        result = "관심 종목을 지정해주세요.\n예: `관심종목 BTC ETH XRP`"
                    await message.reply(f"```\n{result}\n```")
                    return True
                except Exception as e:
                    await message.reply(f"관심종목 설정 실패: {e}")
                    return True

        # ── 가격 알림 ──
        if any(w in p for w in ["알림", "alert"]) and any(w in p for w in ["가격", "price", "이상", "이하", "도달"]):
            if crypto_mcp_server:
                try:

                    ticker = self._extract_ticker(p)
                    nums = re.findall(r'[\d,]+(?:\.\d+)?', prompt.replace(",", ""))
                    above = float(nums[0]) if len(nums) >= 1 else 0
                    below = float(nums[1]) if len(nums) >= 2 else 0
                    result = await crypto_mcp_server.set_price_alert(ticker, above=above, below=below)
                    await message.reply(f"```\n{result}\n```")
                    return True
                except Exception as e:
                    await message.reply(f"가격 알림 설정 실패: {e}")
                    return True

        # ── 트레일링 스탑 ──
        if any(w in p for w in ["트레일링", "trailing", "추적손절"]):
            if crypto_mcp_server:
                try:

                    ticker = self._extract_ticker(p)
                    pct_match = re.search(r'(\d+(?:\.\d+)?)\s*%?', prompt)
                    trail_pct = float(pct_match.group(1)) if pct_match else 5.0
                    result = await crypto_mcp_server.set_trailing_stop_tool(ticker, trail_pct)
                    await message.reply(f"```\n{result}\n```")
                    return True
                except Exception as e:
                    await message.reply(f"트레일링 스탑 설정 실패: {e}")
                    return True

        # ── 스마트 매매 ──
        if any(w in p for w in ["스마트", "smart"]) and any(w in p for w in ["매매", "trade", "매수", "설정", "모드"]):
            if crypto_mcp_server:
                try:
                    if any(w in p for w in ["설정", "세팅", "모드"]):
                        enabled = not any(w in p for w in ["끄", "off", "비활성"])
                        result = await crypto_mcp_server.set_smart_mode(enabled=enabled)
                    elif any(w in p for w in ["실행", "지금", "now"]):
                        result = await crypto_mcp_server.smart_trade_now()
                    else:
                        result = await crypto_mcp_server.smart_trade_now()
                    await message.reply(f"**스마트 매매**\n```\n{result}\n```")
                    return True
                except Exception as e:
                    await message.reply(f"스마트 매매 실패: {e}")
                    return True

        # ── 매매 사이클 실행 ──
        if any(w in p for w in ["사이클", "cycle"]) and any(w in p for w in ["실행", "run", "돌려"]):
            if crypto_mcp_server:
                try:
                    result = await crypto_mcp_server.run_trade_cycle()
                    await message.reply(f"**매매 사이클**\n```\n{result}\n```")
                    return True
                except Exception as e:
                    await message.reply(f"매매 사이클 실패: {e}")
                    return True

        # ── 안전 한도 설정 ──
        if any(w in p for w in ["안전", "safety", "한도"]) and any(w in p for w in ["설정", "변경", "limit"]):
            if crypto_mcp_server:
                try:
                    result = await crypto_mcp_server.set_safety_limits()
                    await message.reply(f"```\n{result}\n```")
                    return True
                except Exception as e:
                    await message.reply(f"안전 한도 설정 실패: {e}")
                    return True

        # ── 보안 상태 ──
        if any(w in p for w in ["보안", "security"]) and any(w in p for w in ["상태", "status", "확인", "체크"]):
            if crypto_mcp_server:
                try:
                    result = await crypto_mcp_server.security_status()
                    await message.reply(f"**보안 상태**\n```\n{result}\n```")
                    return True
                except Exception as e:
                    await message.reply(f"보안 상태 조회 실패: {e}")
                    return True

        # ── 크립토 도움말 ──
        if any(w in p for w in ["크립토", "crypto"]) and any(w in p for w in ["도움", "help", "명령"]):
            if crypto_mcp_server:
                try:
                    result = await crypto_mcp_server.crypto_help()
                    if len(result) > SEARCH_RESULT_LIMIT:
                        result = result[:SEARCH_RESULT_LIMIT] + "\n..."
                    await message.reply(f"```\n{result}\n```")
                    return True
                except Exception as e:
                    await message.reply(f"크립토 도움말 실패: {e}")
                    return True

        # ══════════════════════════════════════════
        # ── SC2 고급 기능 (sc2_mcp_server) ──
        # ══════════════════════════════════════════

        # ── SC2 로그 내용 읽기 ──
        if any(w in p for w in ["로그내용", "로그읽기", "logfile"]) or ("로그" in p and any(w in p for w in ["읽어", "내용", "보여줘", "열어"])):
            if sc2_mcp_server:
                try:
                    tokens = prompt.strip().split()
                    filename = next((t for t in tokens if "." in t and not t.startswith(".")), "")
                    if filename:
                        result = await sc2_mcp_server.read_log_content(filename)
                    else:
                        result = await sc2_mcp_server.list_bot_logs()
                        result = "파일명을 지정하세요.\n예: `로그 읽어 game_2024.log`\n\n" + result
                    if len(result) > SEARCH_RESULT_LIMIT:
                        result = result[:SEARCH_RESULT_LIMIT] + "\n..."
                    await message.reply(f"**SC2 로그**\n```\n{result}\n```")
                    return True
                except Exception as e:
                    await message.reply(f"로그 읽기 실패: {e}")
                    return True

        # ── SC2 테스트 게임 ──
        if any(w in p for w in ["테스트게임", "연습게임"]) or ("sc2" in p and any(w in p for w in ["테스트", "시작", "연습"])):
            if sc2_mcp_server:
                try:
                    result = await sc2_mcp_server.run_sc2_test_game()
                    await message.reply(f"**SC2 테스트 게임**\n```\n{result}\n```")
                    return True
                except Exception as e:
                    await message.reply(f"SC2 테스트 게임 실행 실패: {e}")
                    return True

        # ── SC2 공격성 조절 ──
        if any(w in p for w in ["공격성", "aggression", "어그로"]) or ("sc2" in p and any(w in p for w in ["공격", "수비", "밸런스"])):
            if sc2_mcp_server:
                try:
                    if any(w in p for w in ["공격", "aggressive", "어그로"]):
                        level = "aggressive"
                    elif any(w in p for w in ["수비", "passive", "방어"]):
                        level = "passive"
                    else:
                        level = "balanced"
                    result = await sc2_mcp_server.set_aggression_level(level)
                    await message.reply(f"**SC2 공격성**\n```\n{result}\n```")
                    return True
                except Exception as e:
                    await message.reply(f"공격성 설정 실패: {e}")
                    return True

        # ══════════════════════════════════════════
        # ── 시스템 고급 기능 ──
        # ══════════════════════════════════════════

        # ── SSH 원격 실행 (관리자 전용) ──
        if any(w in p for w in ["ssh", "원격"]) and any(w in p for w in ["실행", "접속", "명령", "연결"]):
            if not _is_authorized(message):
                await message.reply("🔒 SSH 원격 실행은 관리자 권한이 필요합니다.")
                return True
            if system_mcp_server:
                try:
                    tokens = prompt.strip().split()
                    # ssh user@host command...
                    host = ""
                    user = ""
                    cmd_parts = []
                    for t in tokens:
                        if "@" in t:
                            parts = t.split("@", 1)
                            user = parts[0]
                            host = parts[1]
                        elif not host and ("." in t or ":" in t) and t.lower() not in ("ssh", "원격"):
                            host = t
                        elif host:
                            cmd_parts.append(t)
                    if not host:
                        await message.reply("SSH 접속 정보를 지정하세요.\n예: `ssh user@host ls -la`")
                        return True
                    command = " ".join(cmd_parts) if cmd_parts else "echo connected"
                    # SSH 명령어 화이트리스트 검증 (Fix 2-1)
                    if not _is_ssh_command_safe(command):
                        await message.reply(
                            "보안상 허용되지 않은 SSH 명령어입니다.\n"
                            f"허용 명령어: `{'`, `'.join(sorted(list(_SSH_ALLOWED_COMMANDS)[:15]))}` 등"
                        )
                        return True
                    result = await system_mcp_server.ssh_execute(host, command, user=user)
                    if len(result) > SEARCH_RESULT_LIMIT:
                        result = result[:SEARCH_RESULT_LIMIT] + "\n..."
                    await message.reply(f"**SSH 결과 ({host})**\n```\n{result}\n```")
                    return True
                except Exception as e:
                    await message.reply(f"SSH 실행 실패: {e}")
                    return True

        # ── MCP 도구 목록 ──
        if any(w in p for w in ["mcp", "도구목록", "tool"]) and any(w in p for w in ["목록", "list", "도구", "tool"]):
            if system_mcp_server:
                try:
                    result = await system_mcp_server.list_mcp_tools()
                    if len(result) > SEARCH_RESULT_LIMIT:
                        result = result[:SEARCH_RESULT_LIMIT] + "\n..."
                    await message.reply(f"**MCP 도구 목록**\n```\n{result}\n```")
                    return True
                except Exception as e:
                    await message.reply(f"MCP 도구 목록 조회 실패: {e}")
                    return True

        # ── 스마트홈 ──
        if any(w in p for w in ["스마트홈", "smarthome", "iot"]) or (any(w in p for w in ["조명", "에어컨", "tv", "전등"]) and any(w in p for w in ["켜", "꺼", "on", "off", "설정"])):
            if system_mcp_server:
                try:
                    device = "light"  # default
                    action = "status"
                    for kw, dev in [("조명", "light"), ("전등", "light"), ("에어컨", "ac"), ("tv", "tv"), ("티비", "tv")]:
                        if kw in p:
                            device = dev
                            break
                    if any(w in p for w in ["켜", "on"]):
                        action = "on"
                    elif any(w in p for w in ["꺼", "off"]):
                        action = "off"
                    result = await system_mcp_server.smart_home_control(device, action)
                    await message.reply(f"**스마트홈**\n```\n{result}\n```")
                    return True
                except Exception as e:
                    await message.reply(f"스마트홈 제어 실패: {e}")
                    return True

        # ── 파일 검색 ──
        if any(w in p for w in ["파일", "file", "찾기"]) and system_mcp_server:
            try:

                # "파일 *.log 찾아줘" → pattern=*.log, directory=현재
                # "파일 검색 d:/logs *.txt" → directory=d:/logs, pattern=*.txt
                tokens = prompt.strip().split()
                # 키워드 제거
                filtered = [t for t in tokens if t.lower() not in ("파일", "file", "찾기", "검색", "찾아", "찾아줘", "보여줘", "목록")]
                directory = "."
                pattern = "*"
                for t in filtered:
                    if "/" in t or "\\" in t or ":" in t:
                        directory = t
                    elif "*" in t or "?" in t or "." in t:
                        pattern = t
                    elif t:
                        pattern = f"*{t}*"
                # Path Traversal 차단: ".." 포함된 경로 거부
                if ".." in directory or not _is_path_allowed(directory):
                    await message.reply("보안상 허용되지 않은 경로입니다. 프로젝트 폴더 또는 홈 디렉토리만 검색 가능합니다.")
                    return True
                result = await system_mcp_server.search_files(directory, pattern)
                if len(result) > SEARCH_RESULT_LIMIT:
                    result = result[:SEARCH_RESULT_LIMIT] + "\n... (결과가 잘림)"
                await message.reply(f"**파일 검색 결과**\n```\n{result}\n```")
                return True
            except Exception as e:
                logger.error(f"File search: {e}")
                await message.reply(f"파일 검색 실패: {e}")
                return True

        # ── 프로그램 실행 ──
        # "실행" requires a known program name nearby to avoid overly broad matching
        _prog_names_kr = {"메모장", "노트패드", "계산기", "탐색기", "그림판", "페인트",
                          "크롬", "구글", "파이어폭스", "엣지", "터미널", "코드",
                          "작업관리자", "캡처도구", "파워셸"}
        _prog_names_en = {"notepad", "calc", "explorer", "mspaint", "cmd",
                          "powershell", "code", "chrome", "firefox", "msedge",
                          "taskmgr", "snip", "winterm", "edge", "vscode"}
        _has_prog_name = any(name in p for name in _prog_names_kr) or any(name in p for name in _prog_names_en)
        if (any(w in p for w in ["열어", "run ", "open "]) or ("실행" in p and _has_prog_name)) and system_mcp_server:
            try:
                # "크롬 열어" → name=chrome
                # "메모장 실행해줘" → name=notepad
                _prog_map = {
                    "메모장": "notepad", "노트패드": "notepad",
                    "계산기": "calc", "탐색기": "explorer",
                    "그림판": "mspaint", "페인트": "mspaint",
                    "크롬": "chrome", "구글": "chrome",
                    "파이어폭스": "firefox", "엣지": "msedge", "edge": "msedge",
                    "터미널": "winterm", "코드": "code", "vscode": "code",
                    "작업관리자": "taskmgr", "캡처도구": "snip",
                    "cmd": "cmd", "파워셸": "powershell", "powershell": "powershell",
                }
                tokens = prompt.strip().split()
                prog_name = None
                for t in tokens:
                    tl = t.lower().rstrip("을를이가해줘좀")
                    if tl in _prog_map:
                        prog_name = _prog_map[tl]
                        break
                    # 직접 이름 (notepad, chrome 등)
                    if tl in ("notepad", "calc", "explorer", "mspaint", "cmd",
                              "powershell", "code", "chrome", "firefox", "msedge",
                              "taskmgr", "snip", "winterm"):
                        prog_name = tl
                        break
                if not prog_name:
                    allowed_kr = ", ".join(sorted(_prog_map.keys()))
                    await message.reply(f"실행할 프로그램을 지정해주세요.\n사용 가능: {allowed_kr}")
                    return True
                result = await system_mcp_server.run_program(prog_name)
                await message.reply(f"**프로그램 실행**\n```\n{result}\n```")
                return True
            except Exception as e:
                logger.error(f"Run program: {e}")
                await message.reply(f"프로그램 실행 실패: {e}")
                return True

        # ── 프로세스 종료 (관리자 전용) ──
        if (any(w in p for w in ["종료", "kill"]) and "pid" in p) or ("프로세스" in p and "종료" in p):
            if not _is_authorized(message):
                await message.reply("🔒 프로세스 종료는 관리자 권한이 필요합니다.")
                return True
            if system_mcp_server:
                try:

                    pid_match = re.search(r'\b(\d{2,})\b', prompt)
                    if pid_match:
                        pid = int(pid_match.group(1))
                        # 시스템 중요 프로세스 보호
                        try:
                            _PROTECTED = {"csrss.exe", "smss.exe", "wininit.exe", "services.exe",
                                          "lsass.exe", "svchost.exe", "explorer.exe", "winlogon.exe",
                                          "system", "registry", "dwm.exe"}
                            proc = psutil.Process(pid) if psutil else None
                            if proc and proc.name().lower() in _PROTECTED:
                                await message.reply(f"시스템 핵심 프로세스({proc.name()})는 종료할 수 없습니다.")
                                return True
                        except Exception:
                            pass  # psutil 실패 시 MCP에서 처리
                        result = await system_mcp_server.kill_process(pid)
                        await message.reply(f"**프로세스 종료**\n```\n{result}\n```")
                    else:
                        await message.reply("종료할 프로세스의 PID를 숫자로 지정해주세요.\n예: `프로세스 1234 종료`")
                    return True
                except Exception as e:
                    logger.error(f"Kill process: {e}")
                    await message.reply(f"프로세스 종료 실패: {e}")
                    return True

        # ── 예약 작업 ──
        if any(w in p for w in ["예약", "스케줄", "schedule"]) and system_mcp_server:
            try:
                if any(w in p for w in ["목록", "리스트", "list", "조회"]):
                    result = await system_mcp_server.list_scheduled_tasks()
                    await message.reply(f"**예약 작업 목록**\n```\n{result}\n```")
                    return True
                elif any(w in p for w in ["취소", "삭제", "cancel", "제거"]):

                    tid_match = re.search(r'[a-f0-9]{8}', p)
                    if tid_match:
                        result = await system_mcp_server.cancel_scheduled_task(tid_match.group(0))
                        await message.reply(f"```\n{result}\n```")
                    else:
                        await message.reply("취소할 작업 ID를 지정해주세요.\n예: `예약 취소 a1b2c3d4`")
                    return True
                else:
                    await message.reply(
                        "**예약 작업 명령어**\n"
                        "`예약 목록` - 현재 예약된 작업 조회\n"
                        "`예약 취소 [작업ID]` - 작업 취소\n\n"
                        "새 작업 예약은 AI에게 요청하세요:\n"
                        "예: `매 5분마다 시스템 상태 체크 예약해줘`"
                    )
                    return True
            except Exception as e:
                logger.error(f"Schedule: {e}")
                await message.reply(f"예약 작업 처리 실패: {e}")
                return True

        # ══════════════════════════════════════════
        # ── PC 원격 제어 (관리자 전용) ──
        # ══════════════════════════════════════════
        if any(w in p for w in ["컴퓨터 꺼", "컴퓨터 종료", "shutdown", "pc 꺼", "재시작", "재부팅",
                                 "restart", "reboot", "절전", "sleep", "잠자기", "잠금", "lock",
                                 "볼륨", "소리", "volume", "음량", "밝기", "brightness"]):
            if not _is_authorized(message):
                await message.reply("🔒 PC 원격 제어는 관리자 권한이 필요합니다.")
                return True

        if any(w in p for w in ["컴퓨터 꺼", "컴퓨터 종료", "shutdown", "pc 꺼"]):
            if system_mcp_server and hasattr(system_mcp_server, 'pc_control'):
                await message.reply("60초 후 컴퓨터가 종료됩니다. 취소하려면 `shutdown /a`")
                result = await system_mcp_server.pc_control("shutdown")
                await message.reply(f"```\n{result}\n```")
                return True

        if any(w in p for w in ["재시작", "재부팅", "restart", "reboot"]):
            if system_mcp_server and hasattr(system_mcp_server, 'pc_control'):
                await message.reply("60초 후 PC가 재시작됩니다.")
                result = await system_mcp_server.pc_control("restart")
                await message.reply(f"```\n{result}\n```")
                return True

        if (any(w in p for w in ["절전", "sleep", "잠자기"]) and "모드" in p) or p.strip() in ["절전", "sleep"]:
            if system_mcp_server and hasattr(system_mcp_server, 'pc_control'):
                result = await system_mcp_server.pc_control("sleep")
                await message.reply(f"```\n{result}\n```")
                return True

        if any(w in p for w in ["잠금", "lock", "화면 잠금"]):
            if system_mcp_server and hasattr(system_mcp_server, 'pc_control'):
                result = await system_mcp_server.pc_control("lock")
                await message.reply(f"```\n{result}\n```")
                return True

        if any(w in p for w in ["볼륨", "소리", "volume", "음량"]):
            if system_mcp_server and hasattr(system_mcp_server, 'pc_control'):
                try:
                    if any(w in p for w in ["올려", "up", "높여", "크게"]):
                        result = await system_mcp_server.pc_control("volume_up")
                    elif any(w in p for w in ["내려", "down", "낮춰", "작게"]):
                        result = await system_mcp_server.pc_control("volume_down")
                    elif any(w in p for w in ["음소거", "mute", "뮤트"]):
                        result = await system_mcp_server.pc_control("volume_mute")
                    else:

                        nums = re.findall(r'\d+', p)
                        if nums:
                            result = await system_mcp_server.pc_control("volume_set", nums[0])
                        else:
                            result = await system_mcp_server.pc_control("volume_up")
                    await message.reply(f"**볼륨 제어**\n```\n{result}\n```")
                    return True
                except Exception as e:
                    await message.reply(f"볼륨 제어 실패: {e}")
                    return True

        if any(w in p for w in ["밝기", "brightness"]):
            if system_mcp_server and hasattr(system_mcp_server, 'pc_control'):
                try:

                    nums = re.findall(r'\d+', p)
                    val = nums[0] if nums else "50"
                    result = await system_mcp_server.pc_control("brightness", val)
                    await message.reply(f"**밝기 조절**\n```\n{result}\n```")
                    return True
                except Exception as e:
                    await message.reply(f"밝기 조절 실패: {e}")
                    return True

        # ══════════════════════════════════════════
        # ── SC2 고급 기능 (Phase 4) ──
        # ══════════════════════════════════════════
        if any(w in p for w in ["리플레이", "replay", "전적분석"]):
            if sc2_mcp_server and hasattr(sc2_mcp_server, 'analyze_replay'):
                try:
                    if any(w in p for w in ["목록", "list", "리스트"]):
                        result = await sc2_mcp_server.list_replays()
                    else:
                        result = await sc2_mcp_server.analyze_replay()
                    if len(result) > SEARCH_RESULT_LIMIT:
                        result = result[:SEARCH_RESULT_LIMIT] + "\n..."
                    await message.reply(f"**SC2 리플레이**\n```\n{result}\n```")
                    return True
                except Exception as e:
                    await message.reply(f"리플레이 분석 실패: {e}")
                    return True

        if any(w in p for w in ["코칭", "coaching", "조언"]) and any(w in p for w in ["sc2", "스타", "게임", "코칭"]):
            if sc2_mcp_server and hasattr(sc2_mcp_server, 'sc2_coaching_check'):
                try:
                    result = await sc2_mcp_server.sc2_coaching_check()
                    await message.reply(f"**SC2 코칭**\n```\n{result}\n```")
                    return True
                except Exception as e:
                    await message.reply(f"코칭 분석 실패: {e}")
                    return True

        if any(w in p for w in ["래더", "ladder", "mmr", "랭킹"]):
            if sc2_mcp_server and hasattr(sc2_mcp_server, 'track_ladder'):
                try:
                    # "래더 Player1" → player_name="Player1"
                    tokens = p.replace("래더", "").replace("ladder", "").replace("mmr", "").replace("랭킹", "").strip().split()
                    player = tokens[0] if tokens else ""
                    server = "kr"
                    for t in tokens:
                        if t.lower() in ["kr", "us", "eu", "cn"]:
                            server = t.lower()
                    result = await sc2_mcp_server.track_ladder(player, server)
                    if len(result) > SEARCH_RESULT_LIMIT:
                        result = result[:SEARCH_RESULT_LIMIT] + "\n..."
                    await message.reply(f"**SC2 래더**\n```\n{result}\n```")
                    return True
                except Exception as e:
                    await message.reply(f"래더 조회 실패: {e}")
                    return True

        # ══════════════════════════════════════════
        # ── Google Calendar (Phase 5) ──
        # ══════════════════════════════════════════
        # "스케줄" 제거: L1830 예약 작업(Windows Task Scheduler)과 충돌 방지
        # Google Calendar는 "일정" / "calendar" 키워드로만 진입
        if any(w in p for w in ["일정", "calendar"]) and not any(w in p for w in ["예약", "schedule", "스케줄"]):
            if calendar_integration:
                try:
                    if any(w in p for w in ["추가", "만들", "생성", "등록", "create"]):
                        # AI에게 위임 (날짜/시간 파싱이 복잡)
                        return False
                    elif any(w in p for w in ["이번주", "주간", "upcoming", "다음"]):
                        result = await calendar_integration.get_upcoming_events(7)
                    else:
                        result = await calendar_integration.get_today_events()
                    await message.reply(f"**일정**\n{result}")
                    return True
                except Exception as e:
                    await message.reply(f"일정 조회 실패: {e}")
                    return True

        # ══════════════════════════════════════════
        # ── Notion (Phase 5) ──
        # ══════════════════════════════════════════
        if any(w in p for w in ["메모", "노트", "notion", "기록"]) and not any(w in p for w in ["메모리", "기억"]):
            if notion_integration:
                try:
                    if any(w in p for w in ["저장", "기록해", "적어", "save", "write"]):
                        # "메모 저장 제목: 내용" 형태 파싱
                        text = p
                        for kw in ["메모", "노트", "저장", "기록해", "적어", "save", "write", "notion"]:
                            text = text.replace(kw, "")
                        text = text.strip()
                        if ":" in text:
                            title, content = text.split(":", 1)
                        else:
                            title = text[:30]
                            content = text
                        result = await notion_integration.save_note(title.strip(), content.strip())
                        await message.reply(result)
                        return True
                    elif any(w in p for w in ["검색", "찾아", "search"]):
                        query = p
                        for kw in ["메모", "노트", "검색", "찾아", "search", "notion"]:
                            query = query.replace(kw, "")
                        result = await notion_integration.search_notes(query.strip())
                        await message.reply(result)
                        return True
                    else:
                        result = await notion_integration.list_recent_notes()
                        await message.reply(result)
                        return True
                except Exception as e:
                    await message.reply(f"Notion 처리 실패: {e}")
                    return True

        # ══════════════════════════════════════════
        # ── Vision: !scan 화면 분석 (Phase 2) ──
        # ══════════════════════════════════════════
        if any(w in p for w in ["!scan", "화면 분석", "스캔", "화면분석", "screen scan"]):
            try:
                custom_prompt = p.replace("!scan", "").replace("화면 분석", "").replace("스캔", "").replace("화면분석", "").strip()
                if not custom_prompt:
                    custom_prompt = "이 화면을 분석해서 한국어로 설명해줘. 주요 내용, UI 요소, 텍스트 등을 요약해."
                result = await self._analyze_screen(custom_prompt)
                if len(result) > 1900:
                    result = result[:1900] + "..."
                await message.reply(f"**화면 분석 결과**\n{result}")
                return True
            except Exception as e:
                await message.reply(f"화면 분석 실패: {e}")
                return True

        # ── Vision: !monitor 모니터링 (Phase 2) ──
        # 3분 주기 화면 감시 (tasks.loop 기반 monitor_task 사용)
        if p.startswith("!monitor") or (any(w in p for w in ["모니터링", "감시"]) and any(w in p for w in ["시작", "중지", "on", "off", "start", "stop"])):
            try:
                if any(w in p for w in ["off", "중지", "끄기", "stop"]):
                    self.monitor_enabled = False
                    self.monitor_channel = None
                    if monitor_task.is_running():
                        monitor_task.stop()
                    await message.reply("👁️ **스마트 감시 모드 종료**")
                    return True
                else:
                    # on / start
                    self.monitor_enabled = True
                    self.monitor_channel = message.channel
                    if not monitor_task.is_running():
                        monitor_task.start(self)
                    await message.reply(f"👁️ **스마트 감시 모드 시작** (3분 주기)\n채널: {message.channel.mention}\n종료: `!monitor off`")
                    return True
            except Exception as e:
                await message.reply(f"모니터링 설정 실패: {e}")
                return True

        # ── Vision: !cctv 웹캠 주기 전송 (Phase 2) ──
        if p.startswith("!cctv") or ("cctv" in p and any(w in p for w in ["시작", "중지", "snap", "캡처"])):
            if system_mcp_server:
                try:
                    if any(w in p for w in ["stop", "중지", "끄기", "off"]):
                        if hasattr(self, '_cctv_task') and self._cctv_task is not None:
                            self._cctv_task.cancel()
                            self._cctv_task = None
                            await message.reply("CCTV 모드를 중지했습니다.")
                        else:
                            await message.reply("실행 중인 CCTV가 없습니다.")
                        return True
                    elif any(w in p for w in ["snap", "캡처", "찍어"]):


                        res = await system_mcp_server.capture_webcam()
                        if "base64," in res:
                            b64_str = res.split("base64,", 1)[-1]
                            img_data = base64.b64decode(b64_str)
                            file = discord.File(io.BytesIO(img_data), filename="cctv_snap.jpg")
                            await message.reply(content="**CCTV 스냅샷**", file=file)
                        else:
                            await message.reply(f"웹캠 결과: {res[:500]}")
                        return True
                    else:
                        # start with interval

                        nums = re.findall(r'\d+', p)
                        interval_min = int(nums[0]) if nums else 5
                        interval_sec = interval_min * 60
                        self._cctv_channel = message.channel

                        async def _cctv_loop():
                            _consecutive_failures = 0
                            while True:
                                if not system_mcp_server:
                                    break
                                try:
                                    res = await system_mcp_server.capture_webcam()
                                    if "base64," in res:
                                        b64_str = res.split("base64,", 1)[-1]
                                        img_data = base64.b64decode(b64_str)
                                        now_str = datetime.now(timezone(timedelta(hours=9))).strftime("%H:%M:%S")
                                        file = discord.File(io.BytesIO(img_data), filename=f"cctv_{now_str}.jpg")
                                        await self._cctv_channel.send(content=f"**CCTV [{now_str}]**", file=file)
                                        _consecutive_failures = 0
                                except asyncio.CancelledError:
                                    break
                                except Exception as e:
                                    logger.error(f"CCTV loop error: {e}")
                                    _consecutive_failures += 1
                                    if _consecutive_failures >= 5:
                                        await self._cctv_channel.send("CCTV 5회 연속 실패 — 자동 중단합니다. `!cctv start`로 재시작하세요.")
                                        break
                                await asyncio.sleep(interval_sec)

                        if self._cctv_task and not self._cctv_task.done():
                            self._cctv_task.cancel()
                            try:
                                await self._cctv_task
                            except asyncio.CancelledError:
                                pass
                        self._cctv_task = asyncio.create_task(_cctv_loop())
                        await message.reply(f"CCTV 모드 시작 (간격: {interval_min}분)\n중지: `!cctv stop`\n즉시 캡처: `!cctv snap`")
                        return True
                except Exception as e:
                    await message.reply(f"CCTV 설정 실패: {e}")
                    return True

        return False

    # ── 첨부파일 읽기 헬퍼 ──
    async def _read_attachment(self, att: discord.Attachment) -> str:
        """첨부파일을 다운로드하여 텍스트로 변환. PDF/Excel/Word/PPT/CSV/JSON/코드/텍스트 지원."""
        fname = att.filename or "unknown"
        ext = os.path.splitext(fname)[1].lower()

        try:
            data = await att.read()
        except Exception as e:
            return f"\n[첨부파일: {fname} - 다운로드 실패: {e}]"

        try:
            # ── PDF ──
            if ext == ".pdf":
                try:
                    from PyPDF2 import PdfReader
                    reader = PdfReader(io.BytesIO(data))
                    pages = reader.pages[:20]
                    text = "\n".join(p.extract_text() or "" for p in pages).strip()
                    if not text:
                        return f"\n[첨부파일: {fname} - PDF에서 텍스트를 추출할 수 없습니다 (이미지 PDF일 수 있음)]"
                    if len(text) > ATTACHMENT_TEXT_LIMIT:
                        text = text[:ATTACHMENT_TEXT_LIMIT] + "\n...(이하 생략)"
                    return f"\n\n--- 📄 첨부파일: {fname} ---\n{text}\n--- 끝 ---\n"
                except ImportError:
                    return f"\n[첨부파일: {fname} - PyPDF2 미설치]"

            # ── Excel (.xlsx, .xls) ──
            elif ext in (".xlsx", ".xls"):
                try:
                    from openpyxl import load_workbook
                    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
                    result_parts = []
                    for sheet_name in wb.sheetnames[:5]:
                        ws = wb[sheet_name]
                        rows = []
                        for i, row in enumerate(ws.iter_rows(values_only=True)):
                            if i >= 100:
                                rows.append("...(이하 생략)")
                                break
                            rows.append("\t".join(str(c) if c is not None else "" for c in row))
                        result_parts.append(f"[시트: {sheet_name}]\n" + "\n".join(rows))
                    wb.close()
                    text = "\n\n".join(result_parts)
                    if len(text) > ATTACHMENT_TEXT_LIMIT:
                        text = text[:ATTACHMENT_TEXT_LIMIT] + "\n...(이하 생략)"
                    return f"\n\n--- 📊 첨부파일: {fname} ---\n{text}\n--- 끝 ---\n"
                except ImportError:
                    return f"\n[첨부파일: {fname} - openpyxl 미설치]"

            # ── CSV ──
            elif ext == ".csv":
                if chardet:
                    detected = chardet.detect(data[:10000])
                    encoding = detected.get("encoding", "utf-8") or "utf-8"
                else:
                    encoding = "utf-8"
                try:
                    text_data = data.decode(encoding, errors="replace")
                except Exception:
                    text_data = data.decode("utf-8", errors="replace")
                reader = csv.reader(io.StringIO(text_data))
                rows = []
                for i, row in enumerate(reader):
                    if i >= 100:
                        rows.append("...(이하 생략)")
                        break
                    rows.append("\t".join(row))
                text = "\n".join(rows)
                if len(text) > ATTACHMENT_TEXT_LIMIT:
                    text = text[:ATTACHMENT_TEXT_LIMIT] + "\n...(이하 생략)"
                return f"\n\n--- 📊 첨부파일: {fname} ---\n{text}\n--- 끝 ---\n"

            # ── Word (.docx) ──
            elif ext == ".docx":
                try:
                    from docx import Document
                    doc = Document(io.BytesIO(data))
                    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                    # 테이블도 추출
                    for table in doc.tables[:10]:
                        for row in table.rows:
                            cells = [c.text.strip() for c in row.cells]
                            paragraphs.append("\t".join(cells))
                    text = "\n".join(paragraphs)
                    if len(text) > ATTACHMENT_TEXT_LIMIT:
                        text = text[:ATTACHMENT_TEXT_LIMIT] + "\n...(이하 생략)"
                    return f"\n\n--- 📝 첨부파일: {fname} ---\n{text}\n--- 끝 ---\n"
                except ImportError:
                    return f"\n[첨부파일: {fname} - python-docx 미설치]"

            # ── PowerPoint (.pptx) ──
            elif ext == ".pptx":
                try:
                    from pptx import Presentation
                    prs = Presentation(io.BytesIO(data))
                    slides_text = []
                    for i, slide in enumerate(prs.slides):
                        if i >= 30:
                            slides_text.append("...(이하 슬라이드 생략)")
                            break
                        parts = []
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                for para in shape.text_frame.paragraphs:
                                    t = para.text.strip()
                                    if t:
                                        parts.append(t)
                            if shape.has_table:
                                for row in shape.table.rows:
                                    cells = [c.text.strip() for c in row.cells]
                                    parts.append("\t".join(cells))
                        if parts:
                            slides_text.append(f"[슬라이드 {i+1}]\n" + "\n".join(parts))
                    text = "\n\n".join(slides_text)
                    if len(text) > ATTACHMENT_TEXT_LIMIT:
                        text = text[:ATTACHMENT_TEXT_LIMIT] + "\n...(이하 생략)"
                    return f"\n\n--- 📊 첨부파일: {fname} ---\n{text}\n--- 끝 ---\n"
                except ImportError:
                    return f"\n[첨부파일: {fname} - python-pptx 미설치]"

            # ── JSON ──
            elif ext == ".json":
                try:
                    text_data = data.decode("utf-8", errors="replace")
                    parsed = json.loads(text_data)
                    text = json.dumps(parsed, indent=2, ensure_ascii=False)
                    if len(text) > ATTACHMENT_TEXT_LIMIT:
                        text = text[:ATTACHMENT_TEXT_LIMIT] + "\n...(이하 생략)"
                    return f"\n\n--- 📋 첨부파일: {fname} ---\n{text}\n--- 끝 ---\n"
                except json.JSONDecodeError:
                    text = data.decode("utf-8", errors="replace")[:ATTACHMENT_TEXT_LIMIT]
                    return f"\n\n--- 📋 첨부파일: {fname} (JSON 파싱 실패, 원본) ---\n{text}\n--- 끝 ---\n"

            # ── 코드/텍스트 파일 (알려진 확장자) ──
            else:
                CODE_EXTS = {
                    ".py", ".js", ".ts", ".jsx", ".tsx", ".java", ".c", ".cpp", ".h",
                    ".cs", ".go", ".rs", ".rb", ".php", ".swift", ".kt", ".scala",
                    ".sh", ".bash", ".zsh", ".bat", ".cmd", ".ps1",
                    ".html", ".css", ".scss", ".less", ".xml", ".svg",
                    ".sql", ".r", ".m", ".lua", ".pl", ".dart",
                    ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf",
                    ".md", ".rst", ".txt", ".log", ".env.example",
                    ".gitignore", ".dockerignore", ".editorconfig",
                }
                is_text = ext in CODE_EXTS or (att.content_type and "text" in att.content_type)

                if not is_text:
                    # 바이너리 여부 간단 체크 (null 바이트 비율)
                    sample = data[:8192]
                    null_count = sample.count(b"\x00")
                    is_text = null_count < len(sample) * 0.05  # 5% 미만이면 텍스트

                if is_text:
                    if chardet:
                        detected = chardet.detect(data[:10000])
                        encoding = detected.get("encoding", "utf-8") or "utf-8"
                    else:
                        encoding = "utf-8"
                    try:
                        text = data.decode(encoding, errors="replace")
                    except Exception:
                        text = data.decode("utf-8", errors="replace")
                    if len(text) > ATTACHMENT_TEXT_LIMIT:
                        text = text[:ATTACHMENT_TEXT_LIMIT] + "\n...(이하 생략)"
                    return f"\n\n--- 📎 첨부파일: {fname} ---\n{text}\n--- 끝 ---\n"
                else:
                    return f"\n[첨부파일: {fname} - 바이너리 파일 (텍스트 추출 불가)]"

        except Exception as e:
            logger.error(f"[ATTACHMENT] {fname} 읽기 실패: {e}", exc_info=True)
            return f"\n[첨부파일: {fname} - 읽기 실패: {type(e).__name__}: {e}]"

    # ── Agent Logic ──
    async def _process_agent_request(self, message: discord.Message, prompt: str):
        """에이전트 루프: 프롬프트 → AI 응답 → 도구 호출 → 결과 반영을 반복."""
        # 동일 메시지에 대한 동시 실행 방지 (asyncio Lock dict)
        _key = f"agent_{message.id}"
        if _key not in self._agent_locks:
            self._agent_locks[_key] = asyncio.Lock()
        _agent_lock = self._agent_locks[_key]
        if _agent_lock.locked():
            logger.warning(f"[AGENT] Duplicate agent request ignored for message {message.id}")
            return
        await _agent_lock.acquire()
        try:
            async with message.channel.typing():
                # 0단계: 첨부파일 읽기 (로컬 응답/AI 모두에서 활용)
                image_url = None
                attachment_text = ""
                if message.attachments:
                    for att in message.attachments:
                        if att.content_type and att.content_type.startswith("image/"):
                            image_url = att.url
                        elif att.size and att.size <= 5_000_000:
                            attachment_text += await self._read_attachment(att)
                        else:
                            sz = (att.size or 0) / 1_000_000
                            attachment_text += f"\n[첨부파일: {att.filename} - 파일이 너무 큽니다 ({sz:.1f}MB, 최대 5MB)]"
                    logger.info(f"[AGENT] 첨부파일 처리: image={image_url is not None}, text_len={len(attachment_text)}")

                # 텍스트 없이 첨부파일만 보낸 경우 → 기본 분석 프롬프트 추가
                if not prompt.strip() and (attachment_text or image_url):
                    if image_url and not attachment_text:
                        prompt = "이 이미지를 분석해주세요. 내용을 설명하고 중요한 정보를 알려주세요."
                    else:
                        prompt = "다음 첨부파일의 내용을 분석하고 요약해주세요."

                if attachment_text:
                    prompt = prompt + attachment_text

                # [MEMORY] 대화 히스토리 추가 (로컬 응답 포함 모든 사용자 입력 기록)
                user_id = str(message.author.id)
                if self.memory:
                    self.memory.add_chat_history(user_id, "user", prompt)

                # 1단계: 로컬 키워드 매칭 (API 불필요)
                logger.info(f"[AGENT] Processing: '{prompt[:80]}'")
                try:
                    local_result = await self._try_local_response(message, prompt)
                    logger.info(f"[AGENT] _try_local_response returned: {local_result}")
                    if local_result:
                        return
                except Exception as local_err:
                    logger.error(f"[AGENT] _try_local_response EXCEPTION: {type(local_err).__name__}: {local_err}", exc_info=True)

                # ★ 1.5단계: Agent Router (도메인 기반 지능형 라우팅) ★
                routing_decision = None
                if self.agent_router:
                    try:
                        routing_decision = self.agent_router.route(prompt, user_id, has_image=bool(image_url))
                        logger.info(
                            f"[AGENT_ROUTER] {routing_decision.domain.name} "
                            f"(conf={routing_decision.confidence:.2f}, model={routing_decision.model_hint})"
                        )
                    except Exception as route_err:
                        logger.error(f"[AGENT_ROUTER] Error: {route_err}")

                # 2단계: AI 모델 호출 (API 필요)
                context = self.memory.get_context_string(user_id) if self.memory else ""

                if self.memory:
                    history = self.memory.get_chat_history(user_id)
                else:
                    history = []

                current_prompt = prompt
                _executed_tools: set = set()  # Deduplication: track (tool_name, tool_args)
                _tool_fail_counts: Dict[str, int] = {}  # 도구별 연속 실패 카운터
                _blocked_tools: set = set()  # 2회 실패 시 차단

                for _ in range(5):
                    # history 전달 + 라우팅 결정 전달
                    response = await self._query_hybrid_model(
                        current_prompt, context, history, image_url, user_id,
                        routing=routing_decision,
                    )

                    if not response:
                        await message.reply("⚠️ AI 응답을 받지 못했습니다. 잠시 후 다시 시도해주세요.")
                        return
                    if response.startswith("⚠️") or response.startswith("❌"):
                        await message.reply(response)
                        return

                    if "TOOL:" in response:
                        # TOOL: 마지막 줄에서만 파싱 (prompt injection 방지)
                        last_line = response.rstrip().rsplit("\n", 1)[-1].strip()
                        match = re.match(r"^TOOL:([a-zA-Z0-9_]+)(?::(.*))?", last_line)
                        if match:
                            tool_name = match.group(1).strip()
                            tool_args = match.group(2).strip() if match.group(2) else ""
                            tool_key = (tool_name, tool_args)
                            if tool_key in _executed_tools:
                                logger.warning(f"[AGENT] Duplicate tool call skipped: {tool_name}({tool_args[:50]})")
                                continue
                            _executed_tools.add(tool_key)
                            # ── 도구 인자 검증 (Fix 2-2) ──
                            valid, err_msg = _validate_tool_args(tool_name, tool_args)
                            if not valid:
                                logger.warning(f"[AGENT] Tool args blocked: {tool_name} - {err_msg}")
                                current_prompt = (
                                    f"Tool '{tool_name}' 인자 검증 실패: {err_msg}. "
                                    f"다른 방법을 시도하거나 직접 답변하세요.\n원래 질문: {prompt}"
                                )
                                continue
                            # ── 연속 실패 도구 차단 (Fix 4-1) ──
                            if tool_name in _blocked_tools:
                                current_prompt = (
                                    f"Tool '{tool_name}' 은 이미 2회 연속 실패하여 차단되었습니다. "
                                    f"다른 방법을 시도하거나 직접 답변하세요.\n원래 질문: {prompt}"
                                )
                                continue
                            result = await self._execute_tool(tool_name, tool_args, message, user_id)

                            # ★ REFLECT: 도구 실패 시 자기 수정 패턴 ★
                            is_error = result and ("Error" in result or "실패" in result or "blocked" in result.lower() or "보안 차단" in result)
                            if is_error:
                                _tool_fail_counts[tool_name] = _tool_fail_counts.get(tool_name, 0) + 1
                                if _tool_fail_counts[tool_name] >= 2:
                                    _blocked_tools.add(tool_name)
                                    logger.warning(f"[AGENT] Tool '{tool_name}' blocked after 2 consecutive failures")
                                current_prompt = (
                                    f"이전 요청: {current_prompt}\n"
                                    f"도구 '{tool_name}' 실패: {result}\n"
                                    f"자기반성: 도구 호출이 실패했습니다. 실패 원인을 분석하고 다른 방법을 시도하거나 직접 답변하세요."
                                )
                            else:
                                current_prompt = f"이전 요청: {current_prompt}\n도구 '{tool_name}' 결과: {result}\n(위 결과를 바탕으로 한국어로 답변하세요)"
                            if len(current_prompt) > PROMPT_TRUNCATE_LIMIT:
                                current_prompt = current_prompt[:2000] + "\n...(중략)...\n" + current_prompt[-4000:]
                            image_url = None
                            continue

                    clean_response = response.split("TOOL:")[0].strip()
                    if clean_response:
                        # [MEMORY] 봇 응답 저장
                        if self.memory:
                            self.memory.add_chat_history(user_id, "assistant", clean_response)

                        if len(clean_response) > DISCORD_MSG_LIMIT:
                            for i in range(0, len(clean_response), 2000):
                                await message.reply(clean_response[i:i+2000])
                        else:
                            await message.reply(clean_response)
                    return

                await message.reply("⚠️ 루프 초과 (생각이 너무 많습니다)")

        except Exception as e:
            logger.error(f"Agent Error: {e}", exc_info=True)
            await message.reply("❌ 내부 오류가 발생했습니다. 잠시 후 다시 시도해주세요.")
        finally:
            _agent_lock.release()
            self._agent_locks.pop(_key, None)
            if create_thread_for_long_conversation and not isinstance(message.channel, discord.Thread):
                try:
                    await create_thread_for_long_conversation(message)
                except Exception as thread_err:
                    logger.debug(f"Thread creation failed: {thread_err}")

    def _is_rate_limited(self, model_name: str) -> bool:
        """모델이 현재 쿨다운 중인지 확인"""
        _cleanup_rate_limits()  # 만료 엔트리 정리
        cooldown_until = _rate_limit_cooldowns.get(model_name)
        if cooldown_until and datetime.now(timezone.utc) < cooldown_until:
            return True
        return False

    async def _analyze_screen(self, prompt="이 화면을 분석해서 한국어로 설명해줘. 주요 내용, UI 요소, 텍스트 등을 요약해."):
        if not GEMINI_API_KEY:
            return "Gemini API Key가 설정되지 않아 화면 분석을 할 수 없습니다."
        if not system_mcp_server:
            return "시스템 모듈을 불러올 수 없습니다."
        try:
            b64_str = await system_mcp_server.capture_screenshot()
            if "base64," in b64_str:
                img_data_b64 = b64_str.split("base64,", 1)[-1]
                url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash:generateContent"
                payload = {
                    "contents": [{"parts": [
                        {"text": prompt},
                        {"inline_data": {"mime_type": "image/jpeg", "data": img_data_b64}}
                    ]}]
                }
                session = await self._get_http_session()
                async with session.post(url, json=payload, headers={"Content-Type": "application/json", "x-goog-api-key": GEMINI_API_KEY}, timeout=aiohttp.ClientTimeout(total=HTTP_TIMEOUT_GEMINI)) as resp:
                    if resp.status == 200:
                        data = await resp.json()
                        candidates = data.get('candidates') or []
                        if not candidates:
                            block_reason = data.get('promptFeedback', {}).get('blockReason', 'UNKNOWN')
                            logger.warning(f"[Gemini] No candidates. blockReason={block_reason}")
                            return f"Gemini 분석 결과 없음 (사유: {block_reason})"
                        text = candidates[0].get('content', {}).get('parts', [{}])[0].get('text')
                        if not text:
                            finish_reason = candidates[0].get('finishReason', 'UNKNOWN')
                            logger.warning(f"[Gemini] Empty text. finishReason={finish_reason}")
                            return f"Gemini 분석 텍스트 비어있음 (finishReason: {finish_reason})"
                        return text
                    else:
                        error_body = await resp.text()
                        logger.error(f"[Gemini] HTTP {resp.status}: {error_body[:300]}")
                        return f"Gemini 분석 실패 (HTTP {resp.status})"
            return "스크린샷 데이터 형식 오류 (base64 헤더 없음)"
        except Exception as e:
            return f"화면 분석 중 오류 발생: {e}"

    def _set_rate_limit(self, model_name: str, seconds: int = None):
        """모델에 쿨다운 설정"""
        cd = seconds or RATE_LIMIT_COOLDOWN_SECONDS
        _rate_limit_cooldowns[model_name] = datetime.now(timezone.utc) + timedelta(seconds=cd)
        logger.warning(f"[RateLimit] {model_name} 쿨다운 {cd}초 설정")

    def _build_fallback_system_prompt(self, context: str, history_str: str) -> str:
        """라우터 미사용 시 기본 시스템 프롬프트 (전체 도구 포함)"""
        try:
            return build_system_prompt(
                role="default", context=context,
                history_str=history_str, categories=None, domain=None,
            )
        except Exception:
            now_kst = datetime.now(timezone(timedelta(hours=9))).strftime("%Y-%m-%d %H:%M:%S (KST)")
            return (
                "[절대 지시사항]\n"
                "당신은 J.A.R.V.I.S., 장선우 사령관의 AI 부관입니다.\n"
                "호칭: 사령관님. 합쇼체, 군대식, 통신 프로토콜 준수.\n"
                "★ 민간 호칭(사장님, 고객님 등) 절대 금지. 이모지 사용 금지. ★\n"
                f"현재 시간(KST): {now_kst}\n"
                "도구 사용 시 **오직** `TOOL:도구명:인자` 만 출력. 사족 금지.\n"
                f"사용자 컨텍스트: {context}\n{history_str}\n"
            )

    async def _query_hybrid_model(self, prompt: str, context: str, history: list, image_url: str, user_id: str, routing=None) -> Optional[str]:
        """Proxy/직접 API를 통한 Claude 모델 쿼리. 캐스케이드 폴백 지원."""
        # ★ 페르소나 강제 주입 — 프록시가 system 메시지를 무시해도 user 메시지로 전달 ★
        prompt = _PERSONA_PREAMBLE + prompt
        # 히스토리 문자열 변환
        history_str = ""
        if history:
            history_str = "대화 내역:\n"
            for msg in history[-HISTORY_TURNS:]:
                role = "User" if msg['role'] == 'user' else "Jarvis"
                history_str += f"{role}: {msg['content'][:200]}\n"

        # ★ AgentRouter 기반 도메인별 시스템 프롬프트 생성 ★
        if routing is not None:
            try:
                system_prompt = build_system_prompt(
                    role=routing.role,
                    context=context,
                    history_str=history_str,
                    categories=routing.tool_categories,
                    domain=routing.domain.value,
                )
            except Exception:
                # fallback: 기존 기본 프롬프트
                system_prompt = self._build_fallback_system_prompt(context, history_str)
        else:
            system_prompt = self._build_fallback_system_prompt(context, history_str)

        # ★ 페르소나 가드: system 프롬프트에 "사령관" 없으면 강제 주입 ★
        if "사령관" not in system_prompt:
            from system_prompts import _JARVIS_IDENTITY
            system_prompt = _JARVIS_IDENTITY + "\n" + system_prompt

        # ★ Model Selector: 모델 힌트 기반 최적 캐스케이드 계획 ★
        model_hint = routing.model_hint if routing else "sonnet"
        model_plan = None
        if self.model_selector:
            try:
                model_plan = self.model_selector.select(
                    model_hint=model_hint,
                    image_required=bool(image_url),
                )
            except Exception as e:
                logger.warning(f"[MODEL_SELECTOR] Error: {e}")

        errors = []

        # 1. Claude API (Official) — ModelSelector 기반 모델 순서 적용
        # Claude는 이미지 분석(Vision)도 지원하므로 image_url 유무와 관계없이 시도
        claude_models = (model_plan.claude_models if model_plan
                         else ["claude-haiku-4-5-20251001", "claude-sonnet-4-5-20250929"])
        if CLAUDE_API_KEY and claude_models:
            all_claude_cooled = all(
                self._is_rate_limited(f"claude-{m}") for m in claude_models
            )
            if not all_claude_cooled:
                for attempt in range(2):
                    try:
                        result = await self._query_claude_api(
                            prompt, system_prompt, models=claude_models,
                            image_url=image_url,
                        )
                        if result:
                            return result
                    except RateLimitError as e:
                        errors.append(f"Claude API: 한도 초과 (쿨다운 {e.retry_after or 60}초)")
                        break
                    except PermissionError as e:
                        errors.append(f"Claude API: {e}")
                        break
                    except Exception as e:
                        errors.append(f"Claude API: {e}")
                        if attempt == 0:
                            await asyncio.sleep(1)
                            continue
                        break
            else:
                errors.append("Claude API: 모든 모델 쿨다운 중")

        # 2. Claude Proxy (Session Key) — model_hint 전달
        proxy_model = model_plan.proxy_model if model_plan else "sonnet"
        if CLAUDE_SESSION_KEY and not image_url and not self._is_rate_limited("claude-proxy"):
            try:
                return await self._query_claude_proxy(
                    prompt, user_id, system_prompt, model=proxy_model,
                )
            except Exception as e:
                errors.append(f"Claude Proxy: {e}")

        # 3. Gemini Fallback — ModelSelector 기반 모델 순서 적용
        gemini_models = (model_plan.gemini_models if model_plan
                         else ["gemini-2.5-flash", "gemini-2.0-flash", "gemini-2.0-flash-lite"])
        if GEMINI_API_KEY:
            for model_name in gemini_models:
                if self._is_rate_limited(f"gemini-{model_name}"):
                    errors.append(f"{model_name}: 쿨다운 중")
                    continue
                try:
                    return await self._query_gemini(prompt, system_prompt, image_url, model=model_name)
                except RateLimitError as e:
                    self._set_rate_limit(f"gemini-{model_name}", e.retry_after or 30)
                    errors.append(f"{model_name}: 한도 초과")
                except PermissionError as e:
                    errors.append(f"Gemini: {e}")
                    break
                except Exception as e:
                    errors.append(f"{model_name}: {e}")

        # 4. Ollama (로컬 - 한도 없음)
        try:
            return await self._query_ollama(prompt, system_prompt)
        except Exception as e:
            errors.append(f"Ollama: {e}")

        # 에러 상세 로그
        logger.warning(f"[AI] All models failed: {errors}")

        # 에러 원인 분류
        error_str = " | ".join(errors)
        has_auth_error = any("키 오류" in e or "401" in e or "403" in e for e in errors)
        has_rate_limit = any("한도 초과" in e or "쿨다운" in e for e in errors)
        has_connection = any("연결 오류" in e or "Connection" in e or "Ollama" in e for e in errors)

        if has_auth_error:
            error_msg = "API 키 오류입니다. `.env.jarvis` 파일의 API 키를 확인해주세요."
        elif has_rate_limit:
            error_msg = "AI 모델 한도 초과입니다. 약 1분 후 자동 복구됩니다."
        elif has_connection:
            error_msg = "AI 서비스 연결 실패입니다. 네트워크를 확인해주세요."
        else:
            error_msg = "AI 모델이 일시적으로 응답하지 않습니다."

        error_msg += "\n\n**AI 없이 사용 가능한 기능:**\n"
        error_msg += "• `날씨 [지역]` • `시세`/`BTC` • `시스템` • `전적` • `시간` • `도움`"

        logger.debug(f"[AI] Error details: {error_str[:300]}")

        return error_msg

    async def _query_claude_api(self, prompt, system, models=None, image_url=None):
        # ModelSelector가 결정한 모델 순서 사용, 없으면 기본값
        claude_models = models or ["claude-haiku-4-5-20251001", "claude-sonnet-4-5-20250929"]
        headers = {"x-api-key": CLAUDE_API_KEY, "anthropic-version": "2023-06-01", "content-type": "application/json"}

        # 이미지가 있으면 multimodal content 구성
        user_content = prompt
        if image_url:
            try:
                session = await self._get_http_session()
                async with session.get(image_url, timeout=aiohttp.ClientTimeout(total=15)) as img_resp:
                    if img_resp.status == 200:
                        img_data = await img_resp.read()
                        b64_img = base64.b64encode(img_data).decode('utf-8')
                        mime = img_resp.content_type or "image/jpeg"
                        user_content = [
                            {"type": "image", "source": {"type": "base64", "media_type": mime, "data": b64_img}},
                            {"type": "text", "text": prompt},
                        ]
            except Exception as e:
                logger.warning(f"[Claude] 이미지 다운로드 실패: {e} — 텍스트만 전송")
                user_content = f"[NOTE: 이미지 다운로드 실패 — 이미지 분석 불가]\n{prompt}"

        last_error = None
        for model in claude_models:
            if self._is_rate_limited(f"claude-{model}"):
                continue
            data = {
                "model": model,
                "max_tokens": 2048,
                "messages": [{"role": "user", "content": user_content}],
                "system": system
            }
            start_time = datetime.now(timezone.utc)
            try:
                session = await self._get_http_session()
                async with session.post(
                    "https://api.anthropic.com/v1/messages",
                    json=data, headers=headers,
                    timeout=aiohttp.ClientTimeout(total=HTTP_TIMEOUT_GEMINI)
                ) as resp:
                    elapsed = (datetime.now(timezone.utc) - start_time).total_seconds() * 1000
                    if resp.status == 200:
                        res = await resp.json()
                        _track_model_result(f"claude-{model}", True, elapsed)
                        logger.info(f"[Claude] {model} 응답 성공 ({elapsed:.0f}ms)")
                        content = res.get('content') or []
                        if content and 'text' in content[0]:
                            return content[0]['text']
                        return str(res)
                    elif resp.status == 429:
                        retry_after = _safe_retry_after(resp.headers, 60)
                        self._set_rate_limit(f"claude-{model}", retry_after)
                        _track_model_result(f"claude-{model}", False)
                        last_error = RateLimitError(f"{model} 한도 초과", retry_after)
                        continue
                    elif resp.status in (401, 403):
                        error_text = await resp.text()
                        _track_model_result(f"claude-{model}", False)
                        # API 키 오류 → 다른 모델도 동일 실패하므로 즉시 중단
                        raise PermissionError(f"Claude API 키 오류 ({resp.status}): {error_text[:150]}")
                    elif resp.status == 529:
                        _track_model_result(f"claude-{model}", False)
                        last_error = Exception(f"{model} 서버 과부하 (529)")
                        await asyncio.sleep(2)
                        continue
                    else:
                        error_text = await resp.text()
                        _track_model_result(f"claude-{model}", False)
                        last_error = Exception(f"{model} Error {resp.status}: {error_text[:200]}")
                        continue
            except asyncio.TimeoutError:
                _track_model_result(f"claude-{model}", False)
                last_error = Exception(f"{model} 타임아웃 (45초)")
                continue
            except aiohttp.ClientError as e:
                _track_model_result(f"claude-{model}", False)
                last_error = Exception(f"{model} 연결 오류: {e}")
                continue

        if isinstance(last_error, RateLimitError):
            raise last_error
        if last_error:
            raise last_error
        raise Exception("사용 가능한 Claude 모델 없음")

    async def _query_claude_proxy(self, prompt, user_id, system="", model="sonnet"):
        # Local Proxy (claude_proxy.js) - 8780 포트 (8765는 mcp_gateway_proxy가 사용)
        proxy_port = os.environ.get("JARVIS_PORT", "8780")
        url = f"http://localhost:{proxy_port}/chat"
        payload = {
            "message": prompt,
            "user": user_id,
            "model": model,
            "system": system
        }
        try:
            session = await self._get_http_session()
            async with session.post(url, json=payload, timeout=aiohttp.ClientTimeout(total=HTTP_TIMEOUT_DEFAULT)) as resp:
                if resp.status == 200:
                    data = await resp.json()
                    # model이 "none"이면 Proxy 내부에서 모든 모델 실패 → 폴백 필요
                    if data.get('model') == 'none' or data.get('error'):
                        raise Exception("Claude Proxy: 모든 upstream 모델 실패")
                    reply = data.get('reply')
                    if reply:
                        return reply
                    raise Exception("Claude Proxy: 빈 응답")
                raise Exception(f"Claude Proxy Error {resp.status}")
        except (aiohttp.ClientConnectorError, aiohttp.ClientError,
                asyncio.TimeoutError, OSError) as e:
            self._set_rate_limit("claude-proxy", 300)
            logger.warning(f"[AI] Claude Proxy 연결 실패 (300초 쿨다운): {e}")
            raise

    async def _query_gemini(self, prompt, system, image_url=None, model="gemini-2.0-flash"):
        url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent"
        parts = [{"text": system + "\nUser: " + prompt}]
        if image_url:
            try:
                session = await self._get_http_session()
                async with session.get(image_url) as resp:
                    if resp.status == 200:
                        img_data = await resp.read()
                        b64_img = base64.b64encode(img_data).decode('utf-8')
                        mime = resp.content_type or "image/jpeg"
                        parts.append({"inline_data": {"mime_type": mime, "data": b64_img}})
            except Exception as e:
                logger.warning(f"[Gemini] 이미지 다운로드 실패: {e}")
                parts.append({"text": "[NOTE: 이미지 다운로드 실패 — 이미지 분석 불가]"})
        payload = {"contents": [{"parts": parts}]}
        session = await self._get_http_session()
        async with session.post(
            url, json=payload,
            headers={"Content-Type": "application/json", "x-goog-api-key": GEMINI_API_KEY},
            timeout=aiohttp.ClientTimeout(total=30)
        ) as resp:
            if resp.status == 200:
                data = await resp.json()
                candidates = data.get('candidates') or []
                if candidates:
                    logger.info(f"[Gemini] {model} 응답 성공")
                    parts_list = candidates[0].get('content', {}).get('parts', [])
                    text = parts_list[0].get('text') if parts_list else None
                    if text:
                        return text
                    raise Exception("Gemini 응답에 텍스트 없음")
                else:
                    raise Exception("No candidates returned")
            elif resp.status == 429:
                retry_after = _safe_retry_after(resp.headers, 30)
                raise RateLimitError(f"{model} 한도 초과", retry_after)
            elif resp.status in (400, 401, 403):
                error_text = await resp.text()
                # 인증/키 오류 → 다른 Gemini 모델도 동일 실패하므로 즉시 중단
                raise PermissionError(f"Gemini API 키 오류 ({resp.status}): {error_text[:150]}")
            elif resp.status == 503:
                raise Exception(f"{model} 서비스 일시 불가 (503)")
            else:
                error_text = await resp.text()
                raise Exception(f"Gemini {model} Error {resp.status}: {error_text[:200]}")

    async def _query_ollama(self, prompt, system):
        # Ollama가 쿨다운 중이면 skip (이전 연결 실패 시 5분 쿨다운)
        if self._is_rate_limited("ollama-local"):
            raise Exception("Ollama 서버 미응답 (쿨다운 중)")
        url = "http://localhost:11434/api/generate"
        payload = {"model": os.environ.get("OLLAMA_MODEL", "qwen2.5:1.5b"), "prompt": system + "\nUser: " + prompt, "stream": False}
        try:
            session = await self._get_http_session()
            async with session.post(
                url, json=payload,
                timeout=aiohttp.ClientTimeout(total=30, connect=5)  # connect 5초로 빠른 실패
            ) as resp:
                if resp.status == 200:
                    data = await resp.json()
                    logger.info("[Ollama] 응답 성공")
                    return data.get('response') or str(data)
                raise Exception(f"Ollama Error {resp.status}")
        except (aiohttp.ClientConnectorError, asyncio.TimeoutError, OSError) as e:
            # Ollama 서버 미실행 → 300초(5분) 쿨다운 (매번 연결 시도 방지)
            self._set_rate_limit("ollama-local", 300)
            logger.warning(f"[Ollama] 서버 미응답, 5분 쿨다운 설정: {e}")
            raise Exception(f"Ollama 서버 연결 실패 (5분 후 재시도): {e}")

    _DANGEROUS_TOOLS = {"ssh_execute", "kill_process", "pc_control", "run_program", "restart_bot", "execute_terminal_command", "execute_python_code", "pc_mouse_click", "write_file", "edit_file", "openclaw_coding", "openclaw_browser", "openclaw_cron", "openclaw_email"}

    async def _execute_tool(self, name, args, message, user_id):
        """도구명과 인자를 받아 해당 도구 실행. 결과 문자열 반환."""
        logger.info(f"Tool Exec: {name}(***)")
        # Authorization check for dangerous tools
        if name in self._DANGEROUS_TOOLS:
            if not _is_authorized(message):
                return f"권한 부족: '{name}' 도구는 관리자만 사용할 수 있습니다."

        _tool_start = time.time()
        _tool_success = True
        _tool_error = ""
        try:
            # ── ToolDispatcher (점진적 이관 — 등록된 도구 우선 실행) ──
            if tool_dispatcher:
                result = await tool_dispatcher.dispatch(
                    name=name, args=args, message=message, user_id=user_id, bot=self,
                )
                if result is not None:
                    return result

            if name == "search_web":
                if web_tools:
                    return await asyncio.to_thread(web_tools.search_web, args)
                return "web_tools 모듈을 불러올 수 없습니다."
            elif name == "get_weather":
                # 실제 날씨 조회 (web_tools 또는 system_mcp_server)
                city = args.strip() if args.strip() else "서울"
                if web_tools and hasattr(web_tools, 'get_weather'):
                    return await asyncio.to_thread(web_tools.get_weather, city)
                if system_mcp_server:
                    return await system_mcp_server.weather(city)
                return f"{city} 날씨 조회 모듈을 사용할 수 없습니다."
            elif name == "get_sc2_stats":
                if sc2_mcp_server:
                    return await sc2_mcp_server.sc2_bot_stats()
                return "SC2 모듈을 불러올 수 없습니다."
            elif name == "get_system_status":
                if system_mcp_server:
                    return await system_mcp_server.system_resources()
                return "시스템 모듈을 불러올 수 없습니다."
            elif name == "get_crypto_price":
                if crypto_mcp_server:
                    return await crypto_mcp_server.coin_price(args)
                return "크립토 모듈을 불러올 수 없습니다."
            elif name == "scan_screen":
                return await self._analyze_screen()
            elif name == "remember":
                if not self.memory:
                    return "메모리 모듈을 불러올 수 없습니다."
                if "|" in args: k, v = args.split("|", 1)
                else: k, v = "info", args
                self.memory.update_user_memory(user_id, k, v)
                return "Saved."
            elif name == "check_git":
                cmd = ["git", "status"]
                allowed_args = args.split()
                if allowed_args:
                    subcmd = allowed_args[0]
                    if subcmd in ["status", "log", "diff", "show", "branch"]:
                        cmd = ["git"] + allowed_args
                    else: return "Blocked Git command."
                try:
                    res = await asyncio.to_thread(
                        subprocess.check_output, cmd,
                        encoding="utf-8", stderr=subprocess.STDOUT,
                        cwd=os.path.dirname(os.path.abspath(__file__))
                    )
                    if len(res) > 1500: res = res[:1500] + "\n...(truncated)"
                    return f"Git Output:\n{res}"
                except subprocess.CalledProcessError as e: return f"Git Error: {e.output}"
            elif name == "get_fortune":
                if web_tools and hasattr(web_tools, 'get_daily_fortune'):
                    return await asyncio.to_thread(web_tools.get_daily_fortune)
                return "운세 모듈을 불러올 수 없습니다."
            elif name == "translate":
                if system_mcp_server:
                    parts = args.split("|", 2) if "|" in args else [args]
                    text = parts[0].strip()
                    if len(parts) > 1:
                        target = parts[1].strip()
                        source = parts[2].strip() if len(parts) > 2 else ("ko" if target == "en" else "en")
                    else:
                        target, source = _detect_language_direction(text)
                    return await system_mcp_server.translate(text, target, source)
                return "번역 모듈을 불러올 수 없습니다."
            elif name == "calculate":
                if system_mcp_server:
                    return await system_mcp_server.calculate(args)
                return "계산기 모듈을 불러올 수 없습니다."
            elif name == "list_features":
                return (
                    "**JARVIS 사용 가능 기능:**\n"
                    "• 날씨, 검색, 시세, 김프, 공포/탐욕, 시장, 호가창\n"
                    "• 전적, 게임상황, 로그, 테스트게임, 공격성\n"
                    "• 시스템, 프로세스, 캡처, 웹캠, 파일검색\n"
                    "• 운세, 번역, 계산, 속도측정, 프로그램실행\n"
                    "• 네트워크, Git, 브리핑, SSH, MCP도구\n"
                    "• 포트폴리오, 거래내역, 매매통계, 자동매매\n"
                    "• 스마트매매, 손절/익절, 관심종목, 가격알림\n"
                    "• 스마트홈, 예약작업, 메모리(기억)\n"
                    "• **!scan** - 화면 AI 분석 (Gemini Vision)\n"
                    "• **!monitor on/off** - 화면 모니터링 (3분 주기)\n"
                    "• **!cctv start/stop/snap** - 웹캠 CCTV\n"
                    "• 볼륨/밝기/절전/잠금/종료 - PC 원격제어\n"
                    "• 리플레이/코칭/래더 - SC2 고급 분석\n"
                    "• 일정(Google Calendar)/메모(Notion) - 비서"
                )
            # ── 새 AI 도구들 ──
            elif name == "portfolio":
                if crypto_mcp_server:
                    return await crypto_mcp_server.portfolio_summary()
                return "크립토 모듈을 불러올 수 없습니다."
            elif name == "recent_trades":
                if crypto_mcp_server:
                    return await crypto_mcp_server.recent_trades()
                return "크립토 모듈을 불러올 수 없습니다."
            elif name == "trade_stats":
                if crypto_mcp_server:
                    period = args.strip() if args.strip() else "all"
                    return await crypto_mcp_server.trade_statistics(period)
                return "크립토 모듈을 불러올 수 없습니다."
            elif name == "analyze_coin":
                if crypto_mcp_server:
                    ticker = args.strip() if args.strip() else "KRW-BTC"
                    return await crypto_mcp_server.analyze_coin_detail(ticker)
                return "크립토 모듈을 불러올 수 없습니다."
            elif name == "auto_trade":
                if crypto_mcp_server:
                    cmd = args.strip().lower()
                    if cmd == "start":
                        return await crypto_mcp_server.start_auto_trade()
                    elif cmd == "stop":
                        return await crypto_mcp_server.stop_auto_trade()
                    else:
                        return await crypto_mcp_server.auto_trade_status()
                return "크립토 모듈을 불러올 수 없습니다."
            elif name == "pending_orders":
                if crypto_mcp_server:
                    return await crypto_mcp_server.pending_orders()
                return "크립토 모듈을 불러올 수 없습니다."
            elif name == "price_alert":
                # 형식: BTC|60000000 (설정), list (목록), BTC|clear (해제)
                parts = args.split("|", 1) if "|" in args else [args.strip(), ""]
                coin_arg = parts[0].strip().upper()
                value_arg = parts[1].strip() if len(parts) > 1 else ""
                async with _price_alert_lock:
                    if coin_arg == "LIST":
                        user_alerts = _price_alert_store.get(user_id, {})
                        if not user_alerts:
                            return "등록된 가격 알림이 없습니다."
                        lines = [f"{t.replace('KRW-','')}: {p:,.0f} KRW" for t, p in user_alerts.items()]
                        return "현재 가격 알림:\n" + "\n".join(lines)
                    ticker = f"KRW-{coin_arg}" if not coin_arg.startswith("KRW-") else coin_arg
                    if value_arg.lower() == "clear":
                        if user_id in _price_alert_store and ticker in _price_alert_store[user_id]:
                            del _price_alert_store[user_id][ticker]
                            if not _price_alert_store[user_id]:
                                del _price_alert_store[user_id]
                            return f"{coin_arg} 가격 알림이 해제되었습니다."
                        return f"{coin_arg} 가격 알림이 없습니다."
                    try:
                        target_price = float(value_arg.replace(",", ""))
                    except ValueError:
                        return "형식: TOOL:price_alert:BTC|60000000 (설정) / list (목록) / BTC|clear (해제)"
                    _price_alert_store.setdefault(user_id, {})[ticker] = target_price
                    return f"{coin_arg} 가격 알림 설정: {target_price:,.0f} KRW 도달 시 DM 알림"
            elif name == "run_program":
                if system_mcp_server:
                    return await system_mcp_server.run_program(args.strip())
                return "시스템 모듈을 불러올 수 없습니다."
            elif name == "kill_process":
                if system_mcp_server:
                    try:
                        pid = int(args.strip())
                        return await system_mcp_server.kill_process(pid)
                    except ValueError:
                        return "PID는 숫자여야 합니다."
                return "시스템 모듈을 불러올 수 없습니다."
            elif name == "search_files":
                if system_mcp_server:
                    parts = args.split("|", 1) if "|" in args else [".", args.strip()]
                    directory = parts[0].strip()
                    pattern = parts[1].strip() if len(parts) > 1 else "*"
                    if not _is_path_allowed(directory):
                        return "보안 차단: 허용되지 않은 경로입니다. 프로젝트 폴더 또는 홈 디렉토리만 검색 가능합니다."
                    return await system_mcp_server.search_files(directory, pattern)
                return "시스템 모듈을 불러올 수 없습니다."
            elif name == "ssh_execute":
                if system_mcp_server:
                    parts = args.split("|", 1) if "|" in args else [args.strip(), "echo connected"]
                    host_part = parts[0].strip()
                    command = parts[1].strip() if len(parts) > 1 else "echo connected"
                    user, host = ("", host_part)
                    if "@" in host_part:
                        user, host = host_part.split("@", 1)
                    return await system_mcp_server.ssh_execute(host, command, user=user)
                return "시스템 모듈을 불러올 수 없습니다."
            # ── OpenClaw Agentic Tools ──
            elif name == "execute_terminal_command":
                if agentic_mcp_server:
                    return await agentic_mcp_server.execute_terminal_command(args)
                return "Agentic 모듈을 불러올 수 없습니다."
            elif name == "execute_python_code":
                is_safe, reason = _ast_check_python_code(args)
                if not is_safe:
                    return f"보안 차단: {reason}. 이 코드 패턴은 허용되지 않습니다."
                if agentic_mcp_server:
                    return await agentic_mcp_server.execute_python_code(args)
                return "Agentic 모듈을 불러올 수 없습니다."
            elif name == "pc_mouse_move":
                if agentic_mcp_server:
                    parts = args.split("|")
                    try:
                        x = int(parts[0].strip())
                        y = int(parts[1].strip())
                        duration = float(parts[2].strip()) if len(parts) > 2 else 0.5
                        return await agentic_mcp_server.computer_use_mouse_move(x, y, duration)
                    except (ValueError, IndexError):
                        return "좌표값이 올바르지 않습니다. (x|y)"
                return "Agentic 모듈을 불러올 수 없습니다."
            elif name == "pc_mouse_click":
                if agentic_mcp_server:
                    parts = args.split("|")
                    btn = parts[0].strip() if len(parts) > 0 and parts[0].strip() else "left"
                    clicks = int(parts[1].strip()) if len(parts) > 1 else 1
                    return await agentic_mcp_server.computer_use_mouse_click(button=btn, clicks=clicks)
                return "Agentic 모듈을 불러올 수 없습니다."
            elif name == "pc_keyboard_type":
                if agentic_mcp_server:
                    return await agentic_mcp_server.computer_use_keyboard_type(args)
                return "Agentic 모듈을 불러올 수 없습니다."
            elif name == "pc_keyboard_press":
                if agentic_mcp_server:
                    return await agentic_mcp_server.computer_use_keyboard_press(args)
                return "Agentic 모듈을 불러올 수 없습니다."
            elif name == "read_file":
                if agentic_mcp_server:
                    return await agentic_mcp_server.read_file(args.strip())
                return "Agentic 모듈을 불러올 수 없습니다."
            elif name == "write_file":
                if agentic_mcp_server:
                    parts = args.split("|", 1) if "|" in args else [args.strip(), ""]
                    path = parts[0].strip()
                    content = parts[1].strip() if len(parts) > 1 else ""
                    return await agentic_mcp_server.write_file(path, content)
                return "Agentic 모듈을 불러올 수 없습니다."
            elif name == "edit_file":
                if agentic_mcp_server:
                    parts = args.split("|", 2)
                    path = parts[0].strip() if len(parts) > 0 else ""
                    old_t = parts[1].strip() if len(parts) > 1 else ""
                    new_t = parts[2].strip() if len(parts) > 2 else ""
                    return await agentic_mcp_server.edit_file(path, old_t, new_t)
                return "Agentic 모듈을 불러올 수 없습니다."
            elif name == "list_dir":
                if agentic_mcp_server:
                    path = args.strip() if args.strip() else "."
                    return await agentic_mcp_server.list_directory(path)
                return "Agentic 모듈을 불러올 수 없습니다."
            # ── Phase 3: PC 제어 ──
            elif name == "pc_control":
                if system_mcp_server and hasattr(system_mcp_server, 'pc_control'):
                    parts = args.split("|", 1) if "|" in args else [args.strip(), ""]
                    action = parts[0].strip()
                    value = parts[1].strip() if len(parts) > 1 else ""
                    return await system_mcp_server.pc_control(action, value)
                return "PC 제어 모듈을 불러올 수 없습니다."
            # ── Phase 4: SC2 고급 ──
            elif name == "analyze_replay":
                if sc2_mcp_server and hasattr(sc2_mcp_server, 'analyze_replay'):
                    return await sc2_mcp_server.analyze_replay(args.strip() if args.strip() else None)
                return "SC2 리플레이 모듈을 불러올 수 없습니다."
            elif name == "list_replays":
                if sc2_mcp_server and hasattr(sc2_mcp_server, 'list_replays'):
                    return await sc2_mcp_server.list_replays()
                return "SC2 리플레이 모듈을 불러올 수 없습니다."
            elif name == "sc2_coaching":
                if sc2_mcp_server and hasattr(sc2_mcp_server, 'sc2_coaching_check'):
                    return await sc2_mcp_server.sc2_coaching_check()
                return "SC2 코칭 모듈을 불러올 수 없습니다."
            elif name == "track_ladder":
                if sc2_mcp_server and hasattr(sc2_mcp_server, 'track_ladder'):
                    parts = args.split("|", 1) if "|" in args else [args.strip(), "kr"]
                    player = parts[0].strip()
                    server = parts[1].strip() if len(parts) > 1 else "kr"
                    return await sc2_mcp_server.track_ladder(player, server)
                return "SC2 래더 모듈을 불러올 수 없습니다."
            # ── Phase 5: Calendar ──
            elif name == "get_today_events":
                if calendar_integration:
                    return await calendar_integration.get_today_events()
                return "캘린더 모듈을 불러올 수 없습니다."
            elif name == "get_upcoming_events":
                if calendar_integration:
                    days = int(args.strip()) if args.strip().isdigit() else 7
                    return await calendar_integration.get_upcoming_events(days)
                return "캘린더 모듈을 불러올 수 없습니다."
            elif name == "create_event":
                if calendar_integration:
                    parts = args.split("|")
                    title = parts[0].strip() if len(parts) > 0 else "새 일정"
                    start = parts[1].strip() if len(parts) > 1 else ""
                    end = parts[2].strip() if len(parts) > 2 else ""
                    return await calendar_integration.create_event(title, start, end)
                return "캘린더 모듈을 불러올 수 없습니다."
            # ── Phase 5: Notion ──
            elif name == "save_note":
                if notion_integration:
                    parts = args.split("|", 1) if "|" in args else [args.strip()[:30], args.strip()]
                    title = parts[0].strip()
                    content = parts[1].strip() if len(parts) > 1 else title
                    return await notion_integration.save_note(title, content)
                return "Notion 모듈을 불러올 수 없습니다."
            elif name == "search_notes":
                if notion_integration:
                    return await notion_integration.search_notes(args.strip())
                return "Notion 모듈을 불러올 수 없습니다."
            elif name == "list_notes":
                if notion_integration:
                    return await notion_integration.list_recent_notes()
                return "Notion 모듈을 불러올 수 없습니다."
            # ── 신규 도구: 봇 상태 ──
            elif name == "bot_status":
                uptime = self.get_uptime()
                model_info = []
                for m, s in _model_stats.items():
                    rate = (s["success"] / s["calls"] * 100) if s["calls"] > 0 else 0
                    avg_ms = (s["total_ms"] / s["success"]) if s["success"] > 0 else 0
                    model_info.append(f"  {m}: {rate:.0f}% 성공 ({s['calls']}회, 평균 {avg_ms:.0f}ms)")

                # ★ Tool Registry 통계 통합 ★
                tool_summary = ""
                if get_tool_registry:
                    tool_summary = f"\n**도구 사용 통계:**\n{get_tool_registry().get_summary()}"

                return (
                    f"**JARVIS 상태**\n"
                    f"• 가동시간: {uptime}\n"
                    f"• 처리 메시지: {self._message_count:,}개\n"
                    f"• 실행 명령어: {self._command_count:,}개\n"
                    f"• 메모리 사용자: {len(self.memory.get_all_users()) if self.memory and hasattr(self.memory, 'get_all_users') else '?'}명\n"
                    f"**AI 모델 통계:**\n" + ("\n".join(model_info) if model_info else "  통계 없음")
                    + tool_summary
                )
            # ── 신규 도구: 메모리 검색 ──
            elif name == "search_memory":
                if not self.memory:
                    return "메모리 모듈을 불러올 수 없습니다."
                if self.memory and hasattr(self.memory, 'search_memory'):
                    results = self.memory.search_memory(args)
                    if results:
                        lines = [f"• [{r['user_id']}] {r['key']}: {r['value']}" for r in results[:10]]
                        return "메모리 검색 결과:\n" + "\n".join(lines)
                    return "검색 결과가 없습니다."
                return "메모리 모듈을 불러올 수 없습니다."
            # ── 신규 도구: 메모리 백업 ──
            elif name == "backup_memory":
                if self.memory and hasattr(self.memory, 'backup'):
                    return self.memory.backup()
                return "메모리 모듈을 불러올 수 없습니다."

            # ══════════════════════════════════════════════════════
            # OpenClaw Tools (AI → Discord 자동 호출)
            # ══════════════════════════════════════════════════════
            elif name.startswith("openclaw_"):
                try:
                    from utils.openclaw_helper import get_openclaw_helper
                    oc = get_openclaw_helper()
                except ImportError:
                    return "OpenClaw 헬퍼를 불러올 수 없습니다."
                if not oc.available:
                    return "OpenClaw CLI가 설치되어 있지 않습니다."

                skill_map = {
                    "openclaw_weather": lambda a: f"Get weather for {a.strip() or '서울'}",
                    "openclaw_youtube": lambda a: f"Search YouTube: {a.strip()}",
                    "openclaw_summarize": lambda a: f"Summarize: {a.strip()}",
                    "openclaw_news": lambda a: f"Get latest news about {a.strip() or 'technology'}",
                    "openclaw_exchange": lambda a: f"Get exchange rate {a.replace('|', ' to ')}" if a.strip() else "Get USD to KRW exchange rate",
                    "openclaw_stock": lambda a: f"Analyze stock {a.strip() or 'AAPL'} with full scoring",
                    "openclaw_image_gen": lambda a: f"Generate image: {a.strip()}",
                    "openclaw_email": lambda a: f"Send email to {a}" if "|" in a else f"Check emails: {a}",
                    "openclaw_notion": lambda a: f"Notion: {a.replace('|', ' - ')}",
                    "openclaw_github": lambda a: f"GitHub {a.replace('|', ' for repo ')}",
                    "openclaw_coding": lambda a: f"Write code: {a.strip()}",
                    "openclaw_browser": lambda a: f"Browser: open {a.split('|')[0].strip()}" if a.strip() else "Browser status",
                    "openclaw_cron": lambda a: f"Cron: {a.replace('|', ' ')}",
                    "openclaw_transcribe": lambda a: f"Transcribe audio: {a.strip()}",
                    "openclaw_calendar": lambda a: f"Calendar: {'show today events' if a.strip() in ('today', 'none', '') else a.strip()}",
                }
                msg_fn = skill_map.get(name)
                if msg_fn is None:
                    return f"Unknown OpenClaw tool: {name}"
                skill_msg = msg_fn(args)
                oc.record_skill_usage(name)
                timeout = 90 if name in ("openclaw_coding", "openclaw_image_gen", "openclaw_transcribe") else 45
                return await oc.run_skill(skill_msg, timeout=timeout)

            return "Unknown Tool"
        except Exception as e:
            _tool_success = False
            _tool_error = str(e)[:200]
            return f"Tool Error: {e}"
        finally:
            # ★ 도구 감사 로그 기록 ★
            if get_tool_registry:
                _elapsed = (time.time() - _tool_start) * 1000
                get_tool_registry().record_call(name, user_id, _tool_success, _elapsed, _tool_error)

    # ── Legacy Prefix Commands ──
    @commands.command(name="scan")
    async def scan_cmd(self, ctx):
        msg = await ctx.send("Scanning...")
        res = await self._execute_tool("scan_screen", "", ctx.message, str(ctx.author.id))
        await msg.edit(content=f"**Analysis Result:**\n{res}")

    @commands.command(name="search")
    async def search_cmd(self, ctx, *, query=""):
        if not query:
            await ctx.send("검색어를 입력해주세요.\n사용법: `!search <검색어>`")
            return
        if web_tools:
            res = await asyncio.to_thread(web_tools.search_web, query)
            await ctx.send(res)
        else:
            await ctx.send("검색 모듈을 불러올 수 없습니다.")

    @commands.command(name="git")
    async def git_cmd(self, ctx):
        try:
            res = await asyncio.to_thread(
                subprocess.check_output, ["git", "status"],
                encoding="utf-8", stderr=subprocess.STDOUT,
                cwd=os.path.dirname(os.path.abspath(__file__))
            )
            if len(res) > 900: res = res[:900] + "..."
            await ctx.send(f"```\n{res}\n```")
        except Exception as e: await ctx.send(f"Git Error: {e}")

    @commands.command(name="briefing")
    async def briefing_cmd(self, ctx):
        """수동으로 모닝 브리핑 생성"""
        if not self._daily_briefing:
            await ctx.send("브리핑 모듈을 불러올 수 없습니다.")
            return
        await ctx.send("모닝 브리핑 생성 중...")
        try:
            report = await self._daily_briefing.generate_briefing_async()
            if len(report) > DISCORD_MSG_LIMIT:
                for i in range(0, len(report), 2000):
                    await ctx.send(report[i:i+2000])
            else:
                await ctx.send(report)
        except Exception as e:
            await ctx.send(f"브리핑 생성 실패: {e}")

    @commands.command(name="voice")
    async def voice_cmd(self, ctx):
        """음성 채널 제어 패널"""
        view = VoiceControlView(self)
        await ctx.send("🔊 **Voice Control Panel**", view=view)

    @commands.command(name="play")
    async def play_cmd(self, ctx, url: str):
        """유튜브 음악 재생"""
        if not ctx.voice_client:
            if ctx.author.voice:
                await ctx.author.voice.channel.connect()
            else:
                await ctx.send("음성 채널에 먼저 입장해주세요.")
                return

        msg = await ctx.send(f"🎵 검색 중: {url}...")

        ydl_opts = {'format': 'bestaudio/best', 'noplaylist': True, 'quiet': True}
        try:
            import yt_dlp
            with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                info = await asyncio.to_thread(ydl.extract_info, url, download=False)
                url2 = info['url']
                title = info['title']
            source = discord.FFmpegPCMAudio(url2, before_options="-reconnect 1 -reconnect_streamed 1 -reconnect_delay_max 5")
            if ctx.voice_client.is_playing():
                ctx.voice_client.stop()
            ctx.voice_client.play(source)
            await msg.edit(content=f"▶️ 재생 시작: **{title}**")
        except Exception as e:
            await msg.edit(content=f"재생 실패: {e}")

    @commands.command(name="stop")
    async def stop_cmd(self, ctx):
        if ctx.voice_client and ctx.voice_client.is_playing():
            ctx.voice_client.stop()
            await ctx.send("⏹️ 재생 정지")

    @commands.command(name="monitor")
    async def monitor_cmd(self, ctx, status: str = "on"):
        """화면 감시 모드 (on/off)"""
        status = status.lower()
        if status == "on":
            self.monitor_enabled = True
            self.monitor_channel = ctx.channel
            if not monitor_task.is_running():
                monitor_task.start(self)
            await ctx.send(f"👁️ **스마트 감시 모드 시작** (3분 주기)\n채널: {ctx.channel.mention}")
        else:
            self.monitor_enabled = False
            self.monitor_channel = None
            if monitor_task.is_running():
                monitor_task.stop()
            await ctx.send("👁️ **스마트 감시 모드 종료**")

# ═══════════════════════════════════════════════════════════════
#  Slash Commands (/price, /balance, /trade, /chart)
# ═══════════════════════════════════════════════════════════════

@app_commands.command(name="price", description="코인 시세를 조회합니다")
@app_commands.describe(coin="조회할 코인 심볼 (예: BTC, ETH). 비우면 관심 코인 전체")
async def price_slash(interaction: discord.Interaction, coin: Optional[str] = None):
    """코인 시세 조회 (Embed + Reaction)"""
    await interaction.response.defer()
    try:
        if not upbit_client:
            await interaction.followup.send(embed=_error_embed("Upbit 클라이언트를 불러올 수 없습니다."))
            return
        if coin:
            ticker = coin.upper()
            if not ticker.startswith("KRW-"):
                ticker = f"KRW-{ticker}"
            price = await asyncio.to_thread(upbit_client.get_current_price, ticker)
            if price is None:
                await interaction.followup.send(embed=_error_embed(f"`{ticker}` 시세를 조회할 수 없습니다."))
                return
            change_pct = None
            try:
                df = await asyncio.to_thread(upbit_client.get_ohlcv, ticker, interval="day", count=2)
                if df is not None and len(df) >= 2:
                    prev_close = df["close"].iloc[-2]
                    if prev_close > 0:
                        change_pct = (price - prev_close) / prev_close * 100
            except Exception:
                pass
            embed = _price_embed(ticker, price, change_pct)
            sent = await interaction.followup.send(embed=embed, wait=True)
            bot_instance = interaction.client
            if isinstance(bot_instance, JarvisBot):
                bot_instance._add_reaction_context(sent.id, {"type": "price_single", "ticker": ticker, "tickers": [ticker]})
            await sent.add_reaction("\U0001F44D")
            await sent.add_reaction("\U0001F4CA")
        else:
            default_tickers = ["KRW-BTC", "KRW-ETH", "KRW-XRP", "KRW-SOL"]
            if upbit_config and hasattr(upbit_config, 'DEFAULT_WATCH_LIST'):
                default_tickers = list(upbit_config.DEFAULT_WATCH_LIST)
            prices = await asyncio.to_thread(upbit_client.get_prices, default_tickers)
            embed = _multi_price_embed(prices)
            sent = await interaction.followup.send(embed=embed, wait=True)
            bot_instance = interaction.client
            if isinstance(bot_instance, JarvisBot):
                bot_instance._add_reaction_context(sent.id, {"type": "price_multi", "tickers": default_tickers})
            await sent.add_reaction("\U0001F44D")
            await sent.add_reaction("\U0001F4CA")
    except Exception as e:
        logger.error(f"/price: {e}")
        await interaction.followup.send(embed=_error_embed(f"시세 조회 실패: {e}"))


@app_commands.command(name="balance", description="포트폴리오 잔고를 확인합니다")
async def balance_slash(interaction: discord.Interaction):
    """잔고 조회 (Embed, ephemeral)"""
    await interaction.response.defer(ephemeral=True)
    try:
        if not upbit_client:
            await interaction.followup.send(embed=_error_embed("Upbit 클라이언트를 불러올 수 없습니다."), ephemeral=True)
            return
        balances = await asyncio.to_thread(upbit_client.get_balances)
        total = await asyncio.to_thread(upbit_client.get_total_balance_krw)
        embed = _balance_embed(balances, total)
        await interaction.followup.send(embed=embed, ephemeral=True)
    except Exception as e:
        logger.error(f"/balance: {e}")
        await interaction.followup.send(embed=_error_embed(f"잔고 조회 실패: {e}"), ephemeral=True)


@app_commands.command(name="trade", description="코인을 매수 또는 매도합니다")
@app_commands.describe(action="매매 유형", coin="코인 심볼 (예: BTC)", amount="금액 (매수: KRW, 매도: 수량)")
@app_commands.choices(action=[
    app_commands.Choice(name="매수 (Buy)", value="buy"),
    app_commands.Choice(name="매도 (Sell)", value="sell"),
])
async def trade_slash_cmd(interaction: discord.Interaction, action: app_commands.Choice[str], coin: str, amount: float):
    """매매 실행 (버튼 UI 포함)"""
    await interaction.response.defer()

    # 권한 체크: 봇 오너, 관리자, 또는 Trader 역할 필요
    has_perm = False
    if BOT_OWNER_ID and str(interaction.user.id) == BOT_OWNER_ID:
        has_perm = True
    elif hasattr(interaction.user, "guild_permissions") and interaction.user.guild_permissions.administrator:
        has_perm = True
    elif hasattr(interaction.user, "roles"):
        user_roles = {r.name.lower() for r in interaction.user.roles}
        if TRADER_ROLE_NAME.lower() in user_roles or (user_roles & ADMIN_ROLE_NAMES):
            has_perm = True
    if not has_perm:
        await interaction.followup.send(embed=_error_embed(f"🔒 매매는 **{TRADER_ROLE_NAME}** 역할이 필요합니다."))
        return

    # 매도 금액 검증
    if amount <= 0:
        await interaction.followup.send(embed=_error_embed("금액은 0보다 커야 합니다."))
        return

    ticker = coin.upper()
    if not ticker.startswith("KRW-"):
        ticker = f"KRW-{ticker}"
    try:
        if not upbit_client:
            await interaction.followup.send(embed=_error_embed("Upbit 클라이언트를 불러올 수 없습니다."))
            return

        # ★ TradeOrchestrator: 승인 게이트 ★
        trade_orch = None
        try:
            from jarvis_features.trade_orchestrator import TradeOrchestrator
            trade_orch = _get_trade_orchestrator()
        except ImportError:
            pass

        if action.value == "buy":
            min_order = 5000
            if upbit_config and hasattr(upbit_config, 'MIN_ORDER_AMOUNT'):
                min_order = upbit_config.MIN_ORDER_AMOUNT
            if amount < min_order:
                await interaction.followup.send(embed=_error_embed(f"최소 주문 금액: **{min_order:,.0f}** KRW"))
                return

            # 승인 게이트: 고액 거래는 확인 필요
            if trade_orch:
                req = trade_orch.create_request(
                    "buy", ticker, amount_krw=amount,
                    user_id=str(interaction.user.id),
                )
                if req.needs_approval:
                    embed = discord.Embed(
                        title="거래 승인 요청",
                        description=req.preview_text(),
                        color=discord.Color.orange(),
                    )
                    embed.set_footer(text=f"ID: {req.request_id[:8]} | 5분 후 자동 취소")
                    view = TradeApprovalView(trade_orch, req, upbit_client, ticker, amount, "buy")
                    await interaction.followup.send(embed=embed, view=view)
                    return

            result = await asyncio.to_thread(upbit_client.buy_market_order, ticker, amount)
            embed = _trade_result_embed("매수", ticker, result)
            embed.add_field(name="주문 금액", value=f"{amount:,.0f} KRW", inline=True)
        else:
            # 매도 승인 게이트
            if trade_orch:
                req = trade_orch.create_request(
                    "sell", ticker, sell_percent=100.0,
                    user_id=str(interaction.user.id),
                )
                if req.needs_approval:
                    embed = discord.Embed(
                        title="거래 승인 요청",
                        description=req.preview_text(),
                        color=discord.Color.orange(),
                    )
                    embed.set_footer(text=f"ID: {req.request_id[:8]} | 5분 후 자동 취소")
                    view = TradeApprovalView(trade_orch, req, upbit_client, ticker, amount, "sell")
                    await interaction.followup.send(embed=embed, view=view)
                    return

            result = await asyncio.to_thread(upbit_client.sell_market_order, ticker, amount)
            embed = _trade_result_embed("매도", ticker, result)
            embed.add_field(name="매도 수량", value=f"{amount:.8g}", inline=True)
        # 버튼 UI 추가
        view = TradeView(ticker, amount) if TradeView else None
        await interaction.followup.send(embed=embed, view=view)
    except Exception as e:
        logger.error(f"/trade: {e}")
        await interaction.followup.send(embed=_error_embed(f"매매 실패: {e}"))


@app_commands.command(name="chart", description="코인 가격 차트를 생성합니다")
@app_commands.describe(coin="코인 심볼 (예: BTC)")
async def chart_slash(interaction: discord.Interaction, coin: str = "BTC"):
    """차트 생성 슬래시 명령"""
    await interaction.response.defer()
    ticker = coin.upper()
    if not ticker.startswith("KRW-"):
        ticker = f"KRW-{ticker}"
    coin_name = ticker.replace("KRW-", "")
    try:
        if not upbit_client:
            await interaction.followup.send(embed=_error_embed("Upbit 클라이언트를 불러올 수 없습니다."))
            return
        df = await asyncio.to_thread(upbit_client.get_ohlcv, ticker, interval="day", count=7)
        if df is None or df.empty:
            await interaction.followup.send(embed=_error_embed(f"{coin_name} 차트 데이터 없음"))
            return
        if generate_price_chart:
            prices = df["close"].tolist()
            timestamps = df.index.tolist()
            buf = generate_price_chart(coin_name, prices, timestamps)
            file = discord.File(buf, filename=f"{coin_name}_chart.png")
            embed = discord.Embed(title=f"\U0001F4C8 {coin_name} 가격 차트", color=discord.Color.dark_gold(), timestamp=datetime.now(timezone.utc))
            embed.set_image(url=f"attachment://{coin_name}_chart.png")
            embed.set_footer(text="JARVIS Crypto | Upbit")
            await interaction.followup.send(embed=embed, file=file)
        else:
            # 텍스트 차트 폴백
            closes = df["close"].values
            min_p, max_p = min(closes), max(closes)
            chart_width = 20
            lines = []
            for i, row in df.iterrows():
                close = row["close"]
                bar_len = int((close - min_p) / (max_p - min_p) * chart_width) if max_p > min_p else chart_width // 2
                bar = "\u2588" * bar_len + "\u2591" * (chart_width - bar_len)
                date_str = i.strftime("%m/%d")
                lines.append(f"`{date_str}` {bar} **{close:,.0f}**")
            embed = discord.Embed(title=f"\U0001F4CA {coin_name} 7일 차트", description="\n".join(lines), color=discord.Color.dark_gold(), timestamp=datetime.now(timezone.utc))
            embed.set_footer(text="JARVIS Crypto | 텍스트 차트")
            await interaction.followup.send(embed=embed)
    except Exception as e:
        logger.error(f"/chart: {e}")
        await interaction.followup.send(embed=_error_embed(f"차트 생성 실패: {e}"))


# ═══════════════════════════════════════════════════════════════
#  Background Tasks
# ═══════════════════════════════════════════════════════════════

@tasks.loop(minutes=1)
async def update_status_task(bot):
    try:
        if upbit_client:
            price = await asyncio.to_thread(upbit_client.get_current_price, "KRW-BTC")
            if price:
                await bot.change_presence(activity=discord.Activity(type=discord.ActivityType.watching, name=f"BTC {price:,.0f} KRW"))
            else:
                logger.debug("[StatusTask] Upbit returned None price")
    except Exception as e:
        logger.warning(f"[StatusTask] Upbit price fetch failed: {e}")
        try:
            await bot.change_presence(activity=discord.Activity(type=discord.ActivityType.playing, name="JARVIS Online"))
        except Exception:
            pass

@update_status_task.before_loop
async def before_status(bot):
    await bot.wait_until_ready()


_last_briefing_date: Optional[str] = None  # Track last briefing date to prevent double-fire

@tasks.loop(minutes=1)
async def daily_briefing_task(bot):
    """매일 아침 08:00 (KST) 자동 브리핑"""
    global _last_briefing_date
    now_kst = datetime.now(timezone(timedelta(hours=9)))
    today_str = now_kst.strftime("%Y-%m-%d")
    if now_kst.hour == 8 and now_kst.minute == 0 and _last_briefing_date != today_str:
        _last_briefing_date = today_str
        target_channel_id = BRIEFING_CHANNEL_ID
        channel = None
        if target_channel_id:
            channel = bot.get_channel(int(target_channel_id))
        if not channel:
            for guild in bot.guilds:
                for ch in guild.text_channels:
                    if ch.permissions_for(guild.me).send_messages:
                        channel = ch
                        break
                if channel: break
        if channel and bot._daily_briefing:
            try:
                logger.info(f"Sending Morning Briefing to {channel.name} ({channel.id})")
                report = await bot._daily_briefing.generate_briefing_async()
                if len(report) > DISCORD_MSG_LIMIT:
                    for i in range(0, len(report), 2000):
                        await channel.send(report[i:i+2000])
                else:
                    await channel.send(report)
            except Exception as e:
                logger.error(f"Briefing Auto-Send Error: {e}")

@daily_briefing_task.before_loop
async def before_briefing(bot):
    await bot.wait_until_ready()


_monitor_last_result: str = ""
_monitor_last_alert_time: float = 0.0
_MONITOR_THROTTLE_SEC = 600  # 10분 내 동일 알림 억제

@tasks.loop(minutes=3)
async def monitor_task(bot):
    global _monitor_last_result, _monitor_last_alert_time
    if hasattr(bot, 'monitor_enabled') and bot.monitor_enabled and hasattr(bot, 'monitor_channel') and bot.monitor_channel:
        try:
            logger.info("Auto-Monitor Scanning...")
            result = await bot._analyze_screen("화면을 감시 중입니다. 특이사항이나 중요한 변화가 있는지 확인해서 보고해주세요. 만약 특이사항이 없다면 '특이사항 없음'이라고만 답해주세요.")
            if "특이사항 없음" not in result and len(result) > 10:
                now = time.monotonic()
                if result == _monitor_last_result and (now - _monitor_last_alert_time) < _MONITOR_THROTTLE_SEC:
                    logger.debug("[Monitor] Duplicate alert suppressed")
                    return
                _monitor_last_result = result
                _monitor_last_alert_time = now
                await bot.monitor_channel.send(f"**스마트 감시 알림**\n{result}")
            else:
                _monitor_last_result = ""
        except Exception as e:
            logger.error(f"Monitor Task Error: {e}")


# ── 포트폴리오 자동 추적 (30분 주기) ──
_last_portfolio_value: float = 0.0

@tasks.loop(minutes=30)
async def portfolio_monitor_task(bot):
    """포트폴리오 변동률 5% 초과 시 사령관에게 DM 알림."""
    global _last_portfolio_value
    if not crypto_mcp_server:
        return
    try:
        summary = await crypto_mcp_server.portfolio_summary()
        # 총 자산 추출 (첫 번째 숫자)
        import re as _re_local
        m = _re_local.search(r'[\d,]+(?:\.\d+)?', summary.replace(',', ''))
        if not m:
            return
        current_value = float(m.group().replace(',', ''))
        if _last_portfolio_value > 0:
            change_pct = abs(current_value - _last_portfolio_value) / _last_portfolio_value * 100
            if change_pct >= 5.0:
                direction = "상승" if current_value > _last_portfolio_value else "하락"
                if BOT_OWNER_ID:
                    try:
                        owner = await bot.fetch_user(int(BOT_OWNER_ID))
                        await owner.send(
                            f"**포트폴리오 변동 알림**\n"
                            f"{direction} {change_pct:.1f}% | "
                            f"{_last_portfolio_value:,.0f} → {current_value:,.0f} KRW"
                        )
                    except Exception:
                        pass
        _last_portfolio_value = current_value
    except Exception as e:
        logger.warning(f"[PortfolioMonitor] {e}")

@portfolio_monitor_task.before_loop
async def before_portfolio_monitor(bot):
    await bot.wait_until_ready()


# ── 가격 알림 체크 (1분 주기) ──
# 구조: {user_id: {"KRW-BTC": 60000000, "KRW-ETH": 5000000}}
_price_alert_store: Dict[str, Dict[str, float]] = {}
_price_alert_lock = asyncio.Lock()

@tasks.loop(minutes=1)
async def price_alert_check_task(bot):
    """등록된 가격 알림 조건 매칭 시 DM 발송."""
    if not _price_alert_store or not upbit_client:
        return
    try:
        async with _price_alert_lock:
            all_tickers = set()
            for alerts in _price_alert_store.values():
                all_tickers.update(alerts.keys())
            if not all_tickers:
                return
        prices = await asyncio.to_thread(upbit_client.get_prices, list(all_tickers))
        if not prices:
            return
        async with _price_alert_lock:
            triggered = []  # (user_id, ticker, target, current)
            for user_id, alerts in list(_price_alert_store.items()):
                for ticker, target_price in list(alerts.items()):
                    current = prices.get(ticker)
                    if current and current >= target_price:
                        triggered.append((user_id, ticker, target_price, current))
                        del alerts[ticker]
                if not alerts:
                    del _price_alert_store[user_id]
        for user_id, ticker, target, current in triggered:
            try:
                user = await bot.fetch_user(int(user_id))
                coin = ticker.replace("KRW-", "")
                await user.send(
                    f"**가격 알림 도달**\n"
                    f"{coin}: {current:,.0f} KRW (목표: {target:,.0f})"
                )
            except Exception:
                pass
    except Exception as e:
        logger.warning(f"[PriceAlert] {e}")

@price_alert_check_task.before_loop
async def before_price_alert(bot):
    await bot.wait_until_ready()


# ═══════════════════════════════════════════════════════════════
#  Main
# ═══════════════════════════════════════════════════════════════

def main():
    if not DISCORD_BOT_TOKEN:
        logger.error("No Token")
        return
    bot = JarvisBot()

    logger.info("JARVIS Discord Bot starting...")
    logger.info(f"  Claude API: {'OK' if CLAUDE_API_KEY else 'N/A'}")
    logger.info(f"  Gemini API: {'OK' if GEMINI_API_KEY else 'N/A'}")
    logger.info(f"  MCP Modules: {MCP_AVAILABLE} (sc2={sc2_mcp_server is not None}, sys={system_mcp_server is not None}, crypto={crypto_mcp_server is not None})")
    logger.info(f"  Advanced Features: {ADVANCED_AVAILABLE}")
    logger.info(f"  Web Tools: {'OK' if web_tools else 'N/A'}")
    logger.info(f"  Memory: {'OK' if MemoryManager else 'N/A'}")

    # ★ Phase 5: Graceful Shutdown ★
    def _graceful_shutdown(sig, frame):
        logger.info(f"[SHUTDOWN] Signal {sig} received. Initiating graceful shutdown...")
        # Save audit log summary
        if get_tool_registry:
            logger.info(f"[SHUTDOWN] Tool stats:\n{get_tool_registry().get_summary()}")
        # Close bot
        asyncio.ensure_future(bot.close())

    signal.signal(signal.SIGINT, _graceful_shutdown)
    signal.signal(signal.SIGTERM, _graceful_shutdown)

    bot.run(DISCORD_BOT_TOKEN)


if __name__ == "__main__":
    main()

"""
UTM Airspace Control Presentation PPTX Generator
=================================================
utm_slides.html 10장 → python-pptx 편집 가능 슬라이드 생성
다크 블루 배경 + Cyan/Green 악센트 테마
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ===== Color Constants =====
BG_DARK   = RGBColor(0x0B, 0x11, 0x20)
BG_CARD   = RGBColor(0x0F, 0x17, 0x2A)
CYAN      = RGBColor(0x00, 0xC8, 0xFF)
GREEN     = RGBColor(0x00, 0xFF, 0x88)
YELLOW    = RGBColor(0xFF, 0xC8, 0x32)
RED       = RGBColor(0xFF, 0x64, 0x64)
PURPLE    = RGBColor(0xA7, 0x8B, 0xFA)
ORANGE    = RGBColor(0xF9, 0x73, 0x16)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xE0, 0xE0, 0xE0)
DIM       = RGBColor(0x94, 0xA3, 0xB8)
DIM2      = RGBColor(0xB0, 0xBE, 0xC5)
BODY_DIM  = RGBColor(0x64, 0x74, 0x8B)
CYAN_LIGHT = RGBColor(0x7D, 0xD3, 0xFC)
PINK      = RGBColor(0xE8, 0x79, 0xF9)
SKY       = RGBColor(0x38, 0xBD, 0xF8)

FONT_BODY = "Malgun Gothic"
FONT_CODE = "Consolas"


# ===== Helper Functions =====

def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name=FONT_BODY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, items, spacing=4):
    """Add a text box with multiple styled paragraphs.
    items: list of (text, font_size, color, bold, font_name)
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        text, fs, clr, bld = item[:4]
        fn = item[4] if len(item) > 4 else FONT_BODY
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(fs)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = fn
        p.space_after = Pt(spacing)
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color=BG_CARD,
                     border_color=None, border_width=2):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()
    return shape


def add_tag(slide, left, top, text, bg_color=CYAN, text_color=BG_DARK, w=None):
    width = w or Inches(1.6)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, Inches(0.32)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = FONT_CODE
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_footer(slide, slide_num, total=10):
    add_text_box(slide, Inches(0.4), Inches(6.9), Inches(3), Inches(0.3),
                 "Swarm Drone Airspace Control", font_size=9, color=BODY_DIM,
                 font_name=FONT_CODE)
    add_text_box(slide, Inches(11), Inches(6.9), Inches(2), Inches(0.3),
                 f"{slide_num} / {total}", font_size=9, color=BODY_DIM,
                 font_name=FONT_CODE, alignment=PP_ALIGN.RIGHT)


def add_top_accent(slide):
    """Gradient-like top bar."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Pt(4)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = CYAN
    bar.line.fill.background()


def add_card_with_header(slide, x, y, w, h, header, header_color, items,
                         border_side="top"):
    """Card box with colored header and bullet items."""
    card = add_rounded_rect(slide, x, y, w, h, BG_CARD, header_color)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.12), w - Inches(0.3),
                 Inches(0.3), header, font_size=14, color=header_color, bold=True)
    txt = "\n".join(f"  {item}" for item in items)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.45), w - Inches(0.3),
                 h - Inches(0.55), txt, font_size=12, color=DIM)
    return card


# ===== Slide Builders =====

def build_slide_1(prs):
    """Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_accent(slide)

    add_text_box(slide, Inches(2), Inches(1.6), Inches(9), Inches(0.8),
                 "군집드론 활용", font_size=48, color=CYAN, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2), Inches(2.3), Inches(9), Inches(0.7),
                 "공역 통제 시스템", font_size=38, color=GREEN, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2), Inches(3.1), Inches(9), Inches(0.5),
                 "Swarm Drone Airspace Control System", font_size=20,
                 color=CYAN_LIGHT, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2), Inches(3.7), Inches(9), Inches(0.5),
                 "Boids 3D 분산 제어  |  TTC 충돌 예측  |  가상 비행 회랑",
                 font_size=16, color=DIM, alignment=PP_ALIGN.CENTER)

    # Tags
    tags = [("Boids 3D", CYAN), ("UTM", GREEN), ("TTC Collision", YELLOW),
            ("Fail-Safe", RED), ("MAVLink", PURPLE)]
    start_x = 2.8
    for i, (label, clr) in enumerate(tags):
        add_tag(slide, Inches(start_x + i * 1.7), Inches(4.5), label, clr, BG_DARK)

    add_text_box(slide, Inches(2), Inches(5.3), Inches(9), Inches(0.4),
                 "Capstone Design 2026", font_size=15, color=BODY_DIM,
                 alignment=PP_ALIGN.CENTER)

    add_footer(slide, 1)


def build_slide_2(prs):
    """Project overview — problem/solution + stats."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_accent(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "프로젝트 개요", font_size=36, color=CYAN, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.15), Inches(10), Inches(0.4),
                 "기존 게임 AI 지휘관 시스템을 실제 드론 군집 공역 통제로 전환합니다.",
                 font_size=16, color=DIM)

    # Problem card
    add_rounded_rect(slide, Inches(0.6), Inches(1.7), Inches(5.8), Inches(2.3),
                     BG_CARD, RED)
    add_text_box(slide, Inches(0.8), Inches(1.8), Inches(5), Inches(0.35),
                 "문제 정의", font_size=16, color=RED, bold=True)
    problems = [
        "- 드론 수 증가 -> 충돌/공역 침범 위험 급증",
        "- 중앙 서버 의존 -> 단일 장애점(SPOF)",
        "- 수동 경로 설정 -> 확장성 한계",
        "- 긴급 상황 대응 -> 실시간 판단 부재",
    ]
    add_text_box(slide, Inches(0.8), Inches(2.2), Inches(5.4), Inches(1.6),
                 "\n".join(problems), font_size=13, color=DIM)

    # Solution card
    add_rounded_rect(slide, Inches(6.7), Inches(1.7), Inches(5.8), Inches(2.3),
                     BG_CARD, GREEN)
    add_text_box(slide, Inches(6.9), Inches(1.8), Inches(5), Inches(0.35),
                 "해결 방안", font_size=16, color=GREEN, bold=True)
    solutions = [
        "- Boids 3D: 분산 자율 군집 제어",
        "- TTC 예측: 시간 기반 선제적 충돌 회피",
        "- 비행 회랑: 고도층/방향 분리 교통 규칙",
        "- Authority FSM: 우선순위 기반 긴급 대응",
    ]
    add_text_box(slide, Inches(6.9), Inches(2.2), Inches(5.4), Inches(1.6),
                 "\n".join(solutions), font_size=13, color=DIM)

    # 4 stat boxes
    stats = [
        ("8", "Boids 힘 벡터", CYAN),
        ("3", "고도층 (Low/Med/High)", GREEN),
        ("5초", "TTC 예측 범위", YELLOW),
        ("5", "Authority 모드", PURPLE),
    ]
    for i, (num, label, clr) in enumerate(stats):
        x = Inches(0.6 + i * 3.15)
        add_rounded_rect(slide, x, Inches(4.2), Inches(2.9), Inches(1.0), BG_CARD)
        add_text_box(slide, x, Inches(4.25), Inches(2.9), Inches(0.55),
                     num, font_size=32, color=clr, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name=FONT_CODE)
        add_text_box(slide, x, Inches(4.8), Inches(2.9), Inches(0.3),
                     label, font_size=11, color=DIM, alignment=PP_ALIGN.CENTER)

    # Flow: SC2 → 알고리즘 추출 → 3D 확장 → 드론 공역 통제
    flow_items = [
        ("SC2 지휘관 봇", RED), ("알고리즘 추출", PURPLE),
        ("3D 확장", CYAN), ("드론 공역 통제", GREEN),
    ]
    y_flow = Inches(5.5)
    for i, (label, clr) in enumerate(flow_items):
        x = Inches(1.0 + i * 3.0)
        add_tag(slide, x, y_flow, label, clr, BG_DARK, w=Inches(2.4))
        if i < 3:
            add_text_box(slide, x + Inches(2.5), y_flow, Inches(0.4), Inches(0.3),
                         "->", font_size=14, color=BODY_DIM, font_name=FONT_CODE)

    add_footer(slide, 2)


def build_slide_3(prs):
    """System architecture — 3 columns + mapping table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_accent(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
                 "시스템 아키텍처", font_size=36, color=CYAN, bold=True)

    # Flow bar
    flow = [("Sensor Input", YELLOW), ("Flight Data Hub", CYAN),
            ("Decision Engine", PURPLE), ("MAVLink Cmd", GREEN)]
    for i, (label, clr) in enumerate(flow):
        x = Inches(0.8 + i * 3.1)
        add_tag(slide, x, Inches(1.1), label, clr, BG_DARK, w=Inches(2.5))
        if i < 3:
            add_text_box(slide, x + Inches(2.55), Inches(1.1), Inches(0.4), Inches(0.3),
                         "->", font_size=12, color=BODY_DIM, font_name=FONT_CODE)

    # 3 module cards
    modules = [
        ("Flight Data Hub", CYAN, "(= SC2 Blackboard)",
         ["DroneState: 위치/속도/헤딩", "공역 상태: 고도층 점유율",
          "충돌 경보: TTC 알림 큐", "회랑 할당: 비행 경로 상태",
          "Authority Mode: 현재 우선순위"]),
        ("Decision Engine", PURPLE, "(= SC2 CombatManager)",
         ["Boids3DController: 군집 제어", "CollisionPredictor: TTC 예측",
          "CorridorManager: 회랑 관리", "AuthorityFSM: 긴급 모드",
          "PathPlanner: 경로 최적화"]),
        ("Communication", GREEN, "(= SC2 API Layer)",
         ["MAVLink: 드론 명령 프로토콜", "GPS/IMU: 위치/자세 데이터",
          "LiDAR: 장애물 감지", "Mesh Network: 드론 간 통신",
          "GCS: 지상 관제 UI"]),
    ]
    for i, (title, clr, origin, items) in enumerate(modules):
        x = Inches(0.5 + i * 4.2)
        y = Inches(1.6)
        add_rounded_rect(slide, x, y, Inches(3.9), Inches(2.6), BG_CARD, clr)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.1), Inches(3.5),
                     Inches(0.3), title, font_size=15, color=clr, bold=True)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.4), Inches(3.5),
                     Inches(0.25), origin, font_size=11, color=DIM)
        txt = "\n".join(f"  {it}" for it in items)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.7), Inches(3.5),
                     Inches(1.8), txt, font_size=12, color=DIM)

    # Mapping table
    add_text_box(slide, Inches(0.8), Inches(4.4), Inches(8), Inches(0.35),
                 "SC2 -> 드론 ATC 모듈 매핑", font_size=14, color=CYAN_LIGHT, bold=True)

    table_data = [
        ("SC2 봇 모듈", "드론 ATC 모듈", "역할"),
        ("Blackboard", "Flight Data Hub", "중앙 상태 저장소"),
        ("BoidsSwarmControl", "Boids3DController", "군집 이동 제어 (8-Force)"),
        ("PotentialFields", "CollisionPredictor", "충돌 예측 & 회피"),
        ("CreepHighway", "CorridorManager", "비행 회랑 관리"),
        ("AuthorityMode FSM", "Authority FSM", "긴급 상황 우선순위"),
    ]
    rows, cols = len(table_data), 3
    tbl_shape = slide.shapes.add_table(rows, cols,
                                        Inches(0.6), Inches(4.8),
                                        Inches(11.5), Inches(2.0))
    tbl = tbl_shape.table
    col_widths = [Inches(3.5), Inches(3.5), Inches(4.5)]
    for ci, w in enumerate(col_widths):
        tbl.columns[ci].width = w

    for r, row_data in enumerate(table_data):
        tbl.rows[r].height = Inches(0.32)
        for c, cell_text in enumerate(row_data):
            cell = tbl.cell(r, c)
            cell.text = cell_text
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.1)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            p = cell.text_frame.paragraphs[0]
            p.font.name = FONT_BODY
            if r == 0:
                p.font.size = Pt(12)
                p.font.bold = True
                p.font.color.rgb = CYAN_LIGHT
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x12, 0x1A, 0x2E)
            else:
                p.font.size = Pt(11)
                p.font.color.rgb = DIM2 if c != 1 else CYAN
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_CARD

    add_footer(slide, 3)


def build_slide_4(prs):
    """Boids 3D — 8-Force vector."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_accent(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
                 "Boids 3D 군집 제어", font_size=36, color=CYAN, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.35),
                 "SC2의 6-Force를 3D로 확장하고 고도 관련 2개 힘을 추가하여 총 8개 힘 벡터를 사용합니다.",
                 font_size=14, color=DIM)

    # Formula box
    add_rounded_rect(slide, Inches(2.5), Inches(1.45), Inches(8), Inches(0.45),
                     RGBColor(0x0A, 0x1A, 0x12), GREEN, border_width=1)
    add_text_box(slide, Inches(2.7), Inches(1.48), Inches(7.5), Inches(0.4),
                 "F_total = Sum(i=1..8) w_i * F_i(drone, neighbors, environment)",
                 font_size=14, color=GREEN, font_name=FONT_CODE,
                 alignment=PP_ALIGN.CENTER)

    # Left column: SC2 forces (5)
    add_rounded_rect(slide, Inches(0.5), Inches(2.1), Inches(6.0), Inches(4.6),
                     BG_CARD, CYAN)
    add_text_box(slide, Inches(0.7), Inches(2.2), Inches(5.5), Inches(0.3),
                 "기존 SC2 힘 (3D 확장)", font_size=14, color=CYAN, bold=True)

    forces_sc2 = [
        ("1. Separation", RED, "w=2.0 / r=5m", "역제곱 반발력, 드론 간 안전 거리 유지"),
        ("2. Alignment", CYAN, "w=1.0 / r=30m", "이웃 속도 벡터 정렬, 편대 방향 통일"),
        ("3. Cohesion", GREEN, "w=1.0 / r=30m", "그룹 중심 응집, 편대 유지"),
        ("4. Target Seeking", YELLOW, "w=1.5", "목표 웨이포인트로 이동, 거리 비례 감쇠"),
        ("5. Obstacle Avoid", PURPLE, "w=2.5 / r=20m", "장애물(건물, 산) 회피, 반발력 기반"),
    ]
    for i, (name, clr, params, desc) in enumerate(forces_sc2):
        y = Inches(2.55 + i * 0.8)
        # Dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                      Inches(0.75), y + Inches(0.05),
                                      Inches(0.15), Inches(0.15))
        dot.fill.solid()
        dot.fill.fore_color.rgb = clr
        dot.line.fill.background()
        add_text_box(slide, Inches(1.0), y - Inches(0.02), Inches(2.5), Inches(0.25),
                     name, font_size=13, color=clr, bold=True)
        add_text_box(slide, Inches(3.5), y - Inches(0.02), Inches(1.5), Inches(0.25),
                     params, font_size=11, color=BODY_DIM, font_name=FONT_CODE)
        add_text_box(slide, Inches(1.0), y + Inches(0.22), Inches(5.0), Inches(0.25),
                     desc, font_size=11, color=DIM)

    # Right column: UTM forces (3) + kinematics
    add_rounded_rect(slide, Inches(6.8), Inches(2.1), Inches(5.8), Inches(4.6),
                     BG_CARD, GREEN)
    add_text_box(slide, Inches(7.0), Inches(2.2), Inches(5.3), Inches(0.3),
                 "신규 UTM 힘", font_size=14, color=GREEN, bold=True)

    forces_utm = [
        ("6. Altitude Hold", SKY, "w=1.0", "지정 고도 유지 (PID P항), 수직 복원력"),
        ("7. Terrain Clearance", ORANGE, "w=3.0", "최소 안전 고도(5m) 미만 시 강제 상승"),
        ("8. Corridor Follow", PINK, "동적", "회랑 중심선 추적, 이탈 시 복원력"),
    ]
    for i, (name, clr, params, desc) in enumerate(forces_utm):
        y = Inches(2.55 + i * 0.8)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                      Inches(7.05), y + Inches(0.05),
                                      Inches(0.15), Inches(0.15))
        dot.fill.solid()
        dot.fill.fore_color.rgb = clr
        dot.line.fill.background()
        add_text_box(slide, Inches(7.3), y - Inches(0.02), Inches(2.5), Inches(0.25),
                     name, font_size=13, color=clr, bold=True)
        add_text_box(slide, Inches(9.8), y - Inches(0.02), Inches(1.2), Inches(0.25),
                     params, font_size=11, color=BODY_DIM, font_name=FONT_CODE)
        add_text_box(slide, Inches(7.3), y + Inches(0.22), Inches(5.0), Inches(0.25),
                     desc, font_size=11, color=DIM)

    # Kinematics section
    add_text_box(slide, Inches(7.0), Inches(5.0), Inches(5.3), Inches(0.3),
                 "드론 운동학 제한", font_size=13, color=YELLOW, bold=True)
    kin_lines = [
        "max_speed: 15 m/s",
        "max_acceleration: 5 m/s^2",
        "max_force: 5 m/s^2 (클램핑)",
        "altitude: 5~120m (법적 제한)",
    ]
    add_text_box(slide, Inches(7.0), Inches(5.3), Inches(5.3), Inches(1.2),
                 "\n".join(f"  {l}" for l in kin_lines), font_size=11, color=DIM)

    add_footer(slide, 4)


def build_slide_5(prs):
    """TTC collision prediction."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_accent(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
                 "TTC 충돌 예측 시스템", font_size=36, color=CYAN, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.35),
                 "기존 거리 기반 회피를 시간 기반 궤적 예측으로 업그레이드하여 선제적 충돌 회피를 구현합니다.",
                 font_size=14, color=DIM)

    # Left: Old vs New
    # Old method
    add_rounded_rect(slide, Inches(0.5), Inches(1.6), Inches(3.8), Inches(2.5),
                     BG_CARD, RED)
    add_text_box(slide, Inches(0.7), Inches(1.7), Inches(3.4), Inches(0.3),
                 "기존 방식 (SC2 Boids)", font_size=13, color=RED, bold=True)
    old_items = [
        "거리 기반 반발력",
        "if distance < threshold -> 반발력 적용",
        "",
        "X  고속 접근 시 반응 지연",
        "X  병렬 이동 드론 불필요한 회피",
        "X  교차 궤적 사전 감지 불가",
    ]
    add_text_box(slide, Inches(0.7), Inches(2.1), Inches(3.4), Inches(1.8),
                 "\n".join(old_items), font_size=11, color=DIM)

    # New method
    add_rounded_rect(slide, Inches(0.5), Inches(4.3), Inches(3.8), Inches(2.3),
                     BG_CARD, GREEN)
    add_text_box(slide, Inches(0.7), Inches(4.4), Inches(3.4), Inches(0.3),
                 "새 방식 (TTC 예측)", font_size=13, color=GREEN, bold=True)
    new_items = [
        "시간 기반 궤적 교차 예측",
        "미래 위치 계산 -> 최소 이격 거리 평가",
        "",
        "V  5초 전 선제적 회피 기동",
        "V  상대 속도 기반 정밀 판단",
        "V  3단계 심각도 분류",
    ]
    add_text_box(slide, Inches(0.7), Inches(4.8), Inches(3.4), Inches(1.6),
                 "\n".join(new_items), font_size=11, color=DIM)

    # Right: TTC formula
    add_rounded_rect(slide, Inches(4.6), Inches(1.6), Inches(4.0), Inches(3.0),
                     BG_CARD, PURPLE)
    add_text_box(slide, Inches(4.8), Inches(1.7), Inches(3.6), Inches(0.3),
                 "TTC 수학 모델", font_size=13, color=PURPLE, bold=True)
    formula = (
        "rel_pos = B.pos - A.pos\n"
        "rel_vel = B.vel - A.vel\n"
        "\n"
        "TTC = -dot(rel_pos, rel_vel)\n"
        "      / dot(rel_vel, rel_vel)\n"
        "\n"
        "min_dist = |A.pos + A.vel*TTC\n"
        "           - (B.pos + B.vel*TTC)|"
    )
    add_rounded_rect(slide, Inches(4.8), Inches(2.1), Inches(3.6), Inches(2.3),
                     RGBColor(0x0A, 0x1A, 0x12), GREEN, border_width=1)
    add_text_box(slide, Inches(5.0), Inches(2.2), Inches(3.2), Inches(2.1),
                 formula, font_size=12, color=GREEN, font_name=FONT_CODE)

    # Alert levels table
    add_rounded_rect(slide, Inches(4.6), Inches(4.8), Inches(4.0), Inches(1.8),
                     BG_CARD, YELLOW)
    add_text_box(slide, Inches(4.8), Inches(4.9), Inches(3.6), Inches(0.3),
                 "3단계 경보 체계", font_size=13, color=YELLOW, bold=True)

    alert_data = [
        ("레벨", "TTC", "대응"),
        ("WARNING", "<= 5.0초", "경로 미세 조정"),
        ("CRITICAL", "<= 3.0초", "회피 기동 시작"),
        ("IMMINENT", "<= 1.0초", "긴급 회피 (수직)"),
    ]
    tbl_shape = slide.shapes.add_table(4, 3,
                                        Inches(4.8), Inches(5.3),
                                        Inches(3.6), Inches(1.2))
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(1.2)
    tbl.columns[1].width = Inches(1.0)
    tbl.columns[2].width = Inches(1.4)
    alert_colors = [None, YELLOW, ORANGE, RED]
    for r, row_data in enumerate(alert_data):
        tbl.rows[r].height = Inches(0.28)
        for c, txt in enumerate(row_data):
            cell = tbl.cell(r, c)
            cell.text = txt
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.08)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            p = cell.text_frame.paragraphs[0]
            p.font.name = FONT_BODY
            p.font.size = Pt(11)
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = CYAN_LIGHT
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x12, 0x1A, 0x2E)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_CARD
                p.font.color.rgb = alert_colors[r] if c == 0 else DIM2
                if c == 0:
                    p.font.bold = True

    # Right side info card
    add_rounded_rect(slide, Inches(9.0), Inches(1.6), Inches(3.8), Inches(5.0),
                     BG_CARD, CYAN)
    add_text_box(slide, Inches(9.2), Inches(1.7), Inches(3.4), Inches(0.3),
                 "핵심 특징", font_size=13, color=CYAN, bold=True)
    features = [
        "- 등속 직선 운동 가정",
        "- 상대 위치/속도 벡터 해석",
        "- O(N^2) 전 쌍 검사",
        "- VoxelGrid 사전 필터링 가능",
        "- 시간 기반 회피 강도 조절:",
        "  strength = (5s - TTC) / 5s",
        "- 정확히 같은 위치 ->",
        "  수직 회피 (0, 0, 1)",
        "- 충돌 시점 중간점 계산",
        "- 긴급도 순 정렬 반환",
    ]
    add_text_box(slide, Inches(9.2), Inches(2.1), Inches(3.4), Inches(4.3),
                 "\n".join(features), font_size=11, color=DIM)

    add_footer(slide, 5)


def build_slide_6(prs):
    """Flight corridor system."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_accent(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
                 "가상 비행 회랑 시스템", font_size=36, color=CYAN, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.35),
                 "SC2 점막 고속도로(CreepHighway)의 A* 경로를 3D 비행 회랑으로 진화시켰습니다.",
                 font_size=14, color=DIM)

    # Left: Altitude layers
    add_rounded_rect(slide, Inches(0.5), Inches(1.6), Inches(6.0), Inches(3.2),
                     BG_CARD, CYAN)
    add_text_box(slide, Inches(0.7), Inches(1.7), Inches(5.5), Inches(0.3),
                 "고도층 분리 (3 Layer)", font_size=14, color=CYAN, bold=True)

    layers = [
        ("LOW", "5~30m", "촬영, 농업, 측량", GREEN),
        ("MEDIUM", "30~60m", "배송, 순찰, 감시", CYAN),
        ("HIGH", "60~120m", "장거리 이동, 중계", PURPLE),
    ]
    for i, (name, alt, desc, clr) in enumerate(layers):
        y = Inches(2.15 + i * 0.55)
        add_rounded_rect(slide, Inches(0.8), y, Inches(5.4), Inches(0.45),
                         RGBColor(0x0A, 0x14, 0x28), clr, border_width=1)
        add_text_box(slide, Inches(0.9), y + Inches(0.05), Inches(1.2), Inches(0.35),
                     name, font_size=12, color=clr, bold=True, font_name=FONT_CODE)
        add_text_box(slide, Inches(2.2), y + Inches(0.05), Inches(3.8), Inches(0.35),
                     f"{alt} -- {desc}", font_size=11, color=DIM)

    # Direction separation
    add_text_box(slide, Inches(0.7), Inches(3.9), Inches(5.5), Inches(0.3),
                 "방향별 고도 분리", font_size=13, color=CYAN_LIGHT, bold=True)
    dir_info = [
        "  Eastbound (0~180 deg) -> 기본 고도",
        "  Westbound (180~360 deg) -> +15m 오프셋",
        "  항공 교통 규칙 (ICAO) 모사",
    ]
    add_text_box(slide, Inches(0.7), Inches(4.2), Inches(5.5), Inches(0.6),
                 "\n".join(dir_info), font_size=11, color=DIM)

    # Right: Corridor structure + altitude profile
    add_rounded_rect(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(2.3),
                     BG_CARD, YELLOW)
    add_text_box(slide, Inches(7.0), Inches(1.7), Inches(5.3), Inches(0.3),
                 "회랑 구조", font_size=14, color=YELLOW, bold=True)
    corridor_items = [
        "FlightCorridor 구성요소:",
        "  waypoints[]  -- 3D 웨이포인트 시퀀스",
        "  width        -- 회랑 폭 (기본 20m)",
        "  altitude_layer -- 할당 고도층",
        "  bidirectional -- 양방향 여부",
    ]
    add_text_box(slide, Inches(7.0), Inches(2.1), Inches(5.3), Inches(1.6),
                 "\n".join(corridor_items), font_size=11, color=DIM)

    # Altitude profile
    add_text_box(slide, Inches(7.0), Inches(4.1), Inches(5.3), Inches(0.3),
                 "고도 프로파일", font_size=13, color=CYAN_LIGHT, bold=True)
    profiles = [
        ("0~15%", "상승 구간 (이륙 -> 순항 고도)", GREEN),
        ("15~85%", "순항 구간 (일정 고도 유지)", CYAN),
        ("85~100%", "하강 구간 (순항 -> 착륙)", YELLOW),
    ]
    for i, (pct, desc, clr) in enumerate(profiles):
        y = Inches(4.45 + i * 0.4)
        add_tag(slide, Inches(7.0), y, pct, clr, BG_DARK, w=Inches(1.2))
        add_text_box(slide, Inches(8.4), y, Inches(4.0), Inches(0.35),
                     desc, font_size=11, color=DIM)

    # Corridor conflict card
    add_rounded_rect(slide, Inches(6.8), Inches(5.7), Inches(5.8), Inches(0.9),
                     BG_CARD, RED)
    add_text_box(slide, Inches(7.0), Inches(5.8), Inches(5.3), Inches(0.25),
                 "회랑 충돌 검출", font_size=13, color=RED, bold=True)
    add_text_box(slide, Inches(7.0), Inches(6.1), Inches(5.3), Inches(0.4),
                 "두 회랑 웨이포인트 교차점 자동 검출 -> 시간 분리 규칙 적용 -> 동적 우회 경로 생성",
                 font_size=11, color=DIM)

    add_footer(slide, 6)


def build_slide_7(prs):
    """Authority Mode FSM."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_accent(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
                 "Authority Mode -- 긴급 대응 FSM", font_size=36, color=CYAN, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.35),
                 "SC2 봇의 Authority Mode를 드론 공역 통제의 우선순위 기반 상태 머신으로 전환합니다.",
                 font_size=14, color=DIM)

    modes = [
        ("EMERGENCY (P0)", RED,
         "트리거: TTC <= 1초, 배터리 부족, GPS 상실",
         "대응: 즉시 정지, 수직 회피, 자동 귀환",
         "SC2: 적 러시 감지 -> 전 유닛 방어"),
        ("DECONFLICT (P1)", ORANGE,
         "트리거: TTC <= 5초, 회랑 교차",
         "대응: 회피 기동, 속도/고도 조정",
         "SC2: 적 근접 -> 전투 태세"),
        ("MISSION (P2)", PURPLE,
         "트리거: 임무 수행 중",
         "대응: 촬영/배송/순찰 임무 수행",
         "SC2: 전략 모드 -> 테크 업그레이드"),
        ("CRUISE (P3)", CYAN,
         "트리거: 이동 중, 위협 없음",
         "대응: Boids 편대 비행, 에너지 최적화",
         "SC2: 경제 모드 -> 자원 수집"),
        ("IDLE (P4)", GREEN,
         "트리거: 대기 중, 착륙 상태",
         "대응: 호버링 / 배터리 충전 대기",
         "SC2: BALANCED -> 균형 유지"),
    ]

    # Left column: 3 modes, Right column: 2 modes + transition rules
    for i, (name, clr, trigger, response, sc2) in enumerate(modes[:3]):
        y = Inches(1.5 + i * 1.65)
        add_rounded_rect(slide, Inches(0.5), y, Inches(6.0), Inches(1.45),
                         BG_CARD, clr)
        add_text_box(slide, Inches(0.7), y + Inches(0.08), Inches(5.5), Inches(0.3),
                     name, font_size=14, color=clr, bold=True)
        info = f"{trigger}\n{response}\n{sc2}"
        add_text_box(slide, Inches(0.7), y + Inches(0.4), Inches(5.5), Inches(0.9),
                     info, font_size=11, color=DIM)

    for i, (name, clr, trigger, response, sc2) in enumerate(modes[3:]):
        y = Inches(1.5 + i * 1.65)
        add_rounded_rect(slide, Inches(6.8), y, Inches(5.8), Inches(1.45),
                         BG_CARD, clr)
        add_text_box(slide, Inches(7.0), y + Inches(0.08), Inches(5.3), Inches(0.3),
                     name, font_size=14, color=clr, bold=True)
        info = f"{trigger}\n{response}\n{sc2}"
        add_text_box(slide, Inches(7.0), y + Inches(0.4), Inches(5.3), Inches(0.9),
                     info, font_size=11, color=DIM)

    # Transition rules card
    add_rounded_rect(slide, Inches(6.8), Inches(4.8), Inches(5.8), Inches(1.7),
                     BG_CARD, YELLOW)
    add_text_box(slide, Inches(7.0), Inches(4.9), Inches(5.3), Inches(0.3),
                 "상태 전이 규칙", font_size=13, color=YELLOW, bold=True)
    rules = [
        "  높은 우선순위가 항상 선점",
        "  EMERGENCY는 모든 상태에서 즉시 전환",
        "  위협 해제 시 이전 상태로 복귀",
        "  배터리 20% 미만 -> 강제 EMERGENCY",
    ]
    add_text_box(slide, Inches(7.0), Inches(5.25), Inches(5.3), Inches(1.1),
                 "\n".join(rules), font_size=11, color=DIM)

    add_footer(slide, 7)


def build_slide_8(prs):
    """Data model — DroneState, Point3D, CollisionAlert, FlightCorridor."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_accent(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
                 "데이터 모델 & 드론 상태", font_size=36, color=CYAN, bold=True)

    # DroneState
    add_rounded_rect(slide, Inches(0.4), Inches(1.2), Inches(5.6), Inches(2.8),
                     BG_CARD, CYAN)
    add_text_box(slide, Inches(0.6), Inches(1.3), Inches(5.2), Inches(0.3),
                 "DroneState (= SC2 Unit)", font_size=13, color=CYAN, bold=True)
    ds_code = (
        "@dataclass\n"
        "class DroneState:\n"
        "  id: int\n"
        "  position: Point3D   # x,y,z (m)\n"
        "  velocity: np.ndarray # [vx,vy,vz]\n"
        "  heading: float       # rad\n"
        "  drone_type: str      # quadrotor\n"
        "  max_speed: float = 15.0\n"
        "  max_accel: float = 5.0\n"
        "  min_alt: float   = 5.0\n"
        "  max_alt: float   = 120.0"
    )
    add_rounded_rect(slide, Inches(0.6), Inches(1.65), Inches(5.2), Inches(2.2),
                     RGBColor(0x0F, 0x17, 0x2A))
    add_text_box(slide, Inches(0.7), Inches(1.7), Inches(5.0), Inches(2.1),
                 ds_code, font_size=11, color=LIGHT, font_name=FONT_CODE)

    # Point3D
    add_rounded_rect(slide, Inches(0.4), Inches(4.2), Inches(5.6), Inches(2.4),
                     BG_CARD, GREEN)
    add_text_box(slide, Inches(0.6), Inches(4.3), Inches(5.2), Inches(0.3),
                 "Point3D (= SC2 Point2)", font_size=13, color=GREEN, bold=True)
    p3d_code = (
        "class Point3D:\n"
        "  x, y, z: float\n"
        "  # distance_to()          -> 3D Euclidean\n"
        "  # horizontal_distance_to() -> 2D\n"
        "  # to_array()             -> numpy vector\n"
        "  # towards()              -> directional move\n"
        "  # offset()               -> coordinate shift"
    )
    add_rounded_rect(slide, Inches(0.6), Inches(4.65), Inches(5.2), Inches(1.8),
                     RGBColor(0x0F, 0x17, 0x2A))
    add_text_box(slide, Inches(0.7), Inches(4.7), Inches(5.0), Inches(1.7),
                 p3d_code, font_size=11, color=LIGHT, font_name=FONT_CODE)

    # CollisionAlert
    add_rounded_rect(slide, Inches(6.3), Inches(1.2), Inches(6.3), Inches(2.4),
                     BG_CARD, RED)
    add_text_box(slide, Inches(6.5), Inches(1.3), Inches(5.9), Inches(0.3),
                 "CollisionAlert", font_size=13, color=RED, bold=True)
    ca_code = (
        "class CollisionAlert:\n"
        "  drone_a_id: int\n"
        "  drone_b_id: int\n"
        "  ttc: float        # time to collision\n"
        "  min_distance: float\n"
        "  conflict_point: Point3D\n"
        "  severity: str\n"
        "   # \"warning\" | \"critical\" | \"imminent\""
    )
    add_rounded_rect(slide, Inches(6.5), Inches(1.65), Inches(5.9), Inches(1.8),
                     RGBColor(0x0F, 0x17, 0x2A))
    add_text_box(slide, Inches(6.6), Inches(1.7), Inches(5.7), Inches(1.7),
                 ca_code, font_size=11, color=LIGHT, font_name=FONT_CODE)

    # FlightCorridor
    add_rounded_rect(slide, Inches(6.3), Inches(3.8), Inches(6.3), Inches(1.8),
                     BG_CARD, YELLOW)
    add_text_box(slide, Inches(6.5), Inches(3.9), Inches(5.9), Inches(0.3),
                 "FlightCorridor", font_size=13, color=YELLOW, bold=True)
    fc_code = (
        "class FlightCorridor:\n"
        "  corridor_id: str\n"
        "  waypoints: List[Point3D]\n"
        "  width: float = 20.0\n"
        "  altitude_layer: str  # low|medium|high\n"
        "  bidirectional: bool"
    )
    add_rounded_rect(slide, Inches(6.5), Inches(4.25), Inches(5.9), Inches(1.2),
                     RGBColor(0x0F, 0x17, 0x2A))
    add_text_box(slide, Inches(6.6), Inches(4.3), Inches(5.7), Inches(1.1),
                 fc_code, font_size=11, color=LIGHT, font_name=FONT_CODE)

    # Spatial indexing
    add_rounded_rect(slide, Inches(6.3), Inches(5.8), Inches(6.3), Inches(0.85),
                     BG_CARD, PURPLE)
    add_text_box(slide, Inches(6.5), Inches(5.9), Inches(5.9), Inches(0.25),
                 "공간 인덱싱", font_size=13, color=PURPLE, bold=True)
    add_text_box(slide, Inches(6.5), Inches(6.15), Inches(5.9), Inches(0.4),
                 "  KDTree3D -- O(log n) 이웃 탐색    |    VoxelGrid -- 3D 공간 분할",
                 font_size=11, color=DIM)

    add_footer(slide, 8)


def build_slide_9(prs):
    """Fail-Safe & self-healing."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_accent(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
                 "Fail-Safe & 자기 복구", font_size=36, color=CYAN, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.35),
                 "SC2 봇의 RuntimeSelfHealing을 드론의 안전 필수 시스템으로 전환합니다.",
                 font_size=14, color=DIM)

    # 3 fail-safe cards
    failsafes = [
        ("배터리 보호", RED,
         ["30% -> 귀환 경고", "20% -> 강제 귀환 (RTH)",
          "10% -> 긴급 착륙", "잔여 비행 시간 계산"]),
        ("통신 두절 대응", YELLOW,
         ["3초 -> 자율 Boids 모드", "10초 -> 호버링 대기",
          "30초 -> 자동 귀환", "메쉬 네트워크 우회"]),
        ("GPS 상실 대응", CYAN,
         ["IMU 관성 항법 전환", "이웃 드론 상대 위치",
          "시각적 오도메트리", "안전 착륙 시퀀스"]),
    ]
    for i, (title, clr, items) in enumerate(failsafes):
        x = Inches(0.4 + i * 4.2)
        add_rounded_rect(slide, x, Inches(1.5), Inches(3.9), Inches(2.4),
                         BG_CARD, clr)
        add_text_box(slide, x + Inches(0.2), Inches(1.6), Inches(3.5), Inches(0.3),
                     title, font_size=14, color=clr, bold=True,
                     alignment=PP_ALIGN.CENTER)
        txt = "\n".join(f"  {it}" for it in items)
        add_text_box(slide, x + Inches(0.2), Inches(2.0), Inches(3.5), Inches(1.7),
                     txt, font_size=12, color=DIM)

    # SC2 -> Drone Fail-Safe mapping
    add_rounded_rect(slide, Inches(0.4), Inches(4.1), Inches(6.0), Inches(2.5),
                     BG_CARD, GREEN)
    add_text_box(slide, Inches(0.6), Inches(4.2), Inches(5.5), Inches(0.3),
                 "SC2 -> 드론 Fail-Safe 매핑", font_size=13, color=GREEN, bold=True)

    fs_data = [
        ("SC2 SelfHealing", "드론 Fail-Safe"),
        ("경제 정체 감지", "배터리 소진 감지"),
        ("생산 중단 복구", "모터 페일오버"),
        ("매니저 크래시 복구", "센서 이중화 전환"),
        ("자원 낭비 감지", "비효율 경로 재계산"),
    ]
    tbl_shape = slide.shapes.add_table(5, 2,
                                        Inches(0.6), Inches(4.6),
                                        Inches(5.6), Inches(1.8))
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(2.8)
    tbl.columns[1].width = Inches(2.8)
    for r, row_data in enumerate(fs_data):
        tbl.rows[r].height = Inches(0.34)
        for c, txt in enumerate(row_data):
            cell = tbl.cell(r, c)
            cell.text = txt
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.08)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            p = cell.text_frame.paragraphs[0]
            p.font.name = FONT_BODY
            p.font.size = Pt(11)
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = CYAN_LIGHT
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x12, 0x1A, 0x2E)
            else:
                p.font.color.rgb = DIM2
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_CARD

    # Safety regulations
    add_rounded_rect(slide, Inches(6.8), Inches(4.1), Inches(5.8), Inches(2.5),
                     BG_CARD, PURPLE)
    add_text_box(slide, Inches(7.0), Inches(4.2), Inches(5.3), Inches(0.3),
                 "안전 규정 준수", font_size=13, color=PURPLE, bold=True)
    regs = [
        "  KALIS -- 한국 무인기 안전 기준",
        "  고도 제한 -- 120m 이하 (법적)",
        "  비행 금지 구역 -- GeoFence 적용",
        "  기체 식별 -- Remote ID 브로드캐스트",
        "  Fail-Safe -- 이중 센서 + 자동 귀환",
    ]
    add_text_box(slide, Inches(7.0), Inches(4.6), Inches(5.3), Inches(1.8),
                 "\n".join(regs), font_size=12, color=DIM)

    add_footer(slide, 9)


def build_slide_10(prs):
    """Conclusion & roadmap."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_accent(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
                 "결론 & 개발 로드맵", font_size=36, color=CYAN, bold=True)

    # Completed items
    add_rounded_rect(slide, Inches(0.4), Inches(1.2), Inches(5.6), Inches(4.2),
                     BG_CARD, CYAN)
    add_text_box(slide, Inches(0.6), Inches(1.3), Inches(5.2), Inches(0.3),
                 "구현 완료", font_size=14, color=CYAN, bold=True)
    completed = [
        "V  Boids 3D 컨트롤러 (8-Force)",
        "V  TTC 충돌 예측 (3단계 경보)",
        "V  가상 비행 회랑 (3 고도층)",
        "V  Authority Mode FSM (5 상태)",
        "V  DroneState / Point3D 데이터 모델",
        "V  KDTree3D 공간 인덱싱",
        "V  VoxelGrid 3D 공간 분할",
        "V  3D 시각화 데모 (Plotly, Ursina)",
        "V  애니메이션 GIF 5종",
    ]
    add_text_box(slide, Inches(0.6), Inches(1.7), Inches(5.2), Inches(3.5),
                 "\n".join(completed), font_size=13, color=DIM)

    # Roadmap
    add_rounded_rect(slide, Inches(6.3), Inches(1.2), Inches(6.3), Inches(3.0),
                     BG_CARD, YELLOW)
    add_text_box(slide, Inches(6.5), Inches(1.3), Inches(5.9), Inches(0.3),
                 "3단계 개발 로드맵", font_size=14, color=YELLOW, bold=True)

    stages = [
        ("S1", "Algorithm (현재)", "Boids 3D + TTC + Corridor\nPython 시뮬레이션 검증 완료", CYAN),
        ("S2", "Simulation (다음)", "ROS 2 + Gazebo 물리 시뮬\n바람/장애물/센서 노이즈 적용", PURPLE),
        ("S3", "Real Deploy (목표)", "Pixhawk + Raspberry Pi\nMAVLink 통신, 실외 비행 테스트", GREEN),
    ]
    for i, (badge, title, desc, clr) in enumerate(stages):
        y = Inches(1.75 + i * 0.8)
        # Badge circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(6.6), y, Inches(0.45), Inches(0.45)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = clr
        circle.line.fill.background()
        # Badge text
        add_text_box(slide, Inches(6.6), y + Inches(0.05), Inches(0.45), Inches(0.35),
                     badge, font_size=10, color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name=FONT_CODE)
        add_text_box(slide, Inches(7.2), y, Inches(2.5), Inches(0.25),
                     title, font_size=13, color=LIGHT, bold=True)
        add_text_box(slide, Inches(7.2), y + Inches(0.25), Inches(5.0), Inches(0.5),
                     desc, font_size=11, color=DIM)

    # Stats
    add_rounded_rect(slide, Inches(6.3), Inches(4.4), Inches(6.3), Inches(1.0),
                     BG_CARD, GREEN)
    add_text_box(slide, Inches(6.5), Inches(4.5), Inches(5.9), Inches(0.25),
                 "프로젝트 규모", font_size=13, color=GREEN, bold=True)
    stat_items = [
        ("8", "Boids 힘 벡터", CYAN),
        ("3", "고도층", YELLOW),
        ("3단계", "경보 체계", GREEN),
        ("5", "Authority 모드", PURPLE),
    ]
    for i, (num, label, clr) in enumerate(stat_items):
        x = Inches(6.5 + i * 1.5)
        add_text_box(slide, x, Inches(4.8), Inches(1.4), Inches(0.35),
                     num, font_size=20, color=clr, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name=FONT_CODE)
        add_text_box(slide, x, Inches(5.1), Inches(1.4), Inches(0.25),
                     label, font_size=10, color=DIM, alignment=PP_ALIGN.CENTER)

    # Thank You
    add_text_box(slide, Inches(2), Inches(5.8), Inches(9), Inches(0.5),
                 "Thank you", font_size=20, color=BODY_DIM,
                 alignment=PP_ALIGN.CENTER)

    tags = [("Boids 3D", CYAN), ("UTM", GREEN), ("TTC", YELLOW),
            ("Fail-Safe", RED), ("Sim-to-Real", PURPLE)]
    for i, (label, clr) in enumerate(tags):
        add_tag(slide, Inches(3.0 + i * 1.6), Inches(6.3), label, clr, BG_DARK)

    add_footer(slide, 10)


# ===== Main =====

def main():
    print("[1/3] Building 10 slides...")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    builders = [
        (build_slide_1, "Title"),
        (build_slide_2, "Project Overview"),
        (build_slide_3, "System Architecture"),
        (build_slide_4, "Boids 3D"),
        (build_slide_5, "TTC Collision"),
        (build_slide_6, "Flight Corridor"),
        (build_slide_7, "Authority Mode"),
        (build_slide_8, "Data Model"),
        (build_slide_9, "Fail-Safe"),
        (build_slide_10, "Conclusion"),
    ]

    for fn, name in builders:
        fn(prs)
        print(f"  [OK] {name}")

    print("[2/3] Saving PPTX...")
    output_path = Path(__file__).resolve().parent.parent / "wicked_zerg_challenger" / "visuals" / "Swarm_Net_Presentation_v2.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    print(f"[3/3] Done!")
    print(f"  File: {output_path}")
    print(f"  Slides: 10")
    print(f"  Theme: Dark Navy + Cyan/Green accents")


if __name__ == "__main__":
    main()

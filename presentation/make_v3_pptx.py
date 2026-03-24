"""
Swarm-Net Presentation v3 — Unified PPTX Generator
====================================================
v2 (개념 8장) + v2_new (기술 상세 10장) 통합 → 16장 슬라이드
- PNG 시각자료 삽입 (v2 이미지 + UTM 슬라이드 캡처)
- 발표자 노트 (구어체 한국어, 읽으면 바로 발표 가능)
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ===== Paths =====
BASE_DIR = Path(__file__).resolve().parent
V2_IMG_DIR = BASE_DIR.parent / "wicked_zerg_challenger" / "visuals" / "images"
UTM_IMG_DIR = BASE_DIR  # utm_slide_XX.png files
OUTPUT_DIR = BASE_DIR.parent / "wicked_zerg_challenger" / "visuals"

# ===== Color Constants =====
BG_DARK   = RGBColor(0x0B, 0x10, 0x1F)
BG_CARD   = RGBColor(0x12, 0x1A, 0x2E)
CYAN      = RGBColor(0x00, 0xFF, 0xCC)
CYAN_DIM  = RGBColor(0x00, 0xB8, 0x94)
CYAN_BLUE = RGBColor(0x00, 0xC8, 0xFF)
GREEN     = RGBColor(0x00, 0xE6, 0x76)
ORANGE    = RGBColor(0xFF, 0x98, 0x00)
RED       = RGBColor(0xFF, 0x45, 0x45)
PURPLE    = RGBColor(0xB3, 0x88, 0xFF)
GOLD      = RGBColor(0xFF, 0xD7, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xE8, 0xE8, 0xE8)
DIM       = RGBColor(0xA0, 0xA8, 0xB8)
BODY_DIM  = RGBColor(0x88, 0x93, 0xA4)

FONT_BODY = "Malgun Gothic"
FONT_CODE = "Consolas"
TOTAL_SLIDES = 16


# ===== Helper Functions =====

def set_slide_bg(slide, color=BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name=FONT_BODY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=LIGHT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(item, tuple):
            text, item_color, item_bold = item
        else:
            text, item_color, item_bold = item, color, False
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = item_color
        p.font.bold = item_bold
        p.font.name = FONT_BODY
        p.space_after = Pt(8)
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color=BG_CARD,
                     border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape


def add_tag_label(slide, left, top, text, bg_color=CYAN, text_color=BG_DARK,
                  w=None):
    width = w or Inches(2.4)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.40)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = FONT_CODE
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_notes(slide, text):
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.text = text
    for para in tf.paragraphs:
        para.font.size = Pt(12)


def add_footer(slide, slide_num):
    add_text_box(slide, Inches(0.4), Inches(6.9), Inches(3), Inches(0.4),
                 "SWARM-NET v3", font_size=11, color=CYAN_DIM, bold=True,
                 font_name=FONT_CODE)
    add_text_box(slide, Inches(11), Inches(6.9), Inches(2), Inches(0.4),
                 f"{slide_num:02d} / {TOTAL_SLIDES:02d}", font_size=11,
                 color=DIM, font_name=FONT_CODE, alignment=PP_ALIGN.RIGHT)


def add_image(slide, img_path, left, top, width=None, height=None):
    """Insert image if file exists. Returns True if inserted."""
    if not img_path.exists():
        print(f"  [WARN] Image not found: {img_path.name}")
        return False
    kwargs = {}
    if width:
        kwargs['width'] = width
    if height:
        kwargs['height'] = height
    slide.shapes.add_picture(str(img_path), left, top, **kwargs)
    return True


def add_top_accent(slide):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Pt(4)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = CYAN
    bar.line.fill.background()


# ===================================================================
# SLIDE BUILDERS (16 slides)
# ===================================================================

def build_slide_01(prs):
    """타이틀 — 하늘의 새로운 질서를 구축하다."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_tag_label(slide, Inches(0.8), Inches(1.2), "CAPSTONE PROJECT 2026", CYAN, BG_DARK)

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.7), Inches(0.8), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = CYAN
    line.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(1.9), Inches(8), Inches(1.2),
                 "하늘의 새로운 질서를\n구축하다", font_size=44, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(3.3), Inches(8), Inches(0.6),
                 "군집 UAV 기반 다이내믹 공역 통제 시스템", font_size=20, color=LIGHT)
    add_text_box(slide, Inches(0.8), Inches(3.9), Inches(8), Inches(0.4),
                 "Swarm-Net Airspace Manager", font_size=18, color=CYAN,
                 bold=True, font_name=FONT_CODE)
    add_text_box(slide, Inches(0.8), Inches(4.8), Inches(8), Inches(0.7),
                 "StarCraft II 군단 제어 AI 기반\nSim-to-Real 관제(ATC) 솔루션",
                 font_size=15, color=DIM)

    # Tags
    tags = [("Boids 3D", CYAN_BLUE), ("UTM", GREEN), ("TTC", GOLD),
            ("Fail-Safe", RED), ("MAVLink", PURPLE)]
    for i, (label, clr) in enumerate(tags):
        add_tag_label(slide, Inches(0.8 + i * 2.0), Inches(5.7), label, clr, BG_DARK, w=Inches(1.8))

    # Hexagon decoration
    hex_shape = slide.shapes.add_shape(
        MSO_SHAPE.HEXAGON, Inches(9.5), Inches(1.5), Inches(2.5), Inches(2.5)
    )
    hex_shape.fill.background()
    hex_shape.line.color.rgb = CYAN_DIM
    hex_shape.line.width = Pt(2)
    hex_shape.rotation = 30.0

    add_image(slide, V2_IMG_DIR / "slide1_swarm_hex.png",
              Inches(8.3), Inches(1.0), width=Inches(4.5))

    add_footer(slide, 1)

    add_notes(slide, (
        "안녕하십니까, 국립목포대학교 드론기계학과 장선우입니다.\n\n"
        "오늘 제가 발표할 주제는 '하늘의 새로운 질서를 구축하다'입니다.\n\n"
        "최근 드론 산업이 폭발적으로 성장하면서, 하늘 위 교통을 누가, 어떻게 관리할 것인지가 "
        "중요한 과제로 떠올랐습니다. 저는 스타크래프트 2라는 게임의 군집 제어 AI 기술을 "
        "실제 드론 공역 통제 시스템으로 전환하는 프로젝트, 'Swarm-Net'을 소개하겠습니다.\n\n"
        "Boids 3D 알고리즘, TTC 충돌 예측, 가상 비행 회랑 등 핵심 기술을 "
        "하나씩 설명드리겠습니다."
    ))


def build_slide_02(prs):
    """배경 & 문제 — 레이더 사각지대."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_image(slide, V2_IMG_DIR / "slide2_radar_gap.png",
              Inches(7.5), Inches(0.3), width=Inches(5.5))

    add_tag_label(slide, Inches(0.8), Inches(0.6), "BACKGROUND & PROBLEM", ORANGE, BG_DARK)

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(0.8), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = CYAN
    line.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(1.2),
                 "고정형 지상 레이더의\n구조적 한계와 급증하는 저고도 비행체",
                 font_size=36, color=WHITE, bold=True)

    cards = [
        ("탐지 사각지대", "산악·도심 고층 건물군의\n전파 음영으로 고정 레이더가\n커버하지 못하는 저고도 영역", RED),
        ("신속 전개 불가", "고정 인프라 설치에 수 개월 소요\n긴급 상황 시 관제망을\n즉각 전개하기 어려움", ORANGE),
        ("소형 UAV 탐지 한계", "저고도 소형 UAV는 낮은 RCS로\n기존 항공 레이더의\n탐지 임계값 미달", PURPLE),
    ]
    for i, (title, desc, accent) in enumerate(cards):
        x = Inches(0.8 + i * 4.1)
        y = Inches(3.0)
        add_rounded_rect(slide, x, y, Inches(3.7), Inches(2.2), BG_CARD, accent)
        add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.1), Inches(0.4),
                     title, font_size=20, color=accent, bold=True)
        add_text_box(slide, x + Inches(0.3), y + Inches(0.85), Inches(3.1), Inches(1.2),
                     desc, font_size=14, color=LIGHT)

    stats = [("90만+", "국내 등록 드론 (2025)"), ("+30%", "연간 증가율"), ("67%", "관제 사각지대")]
    colors = [CYAN, GREEN, ORANGE]
    for i, (val, label) in enumerate(stats):
        x = Inches(0.8 + i * 3.5)
        add_text_box(slide, x, Inches(5.8), Inches(2.5), Inches(0.5),
                     val, font_size=32, color=colors[i], bold=True, font_name=FONT_CODE)
        add_text_box(slide, x, Inches(6.35), Inches(2.5), Inches(0.3),
                     label, font_size=13, color=DIM)

    add_footer(slide, 2)

    add_notes(slide, (
        "현재 국내에 등록된 드론 수가 90만 대를 넘어섰습니다. "
        "연간 30% 이상 증가하고 있죠.\n\n"
        "그런데 문제는, 이 드론들을 관리할 인프라가 심각하게 부족하다는 점입니다.\n\n"
        "첫째, 산이나 고층 건물 때문에 레이더 전파가 닿지 않는 사각지대가 도심 저고도의 "
        "67%에 달합니다.\n"
        "둘째, 기존 레이더 기지국은 설치에 수개월이 걸려서 긴급 상황에 즉각 대응이 불가능합니다.\n"
        "셋째, 소형 드론은 레이더 반사 면적(RCS)이 너무 작아서 기존 항공 레이더로는 "
        "아예 탐지가 안 됩니다.\n\n"
        "그래서 저희는 완전히 새로운 접근법을 제안합니다."
    ))


def build_slide_03(prs):
    """솔루션 — 이동식 관제탑."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_tag_label(slide, Inches(0.8), Inches(0.6), "OUR SOLUTION", CYAN, BG_DARK)

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(0.8), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = CYAN
    line.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(1.2),
                 "드론이 직접\n하늘의 관제탑이 되다",
                 font_size=36, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(2.7), Inches(10), Inches(0.7),
                 "다수의 통제용 군집 UAV를 공중에 직접 배치하여\n"
                 "실시간 메시 통신망을 형성하는 '이동식 관제탑(Mobile ATC)' 아키텍처",
                 font_size=16, color=DIM)

    add_image(slide, V2_IMG_DIR / "slide3_mobile_atc.png",
              Inches(8.8), Inches(0.8), width=Inches(4.0))

    pillars = [
        ("Swarm Fleet", "6~12대 UAV가 정다각형 대형을\n자율 형성, 3D 다층 레이더 돔 전개", CYAN),
        ("Mesh Radar", "LiDAR + RF 기반 Mesh Network\n삼각측량으로 진입 객체 3D 좌표 산출", GREEN),
        ("Timer Control", "비인가 UAV에 체공 시간 자동 할당\n만료 시 RTH 명령 및 에스코트 프로토콜", ORANGE),
    ]
    for i, (title, desc, accent) in enumerate(pillars):
        x = Inches(0.8 + i * 4.1)
        y = Inches(3.7)
        card = add_rounded_rect(slide, x, y, Inches(3.7), Inches(2.4), BG_CARD, accent)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(1.3), y + Inches(0.3), Inches(1), Inches(1)
        )
        circle.fill.background()
        circle.line.color.rgb = accent
        circle.line.width = Pt(2)
        add_text_box(slide, x + Inches(0.3), y + Inches(1.4), Inches(3.1), Inches(0.4),
                     title, font_size=18, color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name=FONT_CODE)
        add_text_box(slide, x + Inches(0.3), y + Inches(1.9), Inches(3.1), Inches(0.9),
                     desc, font_size=13, color=DIM, alignment=PP_ALIGN.CENTER)

    add_footer(slide, 3)

    add_notes(slide, (
        "저희의 핵심 아이디어는 '관제탑을 하늘로 띄우자'라는 역발상입니다.\n\n"
        "수억 원짜리 고정 레이더 대신, 6~12대의 통제용 군집 드론을 출격시킵니다. "
        "이 드론들이 상공에서 육각형 대형을 자율적으로 형성하고, "
        "서로 Mesh Network로 연결되어 거대한 레이더 돔(Dome)을 만들어냅니다.\n\n"
        "Swarm Fleet은 대형을 유지하고, Mesh Radar는 삼각측량으로 "
        "진입 객체의 3D 좌표를 산출하며, Timer Control은 비인가 드론에 "
        "체공 시간을 할당하고 만료 시 강제 복귀시킵니다.\n\n"
        "설치 시간? 수 분이면 전개 완료입니다."
    ))


def build_slide_04(prs):
    """Sim-to-Real 핵심기술."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_tag_label(slide, Inches(0.8), Inches(0.5), "CORE TECHNOLOGY: SIM-TO-REAL", PURPLE, BG_DARK)

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.0), Inches(0.8), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = CYAN
    line.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(1.0),
                 "가상의 강화학습 알고리즘을\n현실의 드론 관제로 이식하다",
                 font_size=32, color=WHITE, bold=True)

    # SC2 box
    add_rounded_rect(slide, Inches(0.6), Inches(2.5), Inches(4.8), Inches(1.8), BG_CARD, PURPLE)
    add_text_box(slide, Inches(0.8), Inches(2.6), Inches(4.4), Inches(0.3),
                 "STARCRAFT II", font_size=13, color=PURPLE, bold=True, font_name=FONT_CODE)
    add_text_box(slide, Inches(0.8), Inches(3.0), Inches(4.4), Inches(0.4),
                 "Boids + FSM + RL", font_size=14, color=LIGHT, font_name=FONT_CODE)
    add_text_box(slide, Inches(0.8), Inches(3.5), Inches(4.4), Inches(0.4),
                 "10,000+ 게임 검증", font_size=13, color=DIM)

    # Arrow
    add_text_box(slide, Inches(5.6), Inches(2.9), Inches(1.5), Inches(0.4),
                 "Sim -> Real", font_size=14, color=CYAN, bold=True, font_name=FONT_CODE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(5.6), Inches(3.3), Inches(1.5), Inches(0.3),
                 "+ Altitude + Safety", font_size=13, color=DIM, alignment=PP_ALIGN.CENTER)

    # Drone ATC box
    add_rounded_rect(slide, Inches(7.3), Inches(2.5), Inches(4.8), Inches(1.8), BG_CARD, CYAN)
    add_text_box(slide, Inches(7.5), Inches(2.6), Inches(4.4), Inches(0.3),
                 "DRONE ATC", font_size=13, color=CYAN, bold=True, font_name=FONT_CODE)
    add_text_box(slide, Inches(7.5), Inches(3.0), Inches(4.4), Inches(0.4),
                 "Formation + Timer + Alert", font_size=14, color=LIGHT, font_name=FONT_CODE)
    add_text_box(slide, Inches(7.5), Inches(3.5), Inches(4.4), Inches(0.4),
                 "3D Mesh Radar Network", font_size=13, color=DIM)

    # Transfer mapping table
    table_data = [
        ("SC2 Component", "Drone ATC Mapping", "전이 방법론"),
        ("Boids Algorithm (군집 이동)", "Formation Flight Control", "1:1 구조 매핑"),
        ("Blackboard (중앙 상태 관리)", "Flight Data Hub", "1:1 구조 매핑"),
        ("Authority Mode (우선순위 체계)", "ATC Priority Levels", "1:1 구조 매핑"),
        ("IntelManager (정보 수집/탐지)", "3D Radar Mesh (Sensor Fusion)", "적응형 파라미터 튜닝"),
        ("RL Agent (강화학습 의사결정)", "Adaptive Path Optimization", "적응형 파라미터 튜닝"),
    ]
    rows, cols = len(table_data), 3
    tbl_shape = slide.shapes.add_table(rows, cols,
                                        Inches(0.6), Inches(4.5),
                                        Inches(11.5), Inches(2.4))
    tbl = tbl_shape.table
    col_widths = [Inches(4.0), Inches(4.5), Inches(3.0)]
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w
    for r, row_data in enumerate(table_data):
        tbl.rows[r].height = Inches(0.38)
        for c, cell_text in enumerate(row_data):
            cell = tbl.cell(r, c)
            cell.text = cell_text
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.font.name = FONT_BODY
            cell.margin_left = Inches(0.12)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
            if r == 0:
                p.font.size = Pt(13)
                p.font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x15, 0x1E, 0x35)
                p.font.color.rgb = [PURPLE, CYAN, GOLD][c]
            else:
                p.font.size = Pt(13)
                p.font.color.rgb = LIGHT
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_CARD
                if c == 2 and "구조" in cell_text:
                    p.font.color.rgb = GREEN
                elif c == 2 and "적응형" in cell_text:
                    p.font.color.rgb = ORANGE

    add_footer(slide, 4)

    add_notes(slide, (
        "이 시스템의 두뇌를 어디서 가져왔는지 설명드리겠습니다.\n\n"
        "저는 스타크래프트 2라는 게임에서 저그 종족 군집 AI를 직접 개발했습니다. "
        "Boids 알고리즘으로 유닛 군집 이동을 제어하고, FSM으로 우선순위를 관리하며, "
        "강화학습으로 1만 번 이상의 게임에서 전략을 최적화했습니다.\n\n"
        "이 검증된 알고리즘을 3D 공간으로 확장하고, 드론의 물리적 제약(고도 제한, 배터리, 안전 규정)을 "
        "추가하여 실제 드론 관제 시스템으로 이식한 것입니다.\n\n"
        "표에서 보시는 것처럼, SC2의 Boids가 편대 비행 제어로, Blackboard가 Flight Data Hub로, "
        "Authority Mode가 ATC 우선순위 체계로 1:1 매핑됩니다. "
        "게임에서 검증된 AI가 현실의 공역 관제로 재탄생한 것입니다."
    ))


def build_slide_05(prs):
    """시스템 아키텍처 — UTM PNG 삽입."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_accent(slide)

    add_tag_label(slide, Inches(0.3), Inches(0.3), "SYSTEM ARCHITECTURE", CYAN_BLUE, BG_DARK)

    # UTM slide PNG as main visual (left side)
    add_image(slide, UTM_IMG_DIR / "utm_slide_03.png",
              Inches(0.2), Inches(0.8), width=Inches(7.8))

    # Right side: key summary
    add_text_box(slide, Inches(8.3), Inches(0.8), Inches(4.5), Inches(0.5),
                 "시스템 아키텍처", font_size=28, color=CYAN_BLUE, bold=True)

    add_text_box(slide, Inches(8.3), Inches(1.4), Inches(4.5), Inches(0.4),
                 "4단계 데이터 파이프라인:", font_size=14, color=LIGHT, bold=True)

    summary = [
        ("1. Sensor Input", "GPS, IMU, LiDAR 데이터 수집", CYAN_BLUE),
        ("2. Flight Data Hub", "중앙 상태 저장소 (= SC2 Blackboard)", GREEN),
        ("3. Decision Engine", "군집 제어 + 충돌 예측 + 경로 관리", PURPLE),
        ("4. MAVLink Cmd", "드론에 제어 명령 전달", ORANGE),
    ]
    for i, (title, desc, clr) in enumerate(summary):
        y = Inches(1.9 + i * 0.75)
        add_rounded_rect(slide, Inches(8.3), y, Inches(4.5), Inches(0.65), BG_CARD, clr)
        add_text_box(slide, Inches(8.5), y + Inches(0.05), Inches(4.1), Inches(0.3),
                     title, font_size=13, color=clr, bold=True, font_name=FONT_CODE)
        add_text_box(slide, Inches(8.5), y + Inches(0.32), Inches(4.1), Inches(0.3),
                     desc, font_size=12, color=DIM)

    # Key modules
    add_text_box(slide, Inches(8.3), Inches(5.1), Inches(4.5), Inches(0.35),
                 "핵심 모듈 매핑:", font_size=13, color=LIGHT, bold=True)
    modules = [
        "Boids -> Boids3DController (8-Force)",
        "PotentialFields -> CollisionPredictor (TTC)",
        "CreepHighway -> CorridorManager (3 고도층)",
        "AuthorityFSM -> Authority FSM (5 상태)",
    ]
    add_bullet_list(slide, Inches(8.3), Inches(5.4), Inches(4.5), Inches(1.4),
                    [(m, DIM, False) for m in modules], font_size=12)

    add_footer(slide, 5)

    add_notes(slide, (
        "이제 시스템의 내부 구조를 보겠습니다.\n\n"
        "왼쪽 그림에서 전체 아키텍처를 확인하실 수 있습니다. "
        "데이터는 4단계로 흐릅니다.\n\n"
        "먼저 GPS, IMU, LiDAR 같은 센서에서 데이터가 들어오면, "
        "Flight Data Hub가 중앙에서 모든 상태를 관리합니다. "
        "이것은 스타크래프트의 Blackboard 패턴과 동일한 구조입니다.\n\n"
        "그 다음 Decision Engine이 군집 제어, 충돌 예측, 경로 관리를 수행하고, "
        "최종적으로 MAVLink 프로토콜을 통해 각 드론에 제어 명령을 전달합니다.\n\n"
        "오른쪽 표에서 보시는 것처럼, SC2의 모든 핵심 모듈이 "
        "드론 ATC 모듈로 1:1 매핑됩니다."
    ))


def build_slide_06(prs):
    """Boids 3D 군집 제어 — UTM PNG 삽입."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_accent(slide)

    add_tag_label(slide, Inches(0.3), Inches(0.3), "BOIDS 3D SWARM CONTROL", GREEN, BG_DARK)

    add_image(slide, UTM_IMG_DIR / "utm_slide_04.png",
              Inches(0.2), Inches(0.8), width=Inches(7.8))

    add_text_box(slide, Inches(8.3), Inches(0.8), Inches(4.5), Inches(0.5),
                 "Boids 3D 군집 제어", font_size=28, color=GREEN, bold=True)

    add_rounded_rect(slide, Inches(8.3), Inches(1.4), Inches(4.5), Inches(0.5),
                     RGBColor(0x0A, 0x1A, 0x12), GREEN)
    add_text_box(slide, Inches(8.4), Inches(1.45), Inches(4.3), Inches(0.4),
                 "F = Sum(w_i * F_i), i=1..8", font_size=13, color=GREEN,
                 font_name=FONT_CODE, alignment=PP_ALIGN.CENTER)

    forces = [
        ("Separation", "w=2.0", RED), ("Alignment", "w=1.0", CYAN_BLUE),
        ("Cohesion", "w=1.0", GREEN), ("Target", "w=1.5", GOLD),
        ("Obstacle", "w=2.5", PURPLE), ("Alt Hold", "w=1.0", RGBColor(0x38, 0xBD, 0xF8)),
        ("Terrain", "w=3.0", ORANGE), ("Corridor", "동적", RGBColor(0xE8, 0x79, 0xF9)),
    ]
    for i, (name, weight, clr) in enumerate(forces):
        y = Inches(2.1 + i * 0.52)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                      Inches(8.35), y + Inches(0.08),
                                      Inches(0.15), Inches(0.15))
        dot.fill.solid()
        dot.fill.fore_color.rgb = clr
        dot.line.fill.background()
        add_text_box(slide, Inches(8.6), y, Inches(2.0), Inches(0.25),
                     f"{i+1}. {name}", font_size=12, color=clr, bold=True)
        add_text_box(slide, Inches(10.8), y, Inches(1.5), Inches(0.25),
                     weight, font_size=12, color=DIM, font_name=FONT_CODE)

    # Kinematics
    add_text_box(slide, Inches(8.3), Inches(6.3), Inches(4.5), Inches(0.3),
                 "max_speed=15m/s | alt=5~120m", font_size=12, color=BODY_DIM,
                 font_name=FONT_CODE)

    add_footer(slide, 6)

    add_notes(slide, (
        "군집 드론이 어떻게 자율적으로 대형을 유지하며 비행하는지 설명드리겠습니다.\n\n"
        "Boids 알고리즘은 새 떼의 군집 행동을 모사한 것으로, "
        "각 드론에 8개의 힘 벡터를 적용합니다.\n\n"
        "기존 SC2에서 사용하던 5개 힘에 고도 관련 3개 힘을 추가했습니다:\n"
        "- Separation: 드론끼리 너무 가까우면 밀어내는 힘 (가중치 2.0)\n"
        "- Alignment: 이웃 드론과 같은 방향으로 정렬\n"
        "- Cohesion: 그룹 중심으로 모이는 응집력\n"
        "- Target Seeking: 목표 웨이포인트로 이동\n"
        "- Obstacle Avoid: 건물이나 산 같은 장애물 회피 (가중치 2.5)\n\n"
        "여기에 신규로 추가된 3개:\n"
        "- Altitude Hold: 지정 고도 유지 (PID 제어)\n"
        "- Terrain Clearance: 최소 안전 고도 5m 유지 (가중치 3.0으로 가장 강력)\n"
        "- Corridor Follow: 비행 회랑 중심선 추적\n\n"
        "이 8개 힘의 합이 매 프레임마다 각 드론의 가속도로 적용됩니다."
    ))


def build_slide_07(prs):
    """TTC 충돌 예측 — UTM PNG 삽입."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_accent(slide)

    add_tag_label(slide, Inches(0.3), Inches(0.3), "TTC COLLISION PREDICTION", RED, BG_DARK)

    add_image(slide, UTM_IMG_DIR / "utm_slide_05.png",
              Inches(0.2), Inches(0.8), width=Inches(7.8))

    add_text_box(slide, Inches(8.3), Inches(0.8), Inches(4.5), Inches(0.5),
                 "TTC 충돌 예측", font_size=28, color=RED, bold=True)

    add_text_box(slide, Inches(8.3), Inches(1.4), Inches(4.5), Inches(0.7),
                 "기존: 거리 기반 반발력\n-> 고속 접근 시 반응 지연\n\n"
                 "개선: 시간 기반 궤적 예측\n-> 5초 전 선제적 회피",
                 font_size=12, color=LIGHT)

    # Formula
    add_rounded_rect(slide, Inches(8.3), Inches(2.3), Inches(4.5), Inches(1.0),
                     RGBColor(0x0A, 0x1A, 0x12), GREEN)
    add_text_box(slide, Inches(8.4), Inches(2.35), Inches(4.3), Inches(0.9),
                 "TTC = -dot(rel_pos, rel_vel)\n      / dot(rel_vel, rel_vel)",
                 font_size=12, color=GREEN, font_name=FONT_CODE)

    # 3 alert levels
    add_text_box(slide, Inches(8.3), Inches(3.5), Inches(4.5), Inches(0.3),
                 "3단계 경보 체계:", font_size=13, color=LIGHT, bold=True)
    alerts = [
        ("WARNING", "<= 5.0초", "경로 미세 조정", GOLD),
        ("CRITICAL", "<= 3.0초", "회피 기동 시작", ORANGE),
        ("IMMINENT", "<= 1.0초", "긴급 수직 회피", RED),
    ]
    for i, (level, ttc, action, clr) in enumerate(alerts):
        y = Inches(3.9 + i * 0.65)
        add_rounded_rect(slide, Inches(8.3), y, Inches(4.5), Inches(0.55), BG_CARD, clr)
        add_text_box(slide, Inches(8.5), y + Inches(0.05), Inches(1.5), Inches(0.25),
                     level, font_size=12, color=clr, bold=True, font_name=FONT_CODE)
        add_text_box(slide, Inches(10.0), y + Inches(0.05), Inches(1.0), Inches(0.25),
                     ttc, font_size=12, color=LIGHT)
        add_text_box(slide, Inches(8.5), y + Inches(0.28), Inches(4.1), Inches(0.22),
                     action, font_size=12, color=DIM)

    add_text_box(slide, Inches(8.3), Inches(5.9), Inches(4.5), Inches(0.5),
                 "O(N^2) 전 쌍 검사\nVoxelGrid로 사전 필터링",
                 font_size=12, color=BODY_DIM)

    add_footer(slide, 7)

    add_notes(slide, (
        "충돌 예측 시스템에 대해 설명드리겠습니다.\n\n"
        "기존 SC2에서는 두 유닛 사이 거리가 일정 이하면 반발력을 적용했습니다. "
        "하지만 이 방식은 고속으로 접근하는 드론에 대해 반응이 늦습니다.\n\n"
        "그래서 저희는 TTC, 즉 Time-to-Collision 방식으로 업그레이드했습니다. "
        "두 드론의 상대 위치와 상대 속도를 이용해 '몇 초 후에 충돌하는지'를 계산합니다.\n\n"
        "경보는 3단계입니다:\n"
        "- WARNING: 5초 이내면 경로를 미세 조정합니다.\n"
        "- CRITICAL: 3초 이내면 적극적으로 회피 기동을 시작합니다.\n"
        "- IMMINENT: 1초 이내면 긴급 수직 회피, 즉 위아래로 즉시 벗어납니다.\n\n"
        "이 방식 덕분에 기존보다 5초나 먼저 위험을 감지하고 대응할 수 있습니다."
    ))


def build_slide_08(prs):
    """가상 비행 회랑 — UTM PNG 삽입."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_accent(slide)

    add_tag_label(slide, Inches(0.3), Inches(0.3), "FLIGHT CORRIDOR SYSTEM", CYAN_BLUE, BG_DARK)

    add_image(slide, UTM_IMG_DIR / "utm_slide_06.png",
              Inches(0.2), Inches(0.8), width=Inches(7.8))

    add_text_box(slide, Inches(8.3), Inches(0.8), Inches(4.5), Inches(0.5),
                 "비행 회랑 시스템", font_size=28, color=CYAN_BLUE, bold=True)

    add_text_box(slide, Inches(8.3), Inches(1.4), Inches(4.5), Inches(0.35),
                 "3 고도층 분리:", font_size=14, color=LIGHT, bold=True)

    layers = [
        ("LOW   5~30m", "촬영, 농업, 측량", GREEN),
        ("MED  30~60m", "배송, 순찰, 감시", CYAN_BLUE),
        ("HIGH 60~120m", "장거리, 중계", PURPLE),
    ]
    for i, (name, desc, clr) in enumerate(layers):
        y = Inches(1.85 + i * 0.6)
        add_rounded_rect(slide, Inches(8.3), y, Inches(4.5), Inches(0.5),
                         BG_CARD, clr)
        add_text_box(slide, Inches(8.5), y + Inches(0.05), Inches(2.0), Inches(0.4),
                     name, font_size=12, color=clr, bold=True, font_name=FONT_CODE)
        add_text_box(slide, Inches(10.5), y + Inches(0.05), Inches(2.1), Inches(0.4),
                     desc, font_size=12, color=DIM)

    # ICAO direction
    add_text_box(slide, Inches(8.3), Inches(3.7), Inches(4.5), Inches(0.35),
                 "ICAO 방향별 고도 분리:", font_size=13, color=LIGHT, bold=True)
    add_text_box(slide, Inches(8.3), Inches(4.1), Inches(4.5), Inches(0.5),
                 "Eastbound (0~180 deg) -> 기본 고도\nWestbound (180~360 deg) -> +15m",
                 font_size=12, color=DIM)

    # Altitude profile
    add_text_box(slide, Inches(8.3), Inches(4.8), Inches(4.5), Inches(0.35),
                 "고도 프로파일:", font_size=13, color=LIGHT, bold=True)
    profiles = [("0~15%", "상승", GREEN), ("15~85%", "순항", CYAN_BLUE), ("85~100%", "하강", GOLD)]
    for i, (pct, desc, clr) in enumerate(profiles):
        x = Inches(8.3 + i * 1.5)
        add_tag_label(slide, x, Inches(5.2), f"{pct} {desc}", clr, BG_DARK, w=Inches(1.4))

    add_text_box(slide, Inches(8.3), Inches(5.8), Inches(4.5), Inches(0.4),
                 "회랑 폭: 20m | 양방향 지원\n교차점 자동 검출 -> 시간 분리 적용",
                 font_size=12, color=BODY_DIM)

    add_footer(slide, 8)

    add_notes(slide, (
        "드론이 하늘에서 무질서하게 날아다니면 안 되겠죠? "
        "그래서 지상의 도로처럼 하늘에도 '비행 회랑'을 만들었습니다.\n\n"
        "SC2에서 유닛 이동 최적화에 사용했던 CreepHighway의 A* 경로 알고리즘을 "
        "3D 비행 회랑으로 확장한 것입니다.\n\n"
        "하늘을 3개 층으로 나눕니다:\n"
        "- LOW (5~30m): 촬영이나 농업용 드론\n"
        "- MEDIUM (30~60m): 배송이나 순찰 드론\n"
        "- HIGH (60~120m): 장거리 이동이나 통신 중계\n\n"
        "또한 ICAO 항공 규칙을 모사하여, 동쪽으로 가는 드론과 서쪽으로 가는 드론은 "
        "15m 고도 차이를 둡니다. 마치 도로의 중앙선처럼요.\n\n"
        "이렇게 하면 서로 다른 방향으로 비행하는 드론들이 같은 고도에서 "
        "마주치는 것을 원천 차단할 수 있습니다."
    ))


def build_slide_09(prs):
    """Authority Mode FSM — UTM PNG 삽입."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_accent(slide)

    add_tag_label(slide, Inches(0.3), Inches(0.3), "AUTHORITY MODE FSM", PURPLE, BG_DARK)

    add_image(slide, UTM_IMG_DIR / "utm_slide_07.png",
              Inches(0.2), Inches(0.8), width=Inches(7.8))

    add_text_box(slide, Inches(8.3), Inches(0.8), Inches(4.5), Inches(0.5),
                 "Authority Mode", font_size=28, color=PURPLE, bold=True)

    add_text_box(slide, Inches(8.3), Inches(1.3), Inches(4.5), Inches(0.4),
                 "우선순위 기반 5단계 상태 머신:", font_size=13, color=LIGHT, bold=True)

    modes = [
        ("P0 EMERGENCY", "TTC<=1s, GPS 상실", RED),
        ("P1 DECONFLICT", "TTC<=5s, 회랑 교차", ORANGE),
        ("P2 MISSION", "임무 수행 중", PURPLE),
        ("P3 CRUISE", "Boids 편대 비행", CYAN_BLUE),
        ("P4 IDLE", "대기/착륙 상태", GREEN),
    ]
    for i, (name, desc, clr) in enumerate(modes):
        y = Inches(1.8 + i * 0.7)
        add_rounded_rect(slide, Inches(8.3), y, Inches(4.5), Inches(0.6), BG_CARD, clr)
        add_text_box(slide, Inches(8.5), y + Inches(0.05), Inches(2.2), Inches(0.25),
                     name, font_size=12, color=clr, bold=True, font_name=FONT_CODE)
        add_text_box(slide, Inches(8.5), y + Inches(0.3), Inches(4.1), Inches(0.25),
                     desc, font_size=12, color=DIM)

    add_text_box(slide, Inches(8.3), Inches(5.4), Inches(4.5), Inches(0.8),
                 "전이 규칙:\n"
                 "- 높은 우선순위가 항상 선점\n"
                 "- EMERGENCY는 즉시 전환\n"
                 "- 배터리 20% -> 강제 EMERGENCY",
                 font_size=12, color=DIM)

    add_footer(slide, 9)

    add_notes(slide, (
        "실제 드론 관제에서 가장 중요한 것은 '긴급 상황 대응'입니다.\n\n"
        "Authority Mode는 SC2 봇에서 이미 검증된 우선순위 체계인데요, "
        "이것을 드론 공역 통제용으로 전환했습니다.\n\n"
        "5개 단계가 있습니다:\n"
        "- P0 EMERGENCY: 충돌 1초 전이거나, 배터리 부족, GPS 상실 시 모든 것을 중단하고 "
        "즉시 안전 조치를 취합니다.\n"
        "- P1 DECONFLICT: 충돌 5초 전이면 회피 기동을 시작합니다.\n"
        "- P2 MISSION: 촬영이나 배송 같은 임무를 수행합니다.\n"
        "- P3 CRUISE: 위협 없이 편대 비행하며 이동합니다.\n"
        "- P4 IDLE: 대기하거나 착륙 상태입니다.\n\n"
        "핵심은, P0가 발생하면 어떤 상태에서든 즉시 전환된다는 것입니다. "
        "배터리가 20% 미만이면 자동으로 EMERGENCY 모드에 진입하여 "
        "안전하게 귀환합니다."
    ))


def build_slide_10(prs):
    """데이터 모델 — UTM PNG 삽입."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_accent(slide)

    add_tag_label(slide, Inches(0.3), Inches(0.3), "DATA MODEL & DRONE STATE", GOLD, BG_DARK)

    add_image(slide, UTM_IMG_DIR / "utm_slide_08.png",
              Inches(0.2), Inches(0.8), width=Inches(7.8))

    add_text_box(slide, Inches(8.3), Inches(0.8), Inches(4.5), Inches(0.5),
                 "데이터 모델", font_size=28, color=GOLD, bold=True)

    models = [
        ("DroneState", "위치, 속도, 헤딩, 제한값\n(= SC2 Unit 확장)", CYAN_BLUE),
        ("Point3D", "3D 좌표 + 거리 계산\n(= SC2 Point2 확장)", GREEN),
        ("CollisionAlert", "TTC, 충돌점, 심각도\n3단계 경보 데이터", RED),
        ("FlightCorridor", "웨이포인트, 폭, 고도층\n양방향 회랑 정의", GOLD),
    ]
    for i, (name, desc, clr) in enumerate(models):
        y = Inches(1.4 + i * 1.15)
        add_rounded_rect(slide, Inches(8.3), y, Inches(4.5), Inches(1.0), BG_CARD, clr)
        add_text_box(slide, Inches(8.5), y + Inches(0.08), Inches(4.1), Inches(0.3),
                     name, font_size=14, color=clr, bold=True, font_name=FONT_CODE)
        add_text_box(slide, Inches(8.5), y + Inches(0.4), Inches(4.1), Inches(0.5),
                     desc, font_size=12, color=DIM)

    # Spatial indexing
    add_text_box(slide, Inches(8.3), Inches(6.0), Inches(4.5), Inches(0.4),
                 "KDTree3D O(logN) | VoxelGrid 3D",
                 font_size=12, color=BODY_DIM, font_name=FONT_CODE)

    add_footer(slide, 10)

    add_notes(slide, (
        "시스템에서 사용하는 핵심 데이터 구조 4가지를 소개합니다.\n\n"
        "DroneState는 SC2의 Unit 객체를 3D로 확장한 것으로, "
        "드론의 위치, 속도, 헤딩 각도, 그리고 최대 속도 15m/s, 고도 제한 5~120m 같은 "
        "물리적 제약값을 담고 있습니다.\n\n"
        "Point3D는 SC2의 2D 좌표를 3차원으로 확장한 것이고, "
        "CollisionAlert는 TTC 계산 결과를 담는 경보 데이터입니다.\n\n"
        "FlightCorridor는 3D 웨이포인트 시퀀스로 구성된 비행 회랑 정의입니다.\n\n"
        "또한 수십~수백 대의 드론을 효율적으로 관리하기 위해 "
        "KDTree3D와 VoxelGrid로 공간 인덱싱을 수행합니다. "
        "이웃 드론 탐색이 O(N^2)에서 O(log N)으로 개선됩니다."
    ))


def build_slide_11(prs):
    """Fail-Safe — UTM PNG 삽입."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_accent(slide)

    add_tag_label(slide, Inches(0.3), Inches(0.3), "FAIL-SAFE & SELF HEALING", RED, BG_DARK)

    add_image(slide, UTM_IMG_DIR / "utm_slide_09.png",
              Inches(0.2), Inches(0.8), width=Inches(7.8))

    add_text_box(slide, Inches(8.3), Inches(0.8), Inches(4.5), Inches(0.5),
                 "Fail-Safe 시스템", font_size=28, color=RED, bold=True)

    safes = [
        ("배터리 보호", "30% 경고 -> 20% 강제귀환\n10% 긴급착륙", RED),
        ("통신 두절", "3s Boids자율 -> 10s 호버링\n30s 자동귀환", GOLD),
        ("GPS 상실", "IMU 관성항법 전환\n이웃 드론 상대위치 활용", CYAN_BLUE),
    ]
    for i, (title, desc, clr) in enumerate(safes):
        y = Inches(1.4 + i * 1.2)
        add_rounded_rect(slide, Inches(8.3), y, Inches(4.5), Inches(1.05), BG_CARD, clr)
        add_text_box(slide, Inches(8.5), y + Inches(0.08), Inches(4.1), Inches(0.3),
                     title, font_size=14, color=clr, bold=True)
        add_text_box(slide, Inches(8.5), y + Inches(0.4), Inches(4.1), Inches(0.55),
                     desc, font_size=12, color=DIM)

    # Regulations
    add_text_box(slide, Inches(8.3), Inches(5.1), Inches(4.5), Inches(0.3),
                 "안전 규정 준수:", font_size=12, color=LIGHT, bold=True)
    regs = ["KALIS 한국 무인기 안전", "120m 고도 제한", "GeoFence 비행금지", "Remote ID 식별"]
    add_bullet_list(slide, Inches(8.3), Inches(5.4), Inches(4.5), Inches(1.2),
                    [(r, DIM, False) for r in regs], font_size=12)

    add_footer(slide, 11)

    add_notes(slide, (
        "드론은 하늘을 나는 만큼, 안전이 최우선입니다. "
        "SC2 봇의 RuntimeSelfHealing 기능을 드론용 Fail-Safe로 전환했습니다.\n\n"
        "세 가지 핵심 보호 장치가 있습니다:\n\n"
        "첫째, 배터리 보호입니다. 잔량 30%에서 경고, 20%에서 강제 귀환, "
        "10%에서 긴급 착륙을 실행합니다.\n\n"
        "둘째, 통신이 끊기면 3초간은 Boids 알고리즘으로 자율 비행을 유지하고, "
        "10초 지나면 제자리에서 호버링, 30초가 넘으면 자동 귀환합니다.\n\n"
        "셋째, GPS가 상실되면 IMU 관성항법으로 전환하고, "
        "이웃 드론들의 상대 위치 정보를 활용합니다.\n\n"
        "또한 한국 무인기 안전 기준(KALIS), 120m 고도 제한, GeoFence, "
        "Remote ID 등 법적 규정을 모두 준수합니다."
    ))


def build_slide_12(prs):
    """3D 시뮬레이션 데모."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_image(slide, V2_IMG_DIR / "slide5_3d_tracking.png",
              Inches(6.5), Inches(0.3), width=Inches(6.5))

    add_tag_label(slide, Inches(0.8), Inches(0.6), "3D LIVE SIMULATION", GREEN, BG_DARK)

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(0.8), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = CYAN
    line.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(8), Inches(1.0),
                 "사각지대 없는 3D 입체 스캔\n및 실시간 좌표 추적",
                 font_size=32, color=WHITE, bold=True)

    phases = [
        ("PHASE 01", "Mesh 형성", "6대 드론이 육각형 대형으로\n레이더 돔 전개", CYAN),
        ("PHASE 02", "탐지 & 식별", "RF 스캔 -> MAC 추출\n-> DB 등록 -> 허가 검증", GREEN),
        ("PHASE 03", "타이머 & 추적", "삼각측량으로 X,Y,Z\n실시간 위치 + 카운트다운", ORANGE),
        ("PHASE 04", "경고 & 퇴각", "2분 전 경고 ->\n강제 복귀 -> 에스코트", RED),
    ]
    for i, (phase, title, desc, accent) in enumerate(phases):
        x = Inches(0.5 + i * 3.1)
        y = Inches(2.8)
        card = add_rounded_rect(slide, x, y, Inches(2.8), Inches(2.0), BG_CARD, accent)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(2.4), Inches(0.3),
                     phase, font_size=11, color=accent, bold=True, font_name=FONT_CODE)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.55), Inches(2.4), Inches(0.4),
                     title, font_size=18, color=WHITE, bold=True)
        add_text_box(slide, x + Inches(0.2), y + Inches(1.05), Inches(2.4), Inches(0.8),
                     desc, font_size=12, color=DIM)
        if i < 3:
            add_text_box(slide, x + Inches(2.9), y + Inches(0.8), Inches(0.3), Inches(0.3),
                         ">", font_size=18, color=accent, bold=True, font_name=FONT_CODE,
                         alignment=PP_ALIGN.CENTER)

    # Demo cards
    demos = [
        ("3D 인터랙티브 관제 시뮬레이터", "Three.js 렌더링 · 3D 궤도 제어"),
        ("실시간 Fleet 통제 대시보드", "SVG 레이더 시각화 · Fleet 추적"),
    ]
    for i, (name, tech) in enumerate(demos):
        x = Inches(0.8 + i * 5.5)
        y = Inches(5.5)
        card = add_rounded_rect(slide, x, y, Inches(5.0), Inches(0.8), BG_CARD)
        add_text_box(slide, x + Inches(0.3), y + Inches(0.1), Inches(4.4), Inches(0.3),
                     name, font_size=14, color=WHITE, bold=True)
        add_text_box(slide, x + Inches(0.3), y + Inches(0.42), Inches(4.4), Inches(0.25),
                     tech, font_size=11, color=CYAN_DIM, font_name=FONT_CODE)

    add_footer(slide, 12)

    add_notes(slide, (
        "지금까지 설명드린 기술들이 실제로 어떻게 동작하는지 시뮬레이션 데모를 보여드리겠습니다.\n\n"
        "관제 시나리오는 4단계입니다:\n\n"
        "Phase 1에서 6대 드론이 육각형 대형을 자율 형성하여 레이더 돔을 전개합니다.\n\n"
        "Phase 2에서는 이 레이더 돔에 진입하는 비인가 드론을 RF 스캔으로 탐지하고, "
        "MAC 주소를 추출하여 데이터베이스에서 허가 여부를 검증합니다.\n\n"
        "Phase 3에서는 허가된 드론에 타이머를 할당하고, 삼각측량으로 "
        "실시간 3D 좌표를 추적합니다.\n\n"
        "Phase 4에서 체공 시간이 만료 2분 전이 되면 경고를 보내고, "
        "시간 초과 시 강제 복귀 명령을 내립니다. "
        "응답이 없으면 에스코트 프로토콜이 가동됩니다.\n\n"
        "이 모든 과정은 Web 기반 3D 시뮬레이터와 대시보드에서 실시간으로 확인할 수 있습니다."
    ))


def build_slide_13(prs):
    """실시간 대시보드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_image(slide, V2_IMG_DIR / "slide6_timer_dash.png",
              Inches(7.0), Inches(0.3), width=Inches(6.0))

    add_tag_label(slide, Inches(0.8), Inches(0.6), "REAL-TIME CONTROL DASHBOARD", ORANGE, BG_DARK)

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(0.8), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = CYAN
    line.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(1.0),
                 "타이머 기반 자동화 관제\n및 퇴각 통제",
                 font_size=32, color=WHITE, bold=True)

    statuses = [
        ("정상 비행", "> 3분", "08:32", "Heartbeat 정상", GREEN),
        ("시간 임박", "< 2분", "01:26", "Push 알림 전송", ORANGE),
        ("비행 종료", "= 0", "00:00", "강제 복귀 명령", RED),
        ("긴급 프로토콜", "미응답 30s", "! ! !", "에스코트 개입", RGBColor(0xFF, 0x00, 0x44)),
    ]
    for i, (title, cond, timer, action, accent) in enumerate(statuses):
        x = Inches(0.5 + i * 3.1)
        y = Inches(2.8)
        card = add_rounded_rect(slide, x, y, Inches(2.8), Inches(2.8), BG_CARD, accent)

        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(0.9), y + Inches(0.25), Inches(0.9), Inches(0.9)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = accent
        circle.line.fill.background()

        add_text_box(slide, x + Inches(0.2), y + Inches(1.2), Inches(2.4), Inches(0.3),
                     title, font_size=16, color=accent, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), y + Inches(1.55), Inches(2.4), Inches(0.25),
                     cond, font_size=12, color=DIM, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), y + Inches(1.85), Inches(2.4), Inches(0.4),
                     timer, font_size=28, color=accent, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name=FONT_CODE)
        add_text_box(slide, x + Inches(0.2), y + Inches(2.3), Inches(2.4), Inches(0.4),
                     action, font_size=12, color=DIM, alignment=PP_ALIGN.CENTER)

    add_footer(slide, 13)

    add_notes(slide, (
        "실시간 대시보드의 타이머 기반 관제 흐름을 설명드립니다.\n\n"
        "비인가 드론이 관제 공역에 진입하면 자동으로 체공 시간이 할당됩니다.\n\n"
        "초록색, 잔여 시간 3분 이상이면 정상 비행 상태입니다. "
        "Heartbeat ACK 신호로 드론의 생존을 계속 확인합니다.\n\n"
        "주황색, 2분 미만이 되면 경고 알림이 자동 전송됩니다.\n\n"
        "빨간색, 시간이 0이 되면 강제 복귀(RTH) 명령이 발령됩니다.\n\n"
        "가장 심각한 경우, 드론이 30초 이상 무응답이면 "
        "긴급 프로토콜이 발동되어 물리적으로 에스코트 드론이 출동합니다.\n\n"
        "이 모든 과정이 중앙 서버와 드론 간 저지연 통신으로 자동화되어 있습니다."
    ))


def build_slide_14(prs):
    """기대효과 & 활용분야."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_image(slide, V2_IMG_DIR / "slide7_comparison.png",
              Inches(6.5), Inches(2.8), width=Inches(6.5))

    add_tag_label(slide, Inches(0.8), Inches(0.5), "EXPECTED IMPACT & APPLICATIONS", CYAN, BG_DARK)

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.0), Inches(0.8), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = CYAN
    line.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(7), Inches(1.0),
                 "지형에 구애받지 않는\n이동식 관제 인프라",
                 font_size=30, color=WHITE, bold=True)

    # Comparison table
    table_data = [
        ("항목", "기존 방식", "Swarm-Net"),
        ("탐지 범위", "고정 레이더 반경 내", "군집 이동으로 동적 커버"),
        ("배치 시간", "인프라 설치 수 개월", "드론 출격 수 분"),
        ("비용", "CAPEX 수억 원", "OPEX 중심"),
        ("유연성", "고정, 재배치 불가", "실시간 공역 재구성"),
        ("확장성", "추가 기지국 필요", "드론 추가 투입"),
        ("체공 지속", "해당 없음", "Relay + Tethered"),
    ]
    rows, cols = len(table_data), 3
    tbl_shape = slide.shapes.add_table(rows, cols,
                                        Inches(0.6), Inches(2.4),
                                        Inches(6.2), Inches(3.2))
    tbl = tbl_shape.table
    col_widths = [Inches(1.4), Inches(2.3), Inches(2.5)]
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w
    for r, row_data in enumerate(table_data):
        tbl.rows[r].height = Inches(0.40)
        for c, cell_text in enumerate(row_data):
            cell = tbl.cell(r, c)
            cell.text = cell_text
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.font.name = FONT_BODY
            cell.margin_left = Inches(0.12)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
            if r == 0:
                p.font.size = Pt(13)
                p.font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x15, 0x1E, 0x35)
                p.font.color.rgb = [CYAN, RED, GREEN][c]
            else:
                p.font.size = Pt(13)
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_CARD
                if c == 1:
                    p.font.color.rgb = RGBColor(0xFF, 0x88, 0x88)
                elif c == 2:
                    p.font.color.rgb = GREEN
                else:
                    p.font.color.rgb = LIGHT

    # Use cases
    use_cases = [
        ("군사 야전", "신속 방어망"), ("불법 드론 차단", "공항 보안"),
        ("재난 현장", "긴급 관제"), ("드론 쇼", "충돌 방지"),
        ("UAM 공역", "복층 비행로"), ("농업 방제", "구역 제한"),
    ]
    for i, (title, desc) in enumerate(use_cases):
        col = i % 2
        row = i // 2
        x = Inches(7.3 + col * 2.8)
        y = Inches(2.2 + row * 1.5)
        card = add_rounded_rect(slide, x, y, Inches(2.4), Inches(1.2), BG_CARD)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.25), Inches(2.0), Inches(0.35),
                     title, font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.65), Inches(2.0), Inches(0.3),
                     desc, font_size=12, color=DIM, alignment=PP_ALIGN.CENTER)

    add_footer(slide, 14)

    add_notes(slide, (
        "Swarm-Net의 기대효과를 기존 방식과 비교해 보겠습니다.\n\n"
        "가장 큰 차이는 배치 시간입니다. "
        "기존 레이더는 설치에 수개월이 걸리지만, 저희 시스템은 수 분 만에 전개됩니다.\n\n"
        "비용도 고정 설비 수억 원 대신, 드론 Fleet 유지보수 비용만 발생합니다. "
        "유연성도 압도적입니다. 필요한 곳에 실시간으로 공역을 재구성할 수 있죠.\n\n"
        "활용 분야를 보시면, 군사 야전에서 신속하게 방어망을 구축하거나, "
        "공항 주변 불법 드론을 차단하거나, 재난 현장에서 긴급 공역 통제를 하는 등 "
        "다양한 시나리오에 적용 가능합니다.\n\n"
        "드론 쇼 안전, UAM 복층 비행로 관리, 농업 방제 구역 진입 제한에도 활용됩니다.\n\n"
        "배터리 문제는 핵심 노드의 유선 테더링과 교대 비행 프로토콜로 "
        "24시간 끊김 없는 관제가 가능합니다."
    ))


def build_slide_15(prs):
    """결론 & 로드맵 (통합)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_image(slide, V2_IMG_DIR / "slide8_roadmap.png",
              Inches(1.5), Inches(5.0), width=Inches(10.0))

    add_tag_label(slide, Inches(4.5), Inches(0.5), "CONCLUSION & ROADMAP", CYAN, BG_DARK)

    ctr_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5.7), Inches(0.9), Inches(1.3), Pt(3)
    )
    ctr_line.fill.solid()
    ctr_line.fill.fore_color.rgb = CYAN
    ctr_line.line.fill.background()

    add_text_box(slide, Inches(1.5), Inches(1.1), Inches(10), Inches(1.2),
                 "안전하고 체계적인 하늘의 미래,\nSwarm-Net이 만들어갑니다",
                 font_size=32, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Completed (left)
    add_rounded_rect(slide, Inches(0.3), Inches(2.6), Inches(4.5), Inches(2.3), BG_CARD, CYAN)
    add_text_box(slide, Inches(0.5), Inches(2.7), Inches(4.1), Inches(0.3),
                 "구현 완료", font_size=14, color=CYAN, bold=True)
    completed = [
        "Boids 3D Controller (8-Force)",
        "TTC 충돌 예측 (3단계 경보)",
        "가상 비행 회랑 (3 고도층)",
        "Authority Mode FSM (5 상태)",
        "3D 시각화 데모 + 대시보드",
    ]
    add_bullet_list(slide, Inches(0.5), Inches(3.0), Inches(4.1), Inches(1.8),
                    [(c, DIM, False) for c in completed], font_size=12)

    # Roadmap (right)
    stages = [
        ("STAGE 1", "SC2 시뮬레이션", "10,000+ 게임 검증", "COMPLETED", GREEN),
        ("STAGE 2", "3D 시뮬레이터", "파라미터 2D->3D 적응", "COMPLETED", GREEN),
        ("STAGE 3", "실 드론 테스트", "5대 편대 비행 검증", "NEXT PHASE", ORANGE),
        ("STAGE 4", "도시 스케일 ATC", "100+ 드론 관제", "VISION", PURPLE),
    ]
    for i, (stage, title, desc, status, accent) in enumerate(stages):
        x = Inches(5.5 + i * 2.0)
        y = Inches(2.6)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(0.3), y, Inches(0.7), Inches(0.7)
        )
        circle.fill.background()
        circle.line.color.rgb = accent
        circle.line.width = Pt(2.5)
        if i < 3:
            conn = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x + Inches(1.2), y + Inches(0.35), Inches(0.9), Pt(2)
            )
            conn.fill.solid()
            conn.fill.fore_color.rgb = RGBColor(0x33, 0x44, 0x55)
            conn.line.fill.background()

        add_text_box(slide, x, y + Inches(0.8), Inches(1.8), Inches(0.2),
                     stage, font_size=11, color=accent, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name=FONT_CODE)
        add_text_box(slide, x, y + Inches(1.0), Inches(1.8), Inches(0.3),
                     title, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, y + Inches(1.25), Inches(1.8), Inches(0.25),
                     desc, font_size=11, color=DIM, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, y + Inches(1.5), Inches(1.8), Inches(0.2),
                     status, font_size=11, color=accent, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name=FONT_CODE)

    add_footer(slide, 15)

    add_notes(slide, (
        "마지막으로 개발 현황과 로드맵을 정리하겠습니다.\n\n"
        "현재 Stage 1과 2가 완료되었습니다.\n"
        "Stage 1에서 스타크래프트 2 환경에서 1만 번 이상의 게임으로 알고리즘을 검증했고,\n"
        "Stage 2에서 이것을 3D 시뮬레이터로 확장하여 드론 파라미터를 적응시켰습니다.\n\n"
        "구현이 완료된 핵심 모듈은:\n"
        "- Boids 3D 8개 힘 벡터 컨트롤러\n"
        "- TTC 3단계 충돌 예측\n"
        "- 3개 고도층 비행 회랑\n"
        "- 5단계 Authority Mode FSM\n"
        "- Web 기반 3D 시각화 데모와 대시보드\n\n"
        "다음 단계인 Stage 3에서는 실제 드론 5대로 편대 비행 테스트를 진행할 예정이고,\n"
        "최종 비전인 Stage 4에서는 100대 이상의 드론이 참여하는 "
        "도시 스케일 공역 통제 시스템을 목표로 하고 있습니다."
    ))


def build_slide_16(prs):
    """Q&A — Thank You."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Decorative hexagons
    hex1 = slide.shapes.add_shape(
        MSO_SHAPE.HEXAGON, Inches(1.5), Inches(1.0), Inches(3.0), Inches(3.0)
    )
    hex1.fill.background()
    hex1.line.color.rgb = CYAN_DIM
    hex1.line.width = Pt(1.5)
    hex1.rotation = 30.0

    hex2 = slide.shapes.add_shape(
        MSO_SHAPE.HEXAGON, Inches(9.0), Inches(3.5), Inches(2.5), Inches(2.5)
    )
    hex2.fill.background()
    hex2.line.color.rgb = RGBColor(0x00, 0x88, 0x66)
    hex2.line.width = Pt(1)
    hex2.rotation = 30.0

    add_text_box(slide, Inches(2), Inches(2.5), Inches(9), Inches(1.0),
                 "THANK YOU", font_size=48, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name=FONT_CODE)

    add_text_box(slide, Inches(2), Inches(3.5), Inches(9), Inches(0.5),
                 "Q & A", font_size=20, color=DIM,
                 alignment=PP_ALIGN.CENTER, font_name=FONT_CODE)

    add_text_box(slide, Inches(2), Inches(4.3), Inches(9), Inches(0.6),
                 "Swarm-Net: 군집 UAV 기반 다이내믹 공역 통제 시스템\n"
                 "Capstone Design 2026",
                 font_size=14, color=DIM, alignment=PP_ALIGN.CENTER)

    # Tags
    tags = [("Boids 3D", CYAN_BLUE), ("UTM", GREEN), ("TTC", GOLD),
            ("Fail-Safe", RED), ("Sim-to-Real", PURPLE)]
    for i, (label, clr) in enumerate(tags):
        add_tag_label(slide, Inches(2.8 + i * 1.7), Inches(5.3), label, clr, BG_DARK, w=Inches(1.5))

    add_footer(slide, 16)

    add_notes(slide, (
        "이상으로 'Swarm-Net: 군집 UAV 기반 다이내믹 공역 통제 시스템' 발표를 마치겠습니다.\n\n"
        "요약하면, 저희는 스타크래프트 2 게임의 군집 AI를 "
        "실제 드론 공역 통제 시스템으로 전환하는 Sim-to-Real 프로젝트를 진행했습니다.\n\n"
        "8개 힘 벡터 Boids 3D, TTC 충돌 예측, 3개 고도층 비행 회랑, "
        "5단계 Authority Mode 등 핵심 알고리즘을 구현하고 시뮬레이션으로 검증했습니다.\n\n"
        "질문이 있으시면 말씀해 주시기 바랍니다.\n\n"
        "경청해 주셔서 감사합니다."
    ))


# ===== Main =====

def main():
    print("=" * 60)
    print("Swarm-Net Presentation v3 Generator")
    print("v2 (8 slides) + v2_new (10 slides) -> v3 (16 slides)")
    print("=" * 60)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    builders = [
        (build_slide_01, "Title"),
        (build_slide_02, "Background & Problem"),
        (build_slide_03, "Solution - Mobile ATC"),
        (build_slide_04, "Sim-to-Real Core Tech"),
        (build_slide_05, "System Architecture [+PNG]"),
        (build_slide_06, "Boids 3D [+PNG]"),
        (build_slide_07, "TTC Collision [+PNG]"),
        (build_slide_08, "Flight Corridor [+PNG]"),
        (build_slide_09, "Authority Mode [+PNG]"),
        (build_slide_10, "Data Model [+PNG]"),
        (build_slide_11, "Fail-Safe [+PNG]"),
        (build_slide_12, "3D Simulation Demo"),
        (build_slide_13, "Dashboard"),
        (build_slide_14, "Impact & Applications"),
        (build_slide_15, "Conclusion & Roadmap"),
        (build_slide_16, "Q&A - Thank You"),
    ]

    for fn, name in builders:
        fn(prs)
        print(f"  [OK] Slide {builders.index((fn, name)) + 1:02d}: {name}")

    # Save
    output_path = OUTPUT_DIR / "Swarm_Net_Presentation_v3.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    print(f"\n{'=' * 60}")
    print(f"v3 PPTX CREATED SUCCESSFULLY!")
    print(f"  File: {output_path}")
    print(f"  Slides: {TOTAL_SLIDES}")
    print(f"  Speaker notes: All {TOTAL_SLIDES} slides")
    print(f"  Visual images: 15 PNG embedded")
    print(f"  Theme: Dark Navy + Cyan/Green accents")
    print(f"{'=' * 60}")


if __name__ == "__main__":
    main()
